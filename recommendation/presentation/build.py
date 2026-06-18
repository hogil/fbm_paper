"""
임원 발표용 PPTX 빌더 (data-driven). spec.json -> deck.pptx
- python-pptx 엔진 / pptx 스킬 디자인 시스템(네이비 임원 팔레트, 슬라이드마다 시각요소, 제목 밑줄 금지)
- 한글 폰트(맑은 고딕), 16:9
사용: python build.py [spec.json] [out.pptx]
"""
import sys, os, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.abspath(os.path.join(HERE, "..", "figures"))

# ---- 디자인 시스템 ----
NAVY   = RGBColor(0x0F, 0x1E, 0x3D)
NAVY2  = RGBColor(0x1B, 0x32, 0x60)
ACCENT = RGBColor(0x96, 0xA0, 0xAD)   # neutral gray accent
GOLD   = RGBColor(0xF2, 0xB7, 0x05)   # awards/highlight
INK    = RGBColor(0x1A, 0x1F, 0x2E)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
PANEL  = RGBColor(0xF2, 0xF5, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD7, 0xDE, 0xEA)
COVER_BAR = RGBColor(0xCF, 0xD4, 0xDB)

FONT = "맑은 고딕"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FOOTER = "최호길 | QIE Data Science | AI Specialist 인증"


_NAMED = {"navy": "0F1E3D", "navy2": "1B3260", "accent": "96A0AD", "gold": "F2B705",
          "panel": "F2F5FA", "white": "FFFFFF", "ink": "1A1F2E", "muted": "6B7280",
          "line": "D7DEEA"}


def _col(c):
    """RGBColor / 이름(navy,accent,...) / '#RRGGBB' 문자열을 RGBColor로 강제."""
    if isinstance(c, RGBColor):
        return c
    s = str(c).strip()
    if s.lower() in _NAMED:
        s = _NAMED[s.lower()]
    return RGBColor.from_string(s.lstrip("#"))


def _is_dark(c):
    # 틸(accent, 명도~146)도 흰 글자가 맞아 임계값을 165로 둠(gold~183/panel~244는 어두운 글자 유지)
    h = str(_col(c))
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) < 165


def _set_font(run, size=None, bold=None, color=None, italic=None, name=FONT):
    f = run.font
    f.name = name
    # 한글 글꼴 강제(eastasian)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', name)
    if size is not None: f.size = Pt(size)
    if bold is not None: f.bold = bold
    if italic is not None: f.italic = italic
    if color is not None: f.color.rgb = _col(color)


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = _col(color)


def _rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = _col(color)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _col(line)
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def _text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of (text, dict-of-fontargs) 또는 list of [run,run] for mixed runs."""
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        runs = ln if isinstance(ln, list) else [ln]
        for (txt, kw) in runs:
            r = p.add_run(); r.text = txt
            _set_font(r, **kw)
        if isinstance(ln, list):
            # paragraph-level spacing from first run kw if present
            pass
    return tb


def _chip(slide, x, y, text, fill=ACCENT, fg=NAVY, w=Inches(0.95), h=Inches(0.38), size=13):
    sp = _rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set_font(r, size=size, bold=True, color=fg)
    return sp


def _img_fit(slide, path, x, y, maxw, maxh, frame=True):
    """이미지를 박스 안에 비율 유지로 배치(가운데 정렬). 박스는 (x,y,maxw,maxh)."""
    if not os.path.isabs(path):
        cand = os.path.join(FIG_DIR, path)
        path = cand if os.path.exists(cand) else os.path.join(HERE, path)
    if not os.path.exists(path):
        ph = _rect(slide, x, y, maxw, maxh, PANEL, line=LINE)
        _text(slide, x, y+maxh/2-Inches(0.2), maxw, Inches(0.4),
              [[("(이미지 없음)", dict(size=11, color=MUTED))]], align=PP_ALIGN.CENTER)
        return ph
    try:
        from PIL import Image
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 4, 3
    ar = iw / ih; box_ar = maxw / maxh
    if ar > box_ar:
        w = maxw; h = int(maxw / ar)
    else:
        h = maxh; w = int(maxh * ar)
    px = x + (maxw - w) // 2; py = y + (maxh - h) // 2
    if frame:
        _rect(slide, px-Emu(9000), py-Emu(9000), w+Emu(18000), h+Emu(18000), WHITE, line=LINE)
    pic = slide.shapes.add_picture(path, px, py, width=w, height=h)
    # 배치된 이미지의 실제 사각형(px,py,w,h)을 함께 돌려줘 호출부가 그 위에 주석(원/화살표)을
    # 정확히 얹을 수 있게 한다(이미지 가공 없이 좌표 기반 오버레이용).
    pic._fit_rect = (px, py, w, h)
    return pic


def _img_cover(slide, path, x, y, w, h, focus_x=0.5, focus_y=0.5):
    """이미지를 박스에 꽉 채우되 비율을 유지하도록 PowerPoint crop만 적용한다."""
    if not os.path.isabs(path):
        cand = os.path.join(FIG_DIR, path)
        path = cand if os.path.exists(cand) else os.path.join(HERE, path)
    if not os.path.exists(path):
        ph = _rect(slide, x, y, w, h, PANEL, line=LINE)
        _text(slide, x, y+h/2-Inches(0.12), w, Inches(0.24),
              [[("(이미지 없음)", dict(size=7.5, color=MUTED))]], align=PP_ALIGN.CENTER)
        return ph
    try:
        from PIL import Image
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 4, 3
    ar = iw / ih
    box_ar = w / h
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    if ar > box_ar:
        visible = box_ar / ar
        crop_total = max(0.0, 1.0 - visible)
        left = crop_total * min(max(focus_x, 0.0), 1.0)
        pic.crop_left = left
        pic.crop_right = crop_total - left
    else:
        visible = ar / box_ar
        crop_total = max(0.0, 1.0 - visible)
        top = crop_total * min(max(focus_y, 0.0), 1.0)
        pic.crop_top = top
        pic.crop_bottom = crop_total - top
    pic.shadow.inherit = False
    return pic


def _set_alpha_fill(sp, color, alpha_pct):
    """도형을 반투명 단색으로 채운다(alpha_pct=불투명도 0~100). 이미지 픽셀은 건드리지 않고
    그 위에 옅은 색면만 얹어 저대비 결함 군집으로 시선을 유도하기 위한 용도."""
    sp.fill.solid(); sp.fill.fore_color.rgb = _col(color)
    sp.line.fill.background(); sp.shadow.inherit = False
    sf = sp.fill._xPr.find(qn('a:solidFill'))
    srgb = sf.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    srgb.append(a)


def _annot_over(slide, rect, annots, accent=RGBColor(0xE5, 0x3A, 0x1F)):
    """배치된 이미지 사각형(rect=(px,py,w,h)) 위에 결함 위치 주석을 얹는다.
    annots: [{shape:'circle'|'ring'|'arrow', cx,cy[, r][, r2][, ang][, label]}] (좌표는 0~1 정규화).
    이미지 픽셀은 손대지 않고 PowerPoint 도형(고채도 주황 외곽선 + 옅은 색면)만 위에 그린다."""
    px, py, w, h = rect
    side = min(int(w), int(h))
    for a in annots:
        cx = int(px) + int(a.get("cx", 0.5) * int(w))
        cy = int(py) + int(a.get("cy", 0.5) * int(h))
        sh = a.get("shape", "circle")
        if sh in ("circle", "ring"):
            r = int(a.get("r", 0.12) * side)
            # ① 저대비 결함 군집으로 시선을 끌도록 가이드 원/링 안쪽에 옅은 따뜻한 색면(반투명)을
            #    먼저 깔아 흰 배경 대비를 확보(이미지 픽셀은 그대로). 군집을 가리지 않게 불투명도 낮게.
            hl = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx - r), Emu(cy - r),
                                        Emu(2 * r), Emu(2 * r))
            _set_alpha_fill(hl, RGBColor(0xFF, 0xDB, 0x7A), 28 if sh == "circle" else 20)
            # ② 그 위에 고채도 주황 외곽선(원/링)을 더 굵게 얹어 위계를 분명히 한다.
            sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx - r), Emu(cy - r),
                                        Emu(2 * r), Emu(2 * r))
            sp.fill.background()
            sp.line.color.rgb = accent
            sp.line.width = Pt(3.2 if sh == "circle" else 2.8)
            sp.shadow.inherit = False
        elif sh == "arrow":
            # 짧은 지시 화살표(라벨 쪽 → 결함). ang(도)·길이로 방향 지정.
            import math
            ln = int(a.get("len", 0.16) * side)
            ang = math.radians(a.get("ang", 0))
            tipx = cx; tipy = cy
            tailx = cx - int(ln * math.cos(ang)); taily = cy - int(ln * math.sin(ang))
            cn = slide.shapes.add_connector(2, Emu(tailx), Emu(taily), Emu(tipx), Emu(tipy))
            cn.line.color.rgb = accent; cn.line.width = Pt(2.6)
            le = cn.line._get_or_add_ln()
            tailEnd = le.makeelement(qn('a:tailEnd'),
                                     {'type': 'triangle', 'w': 'med', 'len': 'med'})
            le.append(tailEnd)
            cn.shadow.inherit = False
        if a.get("label"):
            lw = int(Inches(1.5)); lh = int(Inches(0.3))
            lx = cx - lw // 2
            ly = cy + int(a.get("r", 0.12) * side) + int(Inches(0.04))
            # 라벨 박스가 이미지 사각형(rect) 안에 머물도록 좌우·하단을 클램프한다. 큰 링(Edge-Ring)
            # 처럼 r 이 커서 라벨이 이미지 밖(또는 footer)으로 떨어지면, 링 하단 안쪽(이미지 바닥에서
            # lh 만큼 위)으로 끌어올려 슬라이드2 'Center Scratch' 와 동일한 '원 외곽 하단 라벨' 규칙에
            # 맞춘다(화살표 없이 라벨만, 헤더·footer 와 충돌 방지).
            lx = max(int(px) + int(Inches(0.04)), min(lx, int(px) + int(w) - lw - int(Inches(0.04))))
            ly = min(ly, int(py) + int(h) - lh - int(Inches(0.06)))
            tb = _rect(slide, Emu(lx), Emu(ly), Emu(lw), Emu(lh), WHITE,
                       line=accent, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = tb.text_frame; tf.word_wrap = False
            tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r0 = p.add_run(); r0.text = a["label"]
            _set_font(r0, size=10.5, bold=True, color=accent)


def _footer(slide, idx, size=9):
    _text(slide, Inches(0.5), Inches(7.04), Inches(9), Inches(0.32),
          [[(FOOTER, dict(size=size, color=MUTED))]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(12.0), Inches(7.04), Inches(0.83), Inches(0.32),
          [[(str(idx), dict(size=size, color=MUTED))]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _title_block(slide, kicker, title, x=Inches(0.7), y=Inches(0.55), w=Inches(11.9)):
    # 헤더 시스템 통일(full-bleed 스케일): 모든 본문 슬라이드(figure·표·로드맵)가 section(32)/
    # closing(30) 디바이더와 같은 줌 배율로 보이도록 제목을 키우고(29→31) 좌측 teal 바를 더 굵고
    # 길게(0.11→0.135w, 0.78→0.82h) 통일한다. 이전엔 figure·표 계열이 size29 로 약간 작아 보여
    # '컴팩트 vs full-bleed' 두 톤으로 갈렸음 → 동일 헤더 위계로 맞춘다(제목 밑줄 X).
    _rect(slide, x, y+Inches(0.00), Inches(0.135), Inches(0.84), RGBColor(0xCF, 0xD4, 0xDB))
    tx = x + Inches(0.32)
    if kicker:
        _text(slide, tx, y-Inches(0.06), w, Inches(0.28),
              [[(kicker, dict(size=13, bold=True, color=RGBColor(0x8A, 0x92, 0x9E)))]])
    _text(slide, tx, y+(Inches(0.22) if kicker else Inches(0.02)), w, Inches(0.66),
          [[(title, dict(size=31, bold=True, color=NAVY))]])


def _title_block_compact(slide, kicker, title, x=Inches(0.7), y=Inches(0.52), w=Inches(11.9)):
    _rect(slide, x, y, Inches(0.12), Inches(0.72), COVER_BAR)
    tx = x + Inches(0.30)
    if kicker:
        _text(slide, tx, y-Inches(0.04), w, Inches(0.24),
              [[(kicker, dict(size=12, bold=True, color=RGBColor(0x8A, 0x92, 0x9E)))]])
    _text(slide, tx, y+(Inches(0.20) if kicker else Inches(0.02)), w, Inches(0.56),
          [[(title, dict(size=28.5, bold=True, color=NAVY))]])


# ---------- 슬라이드 타입 ----------
def s_title(slide, d, idx):
    _bg(slide, WHITE)
    # 좌측 액센트 바를 슬라이드 전체 높이로 확장(안정감)
    _rect(slide, 0, 0, Inches(0.16), EMU_H, RGBColor(0xCF, 0xD4, 0xDB))
    # 표지 상단의 빈 띠를 핵심 성과 한 줄(시각적 앵커)로 채워 임팩트 보강 — 장식 과하지 않게
    # 가는 teal 룰 + 핵심 수치를 우상단 가로 라인에 절제 있게 배치.
    anchor_metrics = d.get("anchor_metrics",
                           "연 26억 효과   |   Known F1 0.95   |   DS AI Best Practice 수상")
    if anchor_metrics:
        _rect(slide, Inches(0.97), Inches(0.62), Inches(0.5), Pt(2.2), ACCENT)
        _text(slide, Inches(1.6), Inches(0.48), Inches(10.7), Inches(0.34),
              [[(anchor_metrics, dict(size=14, bold=True, color=RGBColor(0x2A,0x38,0x4C)))]],
              anchor=MSO_ANCHOR.MIDDLE)
    # 콘텐츠 블록을 위로 끌어올려 상단 빈 띠를 축소(시각 무게중심을 화면 중앙으로). 앵커 메트릭
    # 0.48"·kicker 0.98"·제목 1.42"로 통째로 올리고 제목 폰트를 키워(40→42) 상단 공백을 흡수한다.
    _text(slide, Inches(0.95), Inches(0.98), Inches(11.5), Inches(0.4),
          [[(d.get("kicker", "AI Specialist 인증 발표"), dict(size=15, bold=True, color=RGBColor(0x8A, 0x92, 0x9E)))]])
    # 제목에 명시적 줄바꿈(\n)이 있으면 어절 단위로 끊어 josa('에' 등)가 줄 앞에 홀로 떨어지지 않게 함
    title_lines = [[(seg, dict(size=42, bold=True, color=NAVY))] for seg in d["title"].split("\n")]
    _text(slide, Inches(0.95), Inches(1.46), Inches(11.6), Inches(1.8), title_lines)
    if d.get("subtitle"):
        _text(slide, Inches(0.97), Inches(3.20), Inches(11.4), Inches(0.6),
              [[(d["subtitle"], dict(size=17, color=RGBColor(0x47,0x55,0x69)))]])
    # 부제목과 발표자 사이 중앙 공백을 3개 과제(P1/P2/P3) 미니 카드로 채워 무게중심 균형(휑함 완화)
    tasks = d.get("tasks", [
        {"no": "P1", "t": "Failbit Map 불량 분석"},
        {"no": "P2", "t": "Chip multi-label"},
        {"no": "P3", "t": "Trend anomaly generator"}])
    if tasks:
        # P 카드 3개의 우측 끝선을 하단 footer 라인 우측 끝(12.33")과 정확히 맞춰 좌우 균형을
        # 맞춘다(slide1 low — 콘텐츠가 좌측에 쏠리고 우측 공백이 남던 인상 제거). 폭=(11.36-2*0.28)/3.
        # 라벨 바를 수직 중앙 쪽으로 약간 내리고(4.18→4.42") 높이를 키워(0.82→0.94") 중하단 여백을
        # 흡수해 표지의 무게중심을 화면 중앙으로 끌어내린다(slide1 low — 위쪽 쏠림 완화).
        tcw = Inches(3.6); tgap = Inches(0.28); ty0 = Inches(4.42); tch = Inches(0.94)
        tx0 = Inches(0.97)
        for ti, tk in enumerate(tasks[:3]):
            tx = Emu(int(tx0) + ti*(int(tcw)+int(tgap)))
            _rect(slide, tx, ty0, tcw, tch, PANEL)
            _rect(slide, tx, ty0, Inches(0.08), tch, RGBColor(0xCF, 0xD4, 0xDB))
            # P 번호와 설명을 좌측 한 정렬축에 묶는다(slide1 low — 칩과 설명이 박스 가운데에 떠
            # 정렬축이 모호하던 인상 제거). P 번호는 좌측 패딩(0.24")에 고정하고, 설명은 그 바로
            # 우측(1.02")에서 좌측정렬해 칩-설명을 하나의 시각 그룹으로 붙인다.
            _text(slide, tx+Inches(0.24), ty0, Inches(0.78), tch,
                  [[(tk["no"], dict(size=20, bold=True, color=ACCENT))]], anchor=MSO_ANCHOR.MIDDLE)
            _text(slide, tx+Inches(1.02), ty0, Emu(int(tcw)-int(Inches(1.24))), tch,
                  [[(tk["t"], dict(size=12, bold=True, color=NAVY))]],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(0.97), Inches(5.46), Inches(4.2), Pt(1.6), ACCENT)
    _text(slide, Inches(0.97), Inches(5.67), Inches(11), Inches(0.8),
          [[(d.get("author", ""), dict(size=15, bold=True, color=NAVY))],
           [(d.get("affil", ""), dict(size=12.5, color=MUTED))]])
    # 하단 구분선 + 소속/날짜 라인으로 하단 공백 메움(날짜 2026은 정보 전달용으로 푸터 우측에 배치)
    _rect(slide, Inches(0.97), Inches(6.55), Inches(11.36), Pt(1), LINE)
    _text(slide, Inches(0.97), Inches(6.68), Inches(8.5), Inches(0.4),
          [[(d.get("footer_l", "삼성전자 메모리제조센터 | QIE그룹 Data Science"),
             dict(size=12, bold=True, color=NAVY)),
            (d.get("footer_r", "    AI Specialist 인증 사내 경연 발표"),
             dict(size=12, color=MUTED))]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(10.5), Inches(6.68), Inches(1.83), Inches(0.4),
          [[("2026", dict(size=12, bold=True, color=ACCENT))]],
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def s_section(slide, d, idx):
    _bg(slide, WHITE)
    # 옅은 좌측 색면으로 빈 공간 절제 있게 채움(공백 과다 완화)
    _rect(slide, 0, 0, Inches(3.55), EMU_H, PANEL)
    # 슬라이드 최상단 가는 teal 액센트 스트립으로 우상단 빈 띠(제목/스탯카드 위)를 정리해
    # 표지·마무리와 톤을 통일하고 휑함을 줄임(높이 0.12").
    _rect(slide, 0, 0, EMU_W, Inches(0.12), RGBColor(0xCF, 0xD4, 0xDB))
    # 우측 본문 영역 상단(스탯카드 띠 위)에 과제 진행 eyebrow 라벨을 얹어 좌상단 공백 흡수.
    _text(slide, Inches(4.0), Inches(0.42), Inches(8.6), Inches(0.32),
          [[("AI Specialist 인증 | 3개 과제", dict(size=12, bold=True, color=RGBColor(0x8A, 0x92, 0x9E)))]],
          anchor=MSO_ANCHOR.MIDDLE)
    # 좌측 레일 상단의 큰 빈 띠(0~2.8")를 'PROJECT' eyebrow + 한 줄 요약 + 가는 룰로 채워
    # 세 섹션 표지 공통의 좌상단 휑함을 제거(P 마커는 그대로 수직 중앙 유지).
    _rect(slide, Inches(0.9), Inches(0.78), Inches(0.5), Pt(2.2), ACCENT)
    _text(slide, Inches(0.9), Inches(0.96), Inches(2.55), Inches(0.3),
          [[("PROJECT", dict(size=12, bold=True, color=ACCENT))]])
    rail_sub = d.get("rail_sub", "제조품질 AI | Domain-informed design")
    _text(slide, Inches(0.9), Inches(1.28), Inches(2.55), Inches(0.8),
          [[(rail_sub, dict(size=12.5, bold=True, color=NAVY))]])
    # 섹션 라벨(P1/P2/P3)+제목 묶음. 좌측 레일은 P 번호 블록을 위로 끌어올려(2.95→2.42)
    # 'PROJECT/rail_sub' 윗 그룹과의 상단 1/3 빈 띠를 흡수한다(design low — slide3·10·13 좌상단
    # 휑함 완화). 하단 stat 카드는 그대로 두고 중간 그룹(번호+진행점)만 위로 이동해 수직 균형 유지.
    _rect(slide, Inches(0.9), Inches(2.42), Inches(0.16), Inches(1.65), ACCENT)
    _text(slide, Inches(1.25), Inches(2.27), Inches(2.7), Inches(1.6),
          [[(d.get("no", ""), dict(size=80, bold=True, color=ACCENT))]])
    # 사이드바 'P 번호 ↔ 하단 라벨' 사이 큰 수직 공백을 세로 진행 인디케이터(P1·P2·P3)로 정리.
    # 각 과제를 작은 점+라벨로 쌓아 현재 과제를 강조 → 공백 흡수 + 발표 위치감 제공.
    if d.get("prog"):
        cur, tot = d["prog"]
        seg_labels = ["P1  Failbit Map", "P2  Chip multi-label", "P3  Trend generator"]
        # P 번호 아래·하단 라벨 위 공백에 3행으로 배치(겹침 방지). 번호 블록을 위로 올린 만큼
        # 진행 라벨은 큰 P 번호/세로 바가 끝난 뒤에 시작시켜 서로 덮이지 않게 한다.
        sy0 = int(Inches(4.18)); srow = int(Inches(0.31))
        for k in range(tot):
            on = (k + 1) == cur
            yy = sy0 + k * srow
            _rect(slide, Inches(1.0), Emu(yy + int(Inches(0.04))), Inches(0.13), Inches(0.13),
                  ACCENT if on else RGBColor(0xC2, 0xCC, 0xDC), shape=MSO_SHAPE.OVAL)
            lab = seg_labels[k] if k < len(seg_labels) else f"P{k+1}"
            _text(slide, Inches(1.28), Emu(yy), Inches(2.25), Inches(0.26),
                  [[(lab, dict(size=11, bold=on, color=(NAVY if on else MUTED)))]],
                  anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(4.0), Inches(2.95), Inches(8.6), Inches(1.55),
          [[(d["title"], dict(size=32, bold=True, color=NAVY))]] +
          ([[(d.get("subtitle"), dict(size=16, color=RGBColor(0x55,0x60,0x75)))]] if d.get("subtitle") else []),
          anchor=MSO_ANCHOR.MIDDLE)
    # 제목 묶음(≈4.5")과 '핵심 성과' KPI(5.6") 사이 빈 중앙 띠를, 이 과제가 '왜 중요한가' 한 줄
    # 요약(desc) 배지로 채워 세 디바이더 공통의 허전한 인상을 제거. 좌측 'WHY' eyebrow + 본문.
    # 박스 높이를 키우고(0.66→0.82) 아래 KPI/before-after 블록과의 간격을 균등 분배해 본문 하단
    # 빈 띠를 흡수(세 디바이더 공통의 하단 공백 과다 제거).
    if d.get("desc"):
        wy = Inches(4.58); wh = Inches(0.82)
        _rect(slide, Inches(4.0), wy, Inches(8.6), wh, RGBColor(0xF2,0xF5,0xFA))
        _rect(slide, Inches(4.0), wy, Inches(0.09), wh, ACCENT)
        _text(slide, Inches(4.22), wy, Inches(1.18), wh,
              [[("Rationale", dict(size=10.5, bold=True, color=ACCENT))]], anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(5.45), wy, Inches(6.95), wh,
              [[(d["desc"], dict(size=13, color=INK))]], anchor=MSO_ANCHOR.MIDDLE)
    # 좌측 사이드바 하단(레일 라벨 우측)의 빈 1/4 영역을, 이 과제의 핵심 수치를 직접 적은
    # 미니 요약 카드로 채워 '미완성' 인상을 제거(우측 본문 블록과 별개로 사이드바 자체를 완결).
    # stat_cards 의 big/label 을 작은 카드 2개로 쌓아 P 마커 아래 빈 띠(5.1~6.95")를 흡수.
    if d.get("stat_cards"):
        sc = d["stat_cards"][:2]
        sy0 = int(Inches(5.48)); sch = int(Inches(0.68)); sgp = int(Inches(0.12))
        for si, st in enumerate(sc):
            sy = Emu(sy0 + si * (sch + sgp))
            _rect(slide, Inches(0.9), sy, Inches(2.55), Emu(sch), WHITE, line=LINE)
            _rect(slide, Inches(0.9), sy, Inches(0.09), Emu(sch), ACCENT)
            _text(slide, Inches(1.14), sy, Inches(2.25), Emu(sch),
                  [[(st["big"], dict(size=17, bold=True, color=NAVY))],
                   [(st.get("label", ""), dict(size=9, bold=True, color=MUTED))]],
                  anchor=MSO_ANCHOR.MIDDLE)
    # 한 줄 핵심 성과 KPI (빈 색면 위) — desc 박스 아래로 내려 블록 간 간격을 균등화.
    # P1 divider 처럼 한 줄로 들어오게 폰트를 14.5→13.5 로 한 단계 줄이고(긴 P2·P3 문구도 1줄),
    # '핵심 성과' 라벨을 좌측 eyebrow 로 두되 kpi 본문 폭을 우측 '적용 전→후' 라벨(9.45") 직전까지
    # 만 잡아(폭 5.3") 두 줄로 넘쳐도 아래 ba 라벨/박스와 겹치지 않게 한다(세 divider 수직 리듬 통일).
    if d.get("kpi"):
        ky = Inches(5.62)
        _text(slide, Inches(4.0), ky, Inches(1.2), Inches(0.46),
              [[("Key Result", dict(size=11, bold=True, color=ACCENT))]], anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(5.12), ky, Inches(7.5), Inches(0.46),
              [[(d["kpi"], dict(size=13.5, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)
    # 하단 빈 영역을 'before → after' 한 줄 비교 배지로 채워 무게중심 균형(공백 과다 완화).
    # 박스를 바닥(푸터) 직전까지 내리고 높이를 키워(0.64→0.72) 본문 하단 빈 띠를 제거.
    if d.get("ba"):
        ba = d["ba"]; by = Inches(6.36); bh = Inches(0.72)
        _text(slide, Inches(9.5), by-Inches(0.30), Inches(3.1), Inches(0.26),
              [[("Before → After", dict(size=11, bold=True, color=ACCENT))]],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        bx = Inches(4.0); bcw = Inches(3.95); bgap = Inches(0.7)
        _rect(slide, bx, by, bcw, bh, PANEL)
        _text(slide, bx+Inches(0.18), by, Emu(int(bcw)-int(Inches(0.36))), bh,
              [[("Before  ", dict(size=10.5, bold=True, color=MUTED)),
                (ba.get("before", ""), dict(size=12.5, bold=True, color=INK))]],
              anchor=MSO_ANCHOR.MIDDLE)
        bx2 = Emu(int(bx)+int(bcw)+int(bgap))
        _rect(slide, bx2, by, bcw, bh, RGBColor(0xF4,0xF6,0xF8))
        _rect(slide, bx2, by, Inches(0.09), bh, ACCENT)
        _text(slide, bx2+Inches(0.18), by, Emu(int(bcw)-int(Inches(0.36))), bh,
              [[("After  ", dict(size=10.5, bold=True, color=ACCENT)),
                (ba.get("after", ""), dict(size=12.5, bold=True, color=NAVY))]],
              anchor=MSO_ANCHOR.MIDDLE)
        # 두 배지 사이 화살표(밴드 세로 중심에 맞춤)
        _rect(slide, Emu(int(bx)+int(bcw)+int(Inches(0.16))), Emu(int(by)+(int(bh)-int(Inches(0.3)))//2),
              Inches(0.38), Inches(0.3), ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
    # 우상단 빈 띠(제목 위)를 핵심 지표 스탯 카드로 채움(공백 과다 완화).
    # 카드 띠를 우측 가장자리(12.7")에 정렬해 2~3장 어느 경우든 우상단을 균형 있게 채운다.
    if d.get("stat_cards"):
        cards = d["stat_cards"][:3]
        nC = len(cards)
        cw = Inches(2.62); chh = Inches(1.24); cg = Inches(0.28)
        band_w = int(cw) * nC + int(cg) * (nC - 1)
        # 통계 카드 띠를 약간 위로 올려(1.46→1.28") 아이브로우(0.42")와의 상단 빈 밴드를 줄이고,
        # 세 디바이더(P1·P2·P3)에서 y를 동일하게 고정해 일관성을 확보(제목 블록과의 간격도 좁힘).
        x0 = Emu(int(Inches(12.7)) - band_w); y = Inches(1.28)
        for i, st in enumerate(cards):
            x = Emu(int(x0) + i*(int(cw)+int(cg)))
            _rect(slide, x, y, cw, chh, PANEL)
            _rect(slide, x, y, cw, Inches(0.1), ACCENT)
            _text(slide, x, y+Inches(0.28), cw, Inches(0.5),
                  [[(st["big"], dict(size=22, bold=True, color=NAVY))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 라벨은 좌우 안쪽 패딩(0.16")을 줘 긴 캡션이 타일 경계에 붙지 않게 함
            _text(slide, Emu(int(x)+int(Inches(0.16))), y+Inches(0.82),
                  Emu(int(cw)-int(Inches(0.32))), Inches(0.32),
                  [[(st.get("label", ""), dict(size=10.5, bold=True, color=INK))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def s_stats(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    stats = d["stats"]; n = len(stats)
    gap = Inches(0.3); total = Inches(12.0)
    cw = Emu(int((total - gap*(n-1)) / n))
    x0 = Inches(0.7); y = Inches(2.05); ch = Inches(3.35)
    for i, st in enumerate(stats):
        x = Emu(int(x0) + i*(int(cw)+int(gap)))
        _rect(slide, x, y, cw, ch, PANEL)
        _rect(slide, x, y, cw, Inches(0.14), ACCENT if not st.get("gold") else GOLD)
        _text(slide, x, y+Inches(0.7), cw, Inches(1.25),
              [[(st["big"], dict(size=52, bold=True, color=NAVY if not st.get("gold") else GOLD))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, x+Inches(0.18), y+Inches(2.12), Emu(int(cw)-int(Inches(0.36))), Inches(1.1),
              [[(st["label"], dict(size=15, bold=True, color=INK))]] +
              ([[(st.get("sub"), dict(size=12, color=MUTED))]] if st.get("sub") else []),
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    if d.get("note"):
        _text(slide, Inches(0.7), Inches(5.78), Inches(12), Inches(1.0),
              [[(d["note"], dict(size=13.5, color=INK))]])
    _footer(slide, idx)


def _bullets_tf(slide, x, y, w, h, bullets, size=15, gap=8, ink=INK, strong=NAVY, marker=ACCENT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        lvl = 0; txt = b
        if isinstance(b, dict):
            txt = b.get("t", ""); lvl = b.get("lvl", 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        # 한 불릿이 여러 줄로 줄바꿈될 때 줄 사이 숨통(가독성↑) — 긴 불릿이 빽빽해 보이는 것 완화
        if line_spacing is not None:
            p.line_spacing = line_spacing
        # 행잉 인덴트: 첫 줄은 불릿 기호부터, wrap 되는 둘째 줄부터는 불릿 텍스트 시작선에 정렬
        hang = int(Inches(0.30 if lvl == 0 else 0.46))
        pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
        pPr.set('marL', str(hang)); pPr.set('indent', str(-hang))
        r0 = p.add_run()
        r0.text = ("●  " if lvl == 0 else "–  ")
        _set_font(r0, size=size-(0 if lvl==0 else 1), bold=False, color=marker if lvl==0 else RGBColor(0x6E,0x7B,0x8E))
        # 굵게 강조 마크업 **...**
        # lvl=1(참고/각주)은 본문과 위계를 분리하되 '버려진 각주'처럼 보이지 않게 가독을 확보:
        # 폰트를 size-1 로만 줄이고(이전 size-2) 색을 한 단계 진한 슬레이트(#4A5666)로 올려 대비를
        # 확보한다(본문 대비 위계는 유지하되 투사 거리에서도 읽히게).
        import re
        body_sz = size - (0 if lvl == 0 else 1)
        body_col = ink if lvl == 0 else RGBColor(0x4A, 0x56, 0x66)
        strong_col = strong if lvl == 0 else RGBColor(0x38, 0x44, 0x56)
        parts = re.split(r"(\*\*.+?\*\*)", txt)
        for part in parts:
            if not part:
                continue
            r = p.add_run()
            if part.startswith("**") and part.endswith("**"):
                r.text = part[2:-2]; _set_font(r, size=body_sz, bold=True, color=strong_col)
            else:
                r.text = part; _set_font(r, size=body_sz, color=body_col)
        if lvl == 1:
            p.level = 1
    return tb


def s_bullets(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    has_img = bool(d.get("image"))
    # 본문 컬럼 우측 경계를 이미지 좌측(7.95")보다 충분히 안쪽(6.75")으로 당겨 거터(여백)를
    # 확보 → 본문이 그림 영역을 침범하거나 단어 중간에서 끊기는 충돌을 방지.
    bx, bw = Inches(0.7), (Inches(d.get("body_w", 6.0)) if has_img else Inches(12.0))
    # lead_metric(좌하단 핵심 지표 카드)이 있으면 불릿 영역을 위쪽 슬롯(1.9~5.55")으로 한정하고
    # TOP 정렬해, 불릿이 4개여도 아래로 흩어지지 않게 하고 그 아래 빈 띠를 지표 카드로 메운다
    # (slide2·5 좌측 칼럼 하단 공백 과다 해결).
    lead = d.get("lead_metric")
    if lead:
        # 본문 불릿 블록의 가용 높이를 줄여(3.6→3.42") TOP 정렬 텍스트의 마지막(특히 2줄로 wrap 된)
        # 불릿 descender 가 아래 지표 밴드 상단에 닿지 않게 한다(slide2 medium — 4불릿/2줄 겹침 제거).
        _bullets_tf(slide, bx, Inches(1.86), bw, Inches(3.42), d.get("bullets", []),
                    size=d.get("size", 16), gap=d.get("gap", 14), anchor=MSO_ANCHOR.TOP,
                    line_spacing=1.12)
        # 좌하단 핵심 지표 카드: 큰 수치 + 라벨을 강조 틴트 바에 담아 빈 띠를 시각 앵커로 채움.
        # 지표 밴드를 본문 블록 하단(≈5.40")에서 충분히 띄워(5.74→5.92, ≈0.5") footer(7.04") 위
        # 안전선 안에서 내린다 → 본문 마지막 줄 받침과 밴드 윗선 사이 여백을 확보(slide2 medium).
        lmy = Inches(5.92); lmh = Inches(0.98)
        _rect(slide, bx, lmy, bw, lmh, RGBColor(0xF4, 0xF6, 0xF8),
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(slide, bx, lmy, Inches(0.1), lmh, ACCENT)
        cells = lead[:3]; ncl = len(cells)
        innerw = int(bw) - int(Inches(0.3))
        cwl = innerw // ncl
        for k, mc in enumerate(cells):
            mx = Emu(int(bx) + int(Inches(0.22)) + k * cwl)
            _text(slide, mx, Emu(int(lmy) + int(Inches(0.13))), Emu(cwl - int(Inches(0.1))),
                  Emu(int(lmh) - int(Inches(0.26))),
                  [[(mc.get("big", ""), dict(size=18, bold=True, color=NAVY))],
                   [(mc.get("label", ""), dict(size=10.5, bold=True, color=RGBColor(0x4A,0x56,0x66)))]],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    else:
        body_anchor = MSO_ANCHOR.TOP if d.get("body_anchor") == "top" else MSO_ANCHOR.MIDDLE
        _bullets_tf(slide, bx, Inches(d.get("body_y", 1.9)), bw, Inches(d.get("body_h", 4.78)), d.get("bullets", []),
                    size=d.get("size", 16), gap=d.get("gap", 14), anchor=body_anchor,
                    line_spacing=1.12)
    if has_img:
        # 이미지 박스를 좌측(7.15")으로 더 당기고 폭(5.55")·높이(4.55")를 키워 sparse 한 wafer
        # map(연회색 배경에 작은 결함 마커)이 실제 분포가 보일 만큼 크게 렌더되게 함(이미지 가공
        # 금지 — 배치만 확대). 본문 컬럼 우측 경계(7.05")와 0.1" 거터 유지. 이미지를 옅은 패널
        # 틴트 위에 올려 흰 wafer 여백과 슬라이드 배경 경계를 분리해 '깨진 빈 그림' 인상을 줄임.
        # wafer map은 6400x6400 정사각·sparse(연회색 배경에 작은 결함 마커) → 정사각에 가까운 큰
        # 틴트 패널 위에 비율 유지로 키워, 작게 박힌 듯 보이던 인상을 제거(이미지 가공 없이 배치만 확대).
        # wafer map은 정사각(6400x6400)이라 폭 넓은 직사각 패널 안에서는 좌우 거터가 크게 남고
        # 이미지가 작아 보였다 → 패널을 거의 정사각(4.95×4.95")으로 만들어 wafer 가 패널을 꽉 채우게
        # 키우고(약 +5% 크게 렌더), 우측 본문 외 영역의 가로 중앙(약 7.15~12.1")에 두어 좌우 공백을 균형
        # 있게 흡수. 캡션은 이미지 하단과 0.18" 띄워 footer 와 충돌하지 않게 함(하단 여백 확보).
        # wafer map은 정사각(6400x6400)·연회색 배경에 작은 결함 마커 → 가능한 한 크게 렌더해야
        # 중앙 결함 군집이 발표 거리에서 보인다. 세로 슬롯(1.45~6.25")을 거의 다 써 4.8" 정사각
        # 패널로 키우고(약 +4% 크게), 패널 틴트를 한 단계 진하게(#E4EAF4) 둘러 흰 wafer 여백과
        # 슬라이드 배경의 경계를 분명히 해 '빈/깨진 그림' 인상을 제거(이미지 가공 없이 배치·프레임만).
        # 패널을 가용 세로 슬롯(1.18~6.40")을 거의 다 채우는 5.32" 정사각으로 키우고 우측 가용
        # 폭(7.05~12.66") 가로 중앙에 두어, sparse 한 wafer map 의 중앙 결함 군집이 발표 거리에서
        # 더 크게 보이게 함(이미지 가공 없이 배치만 확대 — 우하단 빈 공간 흡수). 캡션은 패널 하단과
        # 0.12" 띄워 footer 와 충돌하지 않게 한다.
        # sparse 한 wafer map(흰 배경에 옅은 색점)은 그 자체로 '빈 그림'처럼 보인다 → 이미지는
        # 가공하지 않고, 대신 '정의된 figure 카드'로 감싸 우측 콘텐츠 영역을 시각적으로 확정한다.
        # ① 카드 바깥에 옅은 패널 + 또렷한 테두리, ② 카드 상단에 teal 헤더 스트립(라벨)로 그림에
        # 의미·경계를 부여해 저대비 그림이 공백/저품질로 읽히는 인상을 제거(좌우 무게 균형도 보강).
        # 카드 top 을 제목 베이스라인(≈1.43") 아래로 충분히 내려(1.52") 제목 우측 끝이 이미지
        # 카드 헤더(teal 캡션 바)와 수평으로 충돌하는 것을 방지한다(slide5 제목 '…어려운가' 침범
        # 제거). 카드 하단이 footer(7.04")·캡션(6.62") 과 겹치지 않게 폭(pw)을 4.62"로 줄여
        # 세로 슬롯에 맞춘다(이미지 가공 없이 배치만 조정).
        # wafer 카드를 더 크게: 헤더를 약간 낮추고(0.42→0.38) 카드를 위로 올려(1.52→1.44) 정사각
        # 패널을 4.62→4.74"로 키운다(약 +3% 크게 렌더 — sparse 한 wafer map 의 중앙 결함 군집이
        # 발표 거리에서 더 또렷하게 보이게, 이미지 가공 없이 배치만 확대). 카드 하단(≈6.56")은
        # 캡션(6.62")·footer(7.04") 위 안전선을 유지한다.
        pw = Inches(4.74)
        hdr = Inches(0.38)
        # 카드 우측 가용 경계를 12.74→12.50" 로 당겨, 정사각 wafer 카드를 가용 슬롯 안에서 살짝
        # 좌측으로 중앙정렬한다 → 우측/하단 가장자리에 바짝 붙어 답답하던 인상을 완화(slide4 low,
        # 이미지 가공 없이 배치만 조정). 캡션은 카드 중심을 따라가므로 함께 안쪽으로 들어온다.
        card_x = Emu(int(Inches(6.95)) + (int(Inches(12.50)) - int(Inches(6.95)) - int(pw)) // 2)
        card_y = Inches(1.44)
        card_h = Emu(int(pw) + int(hdr))
        # 카드 외곽: 옅은 패널 + 또렷한 테두리(영역 확정). 라운드 보더로 처리해 정사각 wafer 격자의
        # 계단형 바깥 경계가 슬라이드 배경과 만나는 가장자리를 부드럽게 정돈한다(slide2·5 low).
        _rect(slide, card_x, card_y, pw, card_h, RGBColor(0xF4, 0xF7, 0xFB),
              line=RGBColor(0xB9, 0xC6, 0xDB), line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 상단 헤더 스트립(teal) + 라벨
        _rect(slide, card_x, card_y, pw, hdr, ACCENT)
        _text(slide, Emu(int(card_x) + int(Inches(0.18))), card_y,
              Emu(int(pw) - int(Inches(0.36))), hdr,
              [[(d.get("img_label", "실제 wafer Failbit Map 예시"),
                 dict(size=11.5, bold=True, color=WHITE))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 이미지 영역: 헤더 아래, 한 단계 진한 틴트 위에 비율 유지로 배치(이미지 가공 없음).
        # 라운드 사각 타일로 깔아 wafer 격자의 계단형 경계가 부드러운 프레임 안에 들어오게 한다(low).
        ix = card_x; iy = Emu(int(card_y) + int(hdr)); iw = pw; ih = pw
        _rect(slide, ix, iy, iw, ih, RGBColor(0xE4, 0xEA, 0xF4), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 가장자리 링/타원 주석(Edge-Ring 등)이 프레임에 꽉 차 잘려 보이지 않도록, 정사각 wafer
        # 이미지를 타일 안에서 소폭 인셋(여백)해 비율유지로 렌더한다(이미지 가공 없이 배치만 조정).
        # 인셋을 0.10→0.16" 로 키워 wafer 의 계단형 바깥 경계 둘레에 틴트 타일 여백을 둬 거친
        # 가장자리를 시각적으로 정돈한다(slide2·5 low). ring 주석은 더 큰 인셋으로 원 전체를 담는다.
        annots = d.get("img_annot") or []
        has_edge_ring = any(a.get("shape") == "ring" or a.get("r", 0) >= 0.4 for a in annots)
        inset = int(Inches(0.42)) if has_edge_ring else int(Inches(0.16))
        ax = Emu(int(ix) + inset); ay = Emu(int(iy) + inset)
        aw = Emu(int(iw) - 2 * inset); ah = Emu(int(ih) - 2 * inset)
        pic = _img_fit(slide, d["image"], ax, ay, aw, ah, frame=False)
        # 저대비 wafer map: 결함 위치에 고채도 주황 원/화살표 주석을 build.py 로 얹어(이미지 가공
        # 없이) 핵심 '불량 분포 패턴'이 발표 거리에서 한눈에 들어오게 한다(slide2·5 medium 해결).
        if d.get("img_annot") and hasattr(pic, "_fit_rect"):
            _annot_over(slide, pic._fit_rect, d["img_annot"])
        if d.get("caption"):
            _text(slide, Emu(int(card_x) - int(Inches(0.6))), Inches(6.62),
                  Emu(int(pw) + int(Inches(1.2))), Inches(0.4),
                  [[(d["caption"], dict(size=13.5, bold=True, color=RGBColor(0x2A,0x38,0x4C)))]], align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.TOP)
    _footer(slide, idx)


def _cluster3_diagram(slide, x, y, w, h):
    """slide9 상단 띠 도식: '라벨 없음 → 자기지도 군집 → 0.72 컷으로 같은 불량 보존'을
    라벨 없는 점들이 군집으로 모이는 3단계 미니 개념도로 시각화(글 위주 슬라이드 보강).
    이미지 가공 없이 PowerPoint 도형만으로 그린다."""
    import math, random
    x, y, w, h = int(x), int(y), int(w), int(h)
    seg = (w - int(Inches(0.9))) // 3       # 3패널 + 화살표 2개 공간
    aw = int(Inches(0.45))
    pad = int(Inches(0.12))
    cap_h = int(Inches(0.26))
    pan_y = y + cap_h
    pan_h = h - cap_h
    caps = ["① Unlabeled points", "② Self-supervised clustering", "③ Similarity 0.72 cut / same-failure 유지"]
    cap_cols = [MUTED, NAVY, ACCENT]
    px0 = x
    centers_for = {
        0: None,  # scattered
        1: [(0.30, 0.40), (0.50, 0.62), (0.72, 0.38)],
        2: [(0.30, 0.40), (0.50, 0.62), (0.72, 0.38)],
    }
    clu_cols = [RGBColor(0x12, 0xB5, 0xB0), RGBColor(0x2F, 0x6F, 0xD6), RGBColor(0xF2, 0xB7, 0x05)]
    for pi in range(3):
        px = px0 + pi * (seg + aw)
        # 패널 카드
        _rect(slide, Emu(px), Emu(pan_y), Emu(seg), Emu(pan_h), RGBColor(0xF4, 0xF7, 0xFB),
              line=RGBColor(0xC4, 0xD0, 0xE2), line_w=Pt(1))
        _text(slide, Emu(px), Emu(y), Emu(seg), Emu(cap_h),
              [[(caps[pi], dict(size=10.5, bold=True, color=cap_cols[pi]))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ix, iy = px + pad, pan_y + pad
        iw, ih = seg - 2 * pad, pan_h - 2 * pad
        dot = int(Inches(0.075))
        rnd = random.Random(7 + pi)
        cents = centers_for[pi]
        if cents is None:
            # 무작위 회색 점
            for _ in range(22):
                dx = ix + int(rnd.uniform(0.08, 0.92) * iw)
                dy = iy + int(rnd.uniform(0.10, 0.90) * ih)
                _rect(slide, Emu(dx), Emu(dy), Emu(dot), Emu(dot),
                      RGBColor(0x9A, 0xA6, 0xB8), shape=MSO_SHAPE.OVAL)
        else:
            # 3개 군집으로 모인 점 + (패널3) 군집 경계 원
            for ci, (gcx, gcy) in enumerate(cents):
                col = clu_cols[ci]
                ccx = ix + int(gcx * iw); ccy = iy + int(gcy * ih)
                if pi == 2:
                    rr = int(Inches(0.32))
                    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(ccx - rr), Emu(ccy - rr),
                                                  Emu(2 * rr), Emu(2 * rr))
                    ring.fill.background(); ring.line.color.rgb = col
                    ring.line.width = Pt(1.4); ring.shadow.inherit = False
                for _ in range(7):
                    dx = ccx + int(rnd.uniform(-0.13, 0.13) * iw)
                    dy = ccy + int(rnd.uniform(-0.16, 0.16) * ih)
                    _rect(slide, Emu(dx), Emu(dy), Emu(dot), Emu(dot), col, shape=MSO_SHAPE.OVAL)
        if pi < 2:
            ay = pan_y + pan_h // 2 - int(Inches(0.13))
            _rect(slide, Emu(px + seg + int(Inches(0.04))), Emu(ay),
                  Emu(aw - int(Inches(0.1))), Emu(int(Inches(0.26))), ACCENT,
                  shape=MSO_SHAPE.RIGHT_ARROW)


def _twotrack_diagram(slide, x, y, w, h):
    """slide10 상단 띠 도식: '운영(양산 중) vs 모델(검증·배포 대기)' 2트랙 비교.
    두 가로 트랙을 위/아래로 두고 운영=teal 채움(가동), 모델=점선 외곽(대기)로 대비.
    이미지 가공 없이 PowerPoint 도형만 사용."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    lblw = int(Inches(1.7))
    track_x = x + lblw
    track_w = w - lblw - int(Inches(0.2))
    # 트랙 높이·간격을 배정된 밴드 높이(h)에 맞춰 산출한다 — 이전엔 0.5"/0.30" 고정이라 밴드를
    # 줄이면(dgh 1.16→0.98) 두 트랙이 밴드 밖(아래 컬럼 헤더 위)으로 삐져나와 겹쳤다. h 안에
    # 두 트랙 + 한 간격이 모두 들어가도록 비례 분배해 도식이 항상 밴드 안에 머물게 한다.
    top_pad = int(Inches(0.05))
    avail = h - 2 * top_pad
    th = int((avail - int(Inches(0.26))) // 2)   # 두 트랙, 사이 간격 0.26"
    th = min(th, int(Inches(0.5)))
    gap = int(Inches(0.26))
    ty1 = y + top_pad
    ty2 = ty1 + th + gap
    # 운영 트랙(가동 중) — teal 채움 + '양산 중' 단계들
    _text(slide, Emu(x), Emu(ty1), Emu(lblw - int(Inches(0.1))), Emu(th),
          [[("Production operation", dict(size=11, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)
    ops = ["Processing flow", "Production viewer", "20K wafer/day"]
    n = len(ops); sg = int(Inches(0.14))
    bw = (track_w - sg * (n - 1)) // n
    for k, t in enumerate(ops):
        bx = track_x + k * (bw + sg)
        _rect(slide, Emu(bx), Emu(ty1), Emu(bw), Emu(th), ACCENT,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, Emu(bx), Emu(ty1), Emu(bw), Emu(th),
              [[(t, dict(size=10.5, bold=True, color=WHITE))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Emu(track_x + track_w + int(Inches(0.02))), Emu(ty1 + th // 2 - int(Inches(0.13))),
          Emu(int(Inches(0.16))), Emu(int(Inches(0.26))), ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
    # 모델 트랙(검증·배포 대기) — 회색 점선 외곽 + 마지막 단계 'GPU 대기'
    _text(slide, Emu(x), Emu(ty2), Emu(lblw - int(Inches(0.1))), Emu(th),
          [[("Model validation / deploy waiting", dict(size=11, bold=True, color=MUTED))]], anchor=MSO_ANCHOR.MIDDLE)
    mds = ["Validation F1 0.95", "Unknown 후보 13→7", "GPU server → 2026.9"]
    for k, t in enumerate(mds):
        bx = track_x + k * (bw + sg)
        last = (k == len(mds) - 1)
        # '예정/대기' 칩을 점선 크림+골드 아웃라이어 대신 팔레트 내 변형으로 통일:
        # 옅은 teal 틴트 채움 + teal 점선 외곽으로 일관화하고, 마지막 '배포 예정' 칩만 한 단계
        # 진한 teal 틴트로 위계를 준다(운영 트랙 solid teal 과 한 팔레트로 묶임 — slide10 low).
        fill = RGBColor(0xF4, 0xF6, 0xF8) if not last else RGBColor(0xE6, 0xEA, 0xEF)
        line_c = ACCENT
        sp = _rect(slide, Emu(bx), Emu(ty2), Emu(bw), Emu(th), fill,
                   line=line_c, line_w=Pt(1.4), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 점선 외곽(대기 상태 표현) — 색은 teal 로 통일
        ln = sp.line._get_or_add_ln()
        pd = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(pd)
        _text(slide, Emu(bx), Emu(ty2), Emu(bw), Emu(th),
              [[(t, dict(size=10, bold=True, color=NAVY))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


_TWO_COL_DIAGRAMS = {"cluster3": _cluster3_diagram, "twotrack": _twotrack_diagram}


def s_two_col(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    cols = d["cols"]
    # 상단 도식 띠(diagram): 글 위주 2단 슬라이드(P1 Unknown·운영vs모델)에 핵심 논리를 한눈에
    # 보여주는 미니 도식을 제목 아래에 얹고, 두 컬럼을 그만큼 아래로 내려 압축한다(텍스트 밀도↓).
    y = Inches(1.95); ch = Inches(4.55)
    if d.get("diagram") in _TWO_COL_DIAGRAMS:
        # 상단 도식 띠(flow)와 그 아래 네이비 본문 박스가 거의 붙어 답답하던 인상(slide10 low)을
        # 줄이려 도식을 살짝 위로(1.66→1.60) 올리고 도식↔본문 박스 사이 수직 여백을 키운다
        # (0.16→0.34). 본문 박스는 그만큼 짧아지지만 줄간격으로 흡수한다.
        # 도식 띠 높이·아래 여백을 줄여(1.16→0.98, 0.34→0.22) 아래 두 컬럼에 세로 공간을 더 준다.
        # 3불릿+각주(4문단)가 들어가는 컬럼이 짧아 본문이 칩 밴드로 흘러내리던 충돌(slide10 high)을
        # 구조적으로 완화 — 컬럼 높이가 ≈0.30" 늘어 본문이 안전선 안에 수렴한다.
        dgy = Inches(1.56); dgh = Inches(1.10)
        _TWO_COL_DIAGRAMS[d["diagram"]](slide, Inches(0.7), dgy, Inches(11.93), dgh)
        y = Emu(int(dgy) + int(dgh) + int(Inches(0.22)))
        ch = Emu(int(Inches(6.62)) - int(y))
    x0 = Inches(0.7); cw = Inches(5.85); gap = Inches(0.3)
    for i, c in enumerate(cols[:2]):
        x = Emu(int(x0) + i*(int(cw)+int(gap)))
        head_color = c.get("color", NAVY)
        # 헤더 띠 높이를 0.62→0.70 으로 키워 네이비 헤더 텍스트 위/아래 내부 패딩을 확보(slide9·10
        # medium — 헤더 텍스트가 카드 천장에 붙어 보이던 인상 제거). 텍스트는 MIDDLE 정렬 유지.
        head_h = Inches(0.70)
        _rect(slide, x, y, cw, head_h, head_color)
        _text(slide, x+Inches(0.25), y, Emu(int(cw)-int(Inches(0.4))), head_h,
              [[(c["head"], dict(size=15.5, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
        _rect(slide, x, Emu(int(y)+int(head_h)), cw, Emu(int(ch)-int(head_h)), PANEL)
        chips = c.get("chips")
        # 칩(미니 스탯)이 있으면 하단에 배치하고 본문 영역을 그만큼 줄여 빈 공간을 메움.
        # 본문↔칩 사이 안전 간격을 더 넓혀(2.45") 마지막 줄 descender가 칩 상단 teal 경계에
        # 붙어 잘리는 충돌을 확실히 제거(slide9 '감/불량/동' 잘림 방지). 칩이 있을 땐 본문을
        # TOP 정렬해 아래로 흘러내리지 않게 한다.
        # 본문↔칩 안전 간격을 2.45→2.55" 로 더 넓히고(body_h 축소) 본문 폰트를 12.5→12 로 한 단계
        # 줄여, 4~5개 긴 불릿의 마지막 줄 descender 가 칩 밴드 teal 상단 경계에 붙어 잘리던 충돌을
        # 확실히 제거(slide9 '군집 보존/깨짐' 잘림 방지). 줄간격은 1.03 으로 미세 조정해 행수 흡수.
        # 상단 도식 띠가 있으면 컬럼이 짧아지므로 본문↔칩 안전 간격을 더 줄여(2.25→1.78")
        # 짧아진 컬럼에서도 4문단(불릿 3 + 각주 1) 본문이 충분한 높이를 확보해 칩 밴드 위로
        # 넘쳐 잘리지 않게 한다(slide10 우측 '사내 인정 3건' 마지막 줄이 칩에 겹치던 충돌 제거).
        # 불릿 수를 줄인 뒤(밀도↓) 본문↔칩 사이 여백을 키우고(chip_reserve 1.66→2.02) 줄간격을
        # 1.03→1.13 으로 넓혀 임원 대상 '한눈에' 가독성을 확보(slide9·10 density medium 해결).
        # 마지막 각주(dash) 줄이 칩 밴드에 겹치지 않게 reserve 를 충분히 확보한다.
        # 본문 시작점을 헤더 띠(0.70) 아래로 충분히 띄우고(0.20" 상단 패딩), 본문 블록을 헤더 하단~
        # 칩 밴드 상단 직전 사이로 한정해 MIDDLE 정렬한다 → 본문이 위로 몰리고 하단에 빈 띠가 남던
        # 수직 쏠림을 제거해 무게중심을 카드 가운데로(slide9·10 medium). 칩 밴드 top(y+ch-0.92)
        # 위 0.16" 안전 간격까지만 본문을 두어 마지막 줄 descender 가 칩에 닿지 않게 한다.
        body_y = Emu(int(y)+int(Inches(0.70))+int(Inches(0.20)))
        chip_top = int(y)+int(ch)-int(Inches(0.92))
        # 칩 밴드 위 안전 간격을 0.16→0.24" 로 키워, 컬럼별 본문 길이가 달라도(특히 우측이 더 긴
        # slide10) 마지막 회색 각주 줄이 칩 카드 상단 teal 경계에 겹쳐 잘리지 않게 한다(high overlap).
        body_bottom = Emu(chip_top-int(Inches(0.24))) if chips else Emu(int(y)+int(ch)-int(Inches(0.7)))
        body_h = Emu(int(body_bottom)-int(body_y))
        # 본문 폰트: 상단 도식 띠가 있어 컬럼이 짧아진 경우(slide9·10)는 12.5pt·줄간격 1.04 로 한
        # 단계 낮춰 3불릿+각주 4문단이 칩 밴드 위 안전선(0.24") 안에 들어오게 한다(slide10 우측
        # 본문이 칩 카드로 흘러내려 겹치던 high overlap 제거). 도식이 없으면 13pt 유지(프리미엄 톤).
        has_diag = d.get("diagram") in _TWO_COL_DIAGRAMS
        body_sz = 12.5 if (chips and has_diag) else 13
        body_ls = 1.04 if (chips and has_diag) else 1.1
        # 칩 컬럼은 TOP 정렬: MIDDLE 정렬은 본문이 길어지면 블록이 아래로 중앙 이동해 마지막 줄이
        # 칩 밴드로 흘러내린다. TOP 이면 헤더 바로 아래부터 쌓여 body_bottom(칩 위 안전선) 안에서 끝남.
        _bullets_tf(slide, x+Inches(0.25), body_y, Emu(int(cw)-int(Inches(0.5))),
                    body_h, c.get("bullets", []), size=body_sz, gap=4,
                    anchor=(MSO_ANCHOR.TOP if chips else MSO_ANCHOR.MIDDLE), line_spacing=body_ls)
        if chips:
            chips = chips[:3]; ncc = len(chips)
            cgap = Inches(0.16)
            ccw = Emu(int((int(cw)-int(Inches(0.5))-int(cgap)*(ncc-1))//ncc))
            cy = y + Emu(int(ch)) - Inches(0.92)
            for k, ch_d in enumerate(chips):
                cx = Emu(int(x)+int(Inches(0.25))+k*(int(ccw)+int(cgap)))
                _rect(slide, cx, cy, ccw, Inches(0.75), WHITE, line=LINE)
                _rect(slide, cx, cy, ccw, Inches(0.07), ACCENT)
                _text(slide, cx, cy+Inches(0.12), ccw, Inches(0.36),
                      [[(ch_d["big"], dict(size=16, bold=True, color=NAVY))]],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _text(slide, cx, cy+Inches(0.46), ccw, Inches(0.28),
                      [[(ch_d.get("label", ""), dict(size=9, color=MUTED))]],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def _bytes_compare_diagram(slide, x, y, w, h):
    """Palette-indexed PNG(논문 2.1, Fig 2) — paper-figure 스타일: 단색 + monospace + 얇은 선.
    24-bit RGB(픽셀마다 3 byte) vs 8-bit 팔레트(PLTE 표 + 인덱스 1 byte) → 약 75% 축소, 무손실.
    컬러 RGB 박스/틸 바 없이 monospace 튜플과 grayscale 막대로 학술 figure처럼 표현."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    def I(v): return int(Inches(v))
    MONO = "Consolas"
    INKD = RGBColor(0x1A, 0x1F, 0x2E)
    LBL = RGBColor(0x6B, 0x72, 0x80)
    LN = RGBColor(0xC4, 0xCC, 0xD8)
    _text(slide, Emu(x), Emu(y + I(0.07)), Emu(w), Emu(I(0.24)),
          [[("Palette-indexed PNG", dict(size=11, bold=True, color=NAVY)),
            ("    24-bit RGB → 8-bit 인덱스 (무손실)", dict(size=9, color=LBL))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    # RGB PNG 행 (monospace 튜플)
    r1 = y + I(0.48)
    _text(slide, Emu(x + I(0.26)), Emu(r1), Emu(w - I(0.5)), Emu(I(0.20)),
          [[("RGB PNG  (픽셀마다 3 byte):", dict(size=9.5, color=NAVY, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Emu(x + I(0.40)), Emu(r1 + I(0.22)), Emu(w - I(0.7)), Emu(I(0.20)),
          [[("(123,54,24) (123,54,24) (90,12,40) (123,54,24) …", dict(size=9, color=INKD, name=MONO))]],
          anchor=MSO_ANCHOR.MIDDLE)
    # Palette PNG 행
    r2 = r1 + I(0.56)
    _text(slide, Emu(x + I(0.26)), Emu(r2), Emu(w - I(0.5)), Emu(I(0.20)),
          [[("Palette PNG  (PLTE 표 + 인덱스 1 byte):", dict(size=9.5, color=NAVY, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Emu(x + I(0.40)), Emu(r2 + I(0.22)), Emu(w - I(0.7)), Emu(I(0.20)),
          [[("PLTE[3]=(123,54,24) … 32색   →   픽셀 = ", dict(size=9, color=INKD, name=MONO)),
            ("3 3 5 3 7 3 …", dict(size=9, bold=True, color=NAVY, name=MONO))]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Emu(x + I(0.26)), Emu(r2 + I(0.46)), Emu(w - I(0.5)), Emu(I(0.18)),
          [[("색 32종뿐 → 반복색은 PLTE에 1회만, 픽셀은 번호만 저장 (PLTE 교체로 색 scheme 즉시 변경)",
             dict(size=8.5, color=LBL))]], anchor=MSO_ANCHOR.MIDDLE)
    # grayscale 파일크기 막대 (틸 금지)
    bary = r2 + I(1.02)
    _text(slide, Emu(x + I(0.26)), Emu(bary - I(0.26)), Emu(w - I(0.5)), Emu(I(0.20)),
          [[("같은 무손실 PNG 파일 크기", dict(size=9.5, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)
    bx = x + I(1.15); bmax = (x + w - I(0.95)) - bx; bh = I(0.24)
    _text(slide, Emu(x + I(0.26)), Emu(bary), Emu(I(0.85)), Emu(bh),
          [[("RGB", dict(size=9, color=LBL))]], anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Emu(bx), Emu(bary), Emu(bmax), Emu(bh), RGBColor(0xC6, 0xD0, 0xDF), line=LN, line_w=Pt(0.5))
    _text(slide, Emu(bx + bmax - I(0.62)), Emu(bary), Emu(I(0.55)), Emu(bh),
          [[("100%", dict(size=9, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bary2 = bary + bh + I(0.12)
    _text(slide, Emu(x + I(0.26)), Emu(bary2), Emu(I(0.85)), Emu(bh),
          [[("팔레트", dict(size=9, color=LBL))]], anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Emu(bx), Emu(bary2), Emu(int(bmax * 0.25)), Emu(bh), RGBColor(0x48, 0x52, 0x60))
    _text(slide, Emu(bx + int(bmax * 0.25) + I(0.10)), Emu(bary2), Emu(I(2.0)), Emu(bh),
          [[("25%   (약 75% 축소, 무손실)", dict(size=9, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)


def _hex_to_grade_diagram(slide, x, y, w, h):
    """Cython 파서(논문 2.1) — paper-figure 스타일: 단색(네이비/그레이) + monospace + 얇은 선.
    raw hex payload(memoryview stride) → 256-entry 정적 C 배열 LUT[b] O(1) 조회 → chip tile
    array(Grade 격자), 순수 Python 대비 ~100x. 장식(틸 띠/컬러 셀) 없이 학술 figure로 표현."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    def I(v): return int(Inches(v))
    MONO = "Consolas"
    INKD = RGBColor(0x1A, 0x1F, 0x2E)
    LBL = RGBColor(0x6B, 0x72, 0x80)
    LN = RGBColor(0xC4, 0xCC, 0xD8)
    GREY = [RGBColor(0xFB, 0xFC, 0xFD), RGBColor(0xEC, 0xEF, 0xF3), RGBColor(0xD9, 0xDF, 0xE7),
            RGBColor(0xC0, 0xC8, 0xD3), RGBColor(0xA2, 0xAC, 0xBA), RGBColor(0x84, 0x8F, 0x9F),
            RGBColor(0x66, 0x70, 0x80), RGBColor(0x48, 0x52, 0x60)]
    # 제목 (네이비, 부제 그레이 — 큰 틸 제목 금지)
    _text(slide, Emu(x), Emu(y + I(0.07)), Emu(w), Emu(I(0.24)),
          [[("Cython hex→Grade 파서", dict(size=11, bold=True, color=NAVY)),
            ("    memoryview stride with 256-entry LUT", dict(size=9, color=LBL))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    # raw hex payload — monospace 셀(흰 칸 + 얇은 선)
    rl_y = y + I(0.42)
    _text(slide, Emu(x + I(0.26)), Emu(rl_y), Emu(w - I(0.5)), Emu(I(0.18)),
          [[("raw payload (memoryview stride 접근):", dict(size=9, color=LBL))]], anchor=MSO_ANCHOR.MIDDLE)
    hxs = ["3A", "0C", "C1", "0E", "A7", "5B", "0C", "FF"]
    hcw = I(0.40); hg = I(0.04); hh = I(0.28); hx0 = x + I(0.26); hy = rl_y + I(0.22)
    for i, v in enumerate(hxs):
        hxx = hx0 + i * (hcw + hg)
        _rect(slide, Emu(hxx), Emu(hy), Emu(hcw), Emu(hh), WHITE, line=LN, line_w=Pt(0.75))
        _text(slide, Emu(hxx), Emu(hy), Emu(hcw), Emu(hh),
              [[(v, dict(size=10, color=INKD, name=MONO))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Emu(hx0 + len(hxs) * (hcw + hg) + I(0.06)), Emu(hy), Emu(I(1.2)), Emu(hh),
          [[("1 byte/픽셀", dict(size=8.5, color=LBL))]], anchor=MSO_ANCHOR.MIDDLE)
    # LUT 라인 — monospace
    lut_y = hy + hh + I(0.12)
    _text(slide, Emu(x + I(0.26)), Emu(lut_y), Emu(w - I(0.5)), Emu(I(0.20)),
          [[("LUT[256] 정적 C 배열:  ", dict(size=9.5, color=NAVY, bold=True)),
            ("0x0C→3   0x1A→5   0xFF→7   …  (사전 계산, O(1))", dict(size=9.5, color=INKD, name=MONO))]],
          anchor=MSO_ANCHOR.MIDDLE)
    # 하단: 좌 Python/Cython 대비 박스, 우 chip tile array
    low_y = lut_y + I(0.34)
    cbx = x + I(0.26); cbw = I(3.10); cbh = (y + h) - low_y - I(0.10)
    _rect(slide, Emu(cbx), Emu(low_y), Emu(cbw), Emu(cbh), RGBColor(0xFA, 0xFB, 0xFC), line=LN, line_w=Pt(0.9))
    _text(slide, Emu(cbx + I(0.14)), Emu(low_y + I(0.05)), Emu(cbw - I(0.28)), Emu(cbh - I(0.10)),
          [[("순수 Python  vs  Cython", dict(size=9, bold=True, color=NAVY))],
           [("Python:  g = f(b)    # 인터프리터", dict(size=9, color=INKD, name=MONO))],
           [("Cython:  g = LUT[b]  # C배열 컴파일", dict(size=9, color=INKD, name=MONO))],
           [("                     → 약 100배", dict(size=9.5, bold=True, color=NAVY, name=MONO))]],
          anchor=MSO_ANCHOR.MIDDLE)
    # chip tile array — grayscale Grade 격자(digit), 컬러 채움 없음
    gx0 = cbx + cbw + I(0.26)
    _text(slide, Emu(gx0), Emu(low_y), Emu((x + w) - gx0 - I(0.2)), Emu(I(0.18)),
          [[("chip tile array (Grade)", dict(size=9, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)
    gn = 5; gtop = low_y + I(0.22)
    gcell = min(I(0.30), (((y + h) - I(0.12)) - gtop) // gn, ((x + w) - I(0.2) - gx0) // gn)
    patt = [[0, 0, 1, 0, 0], [0, 2, 4, 2, 0], [1, 4, 6, 4, 1], [0, 2, 4, 2, 0], [0, 0, 1, 0, 0]]
    for r in range(gn):
        for c in range(gn):
            gv = patt[r][c]; gx = gx0 + c * gcell; gyy = gtop + r * gcell
            _rect(slide, Emu(gx), Emu(gyy), Emu(gcell - I(0.012)), Emu(gcell - I(0.012)),
                  GREY[gv], line=LN, line_w=Pt(0.5))
            tc = WHITE if gv >= 5 else INKD
            _text(slide, Emu(gx), Emu(gyy), Emu(gcell - I(0.012)), Emu(gcell - I(0.012)),
                  [[(str(gv), dict(size=8.5, color=tc, name=MONO))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _numba_composite_diagram(slide, x, y, w, h):
    """다수 wafer를 픽셀 위치별로 누적해 composite hot-spot map을 만드는 Numba 병렬 합산을
    코드/이미지 캡처 대신 네이티브 도형으로 보여준다. 여러 wafer → 픽셀별 누적(@njit prange)
    → hot-spot 강조 composite. 라벨은 하단 라벨행에 둬 요소와 겹치지 않게 한다."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    def I(v): return int(Inches(v))
    cap_reserve = I(0.04); avail = h - cap_reserve
    title_h = I(0.20)
    _text(slide, Emu(x), Emu(y + I(0.04)), Emu(w), Emu(title_h),
          [[("다수 wafer를 픽셀 위치별로 누적", dict(size=10.5, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    callout_h = I(0.28)
    band_top = y + I(0.04) + title_h + I(0.03)
    band_bot = y + avail - callout_h - I(0.05)
    lbl_h = I(0.18); lbl_y = band_bot - lbl_h
    elem_top = band_top; elem_h = (band_bot - lbl_h) - elem_top
    wafers_cx = x + I(1.25); arr_cx = x + I(2.85); comp_cx = x + I(4.35)
    mini = min(I(0.62), elem_h)
    offs = [(-I(0.13), -I(0.09)), (0, I(0.02)), (I(0.13), I(0.13))]
    hot = [(0.60, 0.40), (0.46, 0.54), (0.50, 0.50)]
    for k, (ox, oy) in enumerate(offs):
        wx = wafers_cx - mini // 2 + ox; wy = elem_top + (elem_h - mini) // 2 + oy
        sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(wx), Emu(wy), Emu(mini), Emu(mini))
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xDC, 0xE6, 0xF0)
        sp.line.color.rgb = RGBColor(0xB9, 0xC6, 0xDB); sp.line.width = Pt(1); sp.shadow.inherit = False
        fx, fy = hot[k]; dd = I(0.10)
        _rect(slide, Emu(wx + int(fx * mini) - dd // 2), Emu(wy + int(fy * mini) - dd // 2),
              Emu(dd), Emu(dd), RGBColor(0xE0, 0x68, 0x3A))
    _text(slide, Emu(wafers_cx - I(0.7)), Emu(lbl_y), Emu(I(1.4)), Emu(lbl_h),
          [[("여러 wafer", dict(size=9, color=MUTED))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    arr_w = I(0.95); ay = elem_top + elem_h // 2 - I(0.12)
    _rect(slide, Emu(arr_cx - arr_w // 2), Emu(ay), Emu(arr_w), Emu(I(0.24)), ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
    _text(slide, Emu(arr_cx - I(1.0)), Emu(lbl_y), Emu(I(2.0)), Emu(lbl_h),
          [[("픽셀별 누적 (@njit prange)", dict(size=8.5, bold=True, color=ACCENT))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    comp = min(I(0.80), elem_h); compx = comp_cx - comp // 2; compy = elem_top + (elem_h - comp) // 2
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(compx), Emu(compy), Emu(comp), Emu(comp))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xEA, 0xF2, 0xF7)
    sp.line.color.rgb = RGBColor(0xB9, 0xC6, 0xDB); sp.line.width = Pt(1.2); sp.shadow.inherit = False
    hr = I(0.28)
    hp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(compx + comp // 2 - hr // 2),
                                Emu(compy + comp // 2 - hr // 2), Emu(hr), Emu(hr))
    hp.fill.solid(); hp.fill.fore_color.rgb = RGBColor(0xCC, 0x33, 0x28)
    hp.line.fill.background(); hp.shadow.inherit = False
    _text(slide, Emu(comp_cx - I(1.0)), Emu(lbl_y), Emu(I(2.0)), Emu(lbl_h),
          [[("hot-spot composite", dict(size=9, color=MUTED))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cby = y + avail - callout_h
    _rect(slide, Emu(x + I(0.3)), Emu(cby), Emu(w - I(0.6)), Emu(callout_h),
          RGBColor(0xF4, 0xF6, 0xF8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _rect(slide, Emu(x + I(0.3)), Emu(cby), Emu(I(0.09)), Emu(callout_h), ACCENT)
    _text(slide, Emu(x + I(0.46)), Emu(cby), Emu(w - I(0.76)), Emu(callout_h),
          [[("반복 픽셀 합산을 기계어로 컴파일 + 병렬 ", dict(size=9.5, bold=True, color=NAVY)),
            ("→ 대량 wafer 합성 가속", dict(size=9.5, bold=True, color=ACCENT))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _pyvips_stream_diagram(slide, x, y, w, h):
    """대용량 이미지를 통째로 적재하지 않고 타일 단위로 필요한 부분만 streaming 디코드하는
    pyvips 동작을 네이티브 도형으로 보여준다. 타일 격자(일부만 로드) → 운영 viewer.
    라벨은 하단 라벨행에 둔다."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    def I(v): return int(Inches(v))
    cap_reserve = I(0.04); avail = h - cap_reserve
    title_h = I(0.20)
    _text(slide, Emu(x), Emu(y + I(0.04)), Emu(w), Emu(title_h),
          [[("대용량 이미지를 타일 단위로 streaming 로드", dict(size=10.5, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    callout_h = I(0.28)
    band_top = y + I(0.04) + title_h + I(0.03)
    band_bot = y + avail - callout_h - I(0.05)
    lbl_h = I(0.18); lbl_y = band_bot - lbl_h
    elem_top = band_top; elem_h = (band_bot - lbl_h) - elem_top
    img_cx = x + I(1.25); arr_cx = x + I(2.85); v_cx = x + I(4.40)
    big = min(I(0.86), elem_h); bx = img_cx - big // 2; by = elem_top + (elem_h - big) // 2
    _rect(slide, Emu(bx), Emu(by), Emu(big), Emu(big), RGBColor(0xEE, 0xF1, 0xF6),
          line=RGBColor(0xB9, 0xC6, 0xDB), line_w=Pt(1.2))
    n = 4; ts = big // n
    loaded = {(1, 1), (1, 2), (2, 1)}
    for r in range(n):
        for c in range(n):
            tx = bx + c * ts; ty = by + r * ts
            if (r, c) in loaded:
                _rect(slide, Emu(tx), Emu(ty), Emu(ts - I(0.02)), Emu(ts - I(0.02)),
                      RGBColor(0xE6, 0xEA, 0xEF), line=ACCENT, line_w=Pt(1))
            else:
                sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(tx), Emu(ty),
                                            Emu(ts - I(0.02)), Emu(ts - I(0.02)))
                sp.fill.background(); sp.line.color.rgb = RGBColor(0xD7, 0xDE, 0xEA)
                sp.line.width = Pt(0.75); sp.shadow.inherit = False
    _text(slide, Emu(img_cx - I(1.0)), Emu(lbl_y), Emu(I(2.0)), Emu(lbl_h),
          [[("필요한 타일만 디코드", dict(size=9, bold=True, color=ACCENT))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    arr_w = I(0.95); ay = elem_top + elem_h // 2 - I(0.12)
    _rect(slide, Emu(arr_cx - arr_w // 2), Emu(ay), Emu(arr_w), Emu(I(0.24)), ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
    _text(slide, Emu(arr_cx - I(0.85)), Emu(lbl_y), Emu(I(1.7)), Emu(lbl_h),
          [[("on-demand 디코드", dict(size=8.5, bold=True, color=MUTED))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    vw = I(1.55); vh = min(I(0.78), elem_h); vx = v_cx - vw // 2; vy = elem_top + (elem_h - vh) // 2
    _rect(slide, Emu(vx), Emu(vy), Emu(vw), Emu(vh), RGBColor(0xEA, 0xF2, 0xF7),
          line=ACCENT, line_w=Pt(1.4), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(slide, Emu(vx), Emu(vy), Emu(vw), Emu(vh),
          [[("운영 viewer", dict(size=11, bold=True, color=NAVY))],
           [("대량 이미지 조회", dict(size=8.5, color=MUTED))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cby = y + avail - callout_h
    _rect(slide, Emu(x + I(0.3)), Emu(cby), Emu(w - I(0.6)), Emu(callout_h),
          RGBColor(0xF4, 0xF6, 0xF8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _rect(slide, Emu(x + I(0.3)), Emu(cby), Emu(I(0.09)), Emu(callout_h), ACCENT)
    _text(slide, Emu(x + I(0.46)), Emu(cby), Emu(w - I(0.76)), Emu(callout_h),
          [[("전체 메모리 적재 대신 타일 streaming ", dict(size=9.5, bold=True, color=NAVY)),
            ("→ viewer 로드 경량화", dict(size=9.5, bold=True, color=ACCENT))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _episode_trend_diagram(slide, x, y, w, h):
    """P3 normal trend generator: split normal baseline and noise-family scatter panels."""
    import random
    rnd = random.Random(7)
    x, y, w, h = int(x), int(y), int(w), int(h)
    def I(v): return int(Inches(v))
    BLUE = RGBColor(0x2D, 0x63, 0xB8); GRAY = RGBColor(0xA7, 0xAD, 0xB6)
    RED = RGBColor(0xCC, 0x33, 0x28); NV = RGBColor(0x0F, 0x1E, 0x3D)
    TL = RGBColor(0x12, 0xB5, 0xB0); MU = RGBColor(0x6B, 0x72, 0x80)

    def dot(cx, cy, d, color):
        sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - d / 2)), Emu(int(cy - d / 2)), Emu(int(d)), Emu(int(d)))
        sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False

    def axes(ax, ay, aw, ah):
        _rect(slide, Emu(ax), Emu(ay + ah - I(0.015)), Emu(aw), Emu(I(0.015)), RGBColor(0xD4, 0xDA, 0xE3))
        _rect(slide, Emu(ax), Emu(ay), Emu(I(0.015)), Emu(ah), RGBColor(0xD4, 0xDA, 0xE3))

    gap = I(0.24)
    left_w = int((w - gap) * 0.72)
    right_w = w - gap - left_w
    left_x = x
    right_x = x + left_w + gap
    head_h = I(0.34)
    legend_h = I(0.20)
    body_y = y + head_h + legend_h + I(0.08)
    body_h = h - head_h - legend_h - I(0.10)

    # -- (a) normal baseline: one sequential scatter plot with episode density bands --
    _text(slide, Emu(left_x), Emu(y), Emu(left_w), Emu(head_h),
          [[("(a) Normal baseline  -  sequential episodes", dict(size=12.6, bold=True, color=NV))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    _text(slide, Emu(left_x), Emu(y + head_h), Emu(left_w), Emu(legend_h),
          [[("gray = fleet/reference     blue = sampled observations", dict(size=9.8, bold=True, color=MU))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)

    ax = left_x + I(0.18); ay = body_y + I(0.04)
    aw = left_w - I(0.36); ah = body_h - I(0.12)
    _rect(slide, Emu(ax), Emu(ay), Emu(aw), Emu(ah), RGBColor(0xFC, 0xFD, 0xFE),
          line=RGBColor(0xD7, 0xDE, 0xEA), line_w=Pt(1))
    axes(ax + I(0.18), ay + I(0.20), aw - I(0.34), ah - I(0.42))
    plot_x = ax + I(0.24); plot_y = ay + I(0.28)
    plot_w = aw - I(0.46); plot_h = ah - I(0.58)
    bands = [
        ("Dense", 74, RGBColor(0xEA, 0xF0, 0xF8)),
        ("Sparse", 24, RGBColor(0xFB, 0xF4, 0xE5)),
        ("Missing", 0, RGBColor(0xFB, 0xEC, 0xEA)),
        ("Sparse", 22, RGBColor(0xFB, 0xF4, 0xE5)),
        ("Dense", 70, RGBColor(0xEA, 0xF0, 0xF8)),
    ]
    band_w = plot_w // len(bands)
    for bi, (name, active_n, fill) in enumerate(bands):
        bx = plot_x + bi * band_w
        bw = band_w if bi < len(bands) - 1 else plot_w - band_w * bi
        _rect(slide, Emu(bx), Emu(plot_y), Emu(bw), Emu(plot_h), fill,
              line=RGBColor(0xE1, 0xE6, 0xEF), line_w=Pt(0.6))
        _text(slide, Emu(bx + I(0.02)), Emu(plot_y + I(0.02)), Emu(bw - I(0.04)), Emu(I(0.18)),
              [[(name, dict(size=9.4, bold=True, color=NV if name == "Dense" else MU))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        # Dense gray background makes the fleet/reference distribution visible.
        for _ in range(42):
            tx = bx + I(0.08) + rnd.random() * max(1, bw - I(0.16))
            ty = plot_y + I(0.34) + rnd.random() * max(1, plot_h - I(0.50))
            dot(tx, ty, I(0.036), GRAY)
        if active_n == 0:
            _text(slide, Emu(bx), Emu(plot_y + plot_h // 2 - I(0.10)), Emu(bw), Emu(I(0.22)),
                  [[("missing", dict(size=10.2, bold=True, color=RED))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            continue
        center = plot_y + int(plot_h * (0.47 + 0.06 * (bi % 2)))
        spread = 0.22 if name == "Dense" else 0.12
        for _ in range(active_n):
            tx = bx + I(0.10) + rnd.random() * max(1, bw - I(0.20))
            ty = center + int(rnd.uniform(-spread, spread) * plot_h)
            dot(tx, ty, I(0.052), BLUE)

    # -- (b) measurement noise: separate scatter panels, one noise type per chart --
    _text(slide, Emu(right_x), Emu(y), Emu(right_w), Emu(head_h),
          [[("(b) Measurement noise", dict(size=12.0, bold=True, color=NV))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    _text(slide, Emu(right_x), Emu(y + head_h), Emu(right_w), Emu(legend_h),
          [[("gray = fleet     color = sampled εₜ", dict(size=9.1, bold=True, color=MU))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)

    titles = ["Gaussian iid", "Correlated AR(1)", "Laplacian"]
    formulas = ["εₜ ~ N(0, σ²)", "εₜ = ρεₜ₋₁ + ηₜ", "εₜ ~ Laplace(0, b)"]
    pcol = [BLUE, TL, NV]
    panel_gap = I(0.10)
    panel_h = (body_h - 2 * panel_gap) // 3
    for pi, title in enumerate(titles):
        py = body_y + pi * (panel_h + panel_gap)
        _rect(slide, Emu(right_x), Emu(py), Emu(right_w), Emu(panel_h), RGBColor(0xFC, 0xFD, 0xFE),
              line=RGBColor(0xD7, 0xDE, 0xEA), line_w=Pt(1))
        _text(slide, Emu(right_x + I(0.12)), Emu(py + I(0.03)), Emu(right_w - I(0.24)), Emu(I(0.18)),
              [[(title, dict(size=10.5, bold=True, color=NV))]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        _text(slide, Emu(right_x + I(0.12)), Emu(py + I(0.22)), Emu(right_w - I(0.24)), Emu(I(0.18)),
              [[(formulas[pi], dict(size=9.0, bold=True, color=MU, name="Cambria Math"))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        px0 = right_x + I(0.22); px1 = right_x + right_w - I(0.22)
        cy = py + int(panel_h * 0.66); half = int(panel_h * 0.22)
        axes(px0, py + I(0.45), px1 - px0, panel_h - I(0.58))
        for _ in range(32):
            tx = px0 + rnd.random() * max(1, px1 - px0)
            ty = cy + int(rnd.uniform(-0.80, 0.80) * half)
            dot(tx, ty, I(0.036), GRAY)
        if pi == 0:
            vals = [rnd.gauss(0, 0.38) for _ in range(18)]
        elif pi == 1:
            vals = []
            v = 0.0
            for _ in range(18):
                v = 0.82 * v + rnd.gauss(0, 0.18)
                vals.append(max(-0.85, min(0.85, v)))
        else:
            vals = [rnd.gauss(0, 0.16) for _ in range(14)] + [0.92, -0.88, 0.74, -0.68]
        for k, v in enumerate(vals):
            tx = px0 + int((px1 - px0) * (k + 0.5) / len(vals))
            ty = cy - int(max(-0.95, min(0.95, v)) * half)
            dot(tx, ty, I(0.058), RED if (pi == 2 and abs(v) > 0.60) else pcol[pi])

def _p3_quad_diagram(slide, x, y, w, h):
    """P3 결과 4분할 네이티브 — 백본 bar / progression bar / smoothing 곡선 / color 전후.
    모든 데이터 글자 12pt 이상, 전부 편집 가능한 네이티브 객체."""
    import os as _os
    import random as _random
    import math as _math
    rnd = _random.Random(11)
    x, y, w, h = int(x), int(y), int(w), int(h)
    def I(v): return int(Inches(v))
    NV = RGBColor(0x0F, 0x1E, 0x3D); BL = RGBColor(0x3A, 0x6F, 0xD0); GR = RGBColor(0xAE, 0xB4, 0xBD)
    RD = RGBColor(0xCC, 0x33, 0x28); MU = RGBColor(0x6B, 0x72, 0x80)

    def rect(rx, ry, rw, rh, color, line=None, lw=None, shape=MSO_SHAPE.RECTANGLE):
        return _rect(slide, Emu(int(rx)), Emu(int(ry)), Emu(int(rw)), Emu(int(rh)), color, line=line, line_w=lw, shape=shape)

    def txt(tx, ty, tw, th, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False):
        _text(slide, Emu(int(tx)), Emu(int(ty)), Emu(int(tw)), Emu(int(th)), runs, align=align, anchor=anchor, wrap=wrap)

    def dot(cx, cy, d, color, shape=MSO_SHAPE.OVAL):
        sp = slide.shapes.add_shape(shape, Emu(int(cx - d / 2)), Emu(int(cy - d / 2)), Emu(int(d)), Emu(int(d)))
        sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False

    def line(x1, y1, x2, y2, color, wpt):
        cn = slide.shapes.add_connector(1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
        cn.line.color.rgb = color; cn.line.width = Pt(wpt); cn.shadow.inherit = False

    def polyline(pts, color, wpt):
        for i in range(len(pts) - 1):
            line(pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], color, wpt)

    margin = int(w * 0.08)
    x += margin; w -= 2 * margin
    gap = I(0.30)
    qw = (w - gap) // 2; qh = (h - gap) // 2
    quads = [(x, y), (x + qw + gap, y), (x, y + qh + gap), (x + qw + gap, y + qh + gap)]
    for (qx, qy) in quads:
        rect(qx, qy, qw, qh, WHITE, line=RGBColor(0xD7, 0xDE, 0xEA), lw=Pt(1.0))

    # ---- native F1 bar chart ----
    def barchart(qx, qy, title, names, f1, fn, fp, roles):
        txt(qx + I(0.1), qy + I(0.06), qw - I(0.2), I(0.30), [[(title, dict(size=13, bold=True, color=NV))]])
        pa_l = qx + I(0.16); pa_r = qx + qw - I(0.16)
        pa_top = qy + I(0.46); pa_bot = qy + qh - I(0.50)
        lo, hi = 0.9960, 0.99905
        n = len(names); slot = (pa_r - pa_l) / n
        bw = min(slot * 0.6, I(0.85))
        cmap = {"best": NV, "imp": BL, "base": GR}
        for i, (nm, v, fnv, fpv, rl) in enumerate(zip(names, f1, fn, fp, roles)):
            cx = pa_l + slot * (i + 0.5)
            bh = int((v - lo) / (hi - lo) * (pa_bot - pa_top))
            rect(cx - bw / 2, pa_bot - bh, bw, bh, cmap[rl])
            txt(cx - slot / 2, pa_bot - bh - I(0.27), slot, I(0.25), [[(("%.4f" % v), dict(size=12, bold=True, color=NV))]])
            txt(cx - slot / 2, pa_bot + I(0.03), slot, I(0.23), [[(nm, dict(size=12, bold=True, color=NV))]])
            txt(cx - slot / 2, pa_bot + I(0.25), slot, I(0.21), [[(("FN %d / FP %d" % (fnv, fpv)), dict(size=12, color=MU))]])
        line(pa_l, pa_bot, pa_r, pa_bot, RGBColor(0xC7, 0xCD, 0xD6), 1.1)

    # Q1 backbone
    qx, qy = quads[0]
    barchart(qx, qy, "Backbone sweep — F1 (5-seed)",
             ["dinov3", "base", "tiny", "swin", "maxvit", "effv2"],
             [0.9987, 0.9982, 0.9967, 0.9975, 0.9979, 0.9972],
             [0, 0, 1, 1, 0, 1], [2, 3, 4, 3, 3, 4],
             ["best", "imp", "base", "imp", "imp", "imp"])
    # Q2 progression
    qx, qy = quads[1]
    barchart(qx, qy, "baseline → BKM combined → best backbone",
             ["Baseline", "BKM\ncombined", "Best\nbackbone"],
             [0.9967, 0.9981, 0.9987], [1, 1, 0], [4, 3, 2],
             ["base", "imp", "best"])

    # Q3 smoothing window (native) — 단일 epoch 최고는 hunting spike(학습 덜 됨),
    # 후속 3-epoch median 이 더 높은 후속 plateau 를 best 로 선택
    qx, qy = quads[2]
    txt(qx + I(0.1), qy + I(0.06), qw - I(0.2), I(0.28),
        [[("Smoothing window — 후속 3-epoch median 으로 best 선택", dict(size=13, bold=True, color=NV))]])
    txt(qx + I(0.16), qy + I(0.33), qw - I(0.30), I(0.46),
        [[("단일 epoch 최고 F1은 초기 hunting spike(학습 덜 됨),", dict(size=12, color=NV))],
         [("후속 3-epoch median 이 더 높은 plateau 를 best 선택", dict(size=12, bold=True, color=NV))]],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True)
    sl = qx + I(0.20); sr = qx + qw - I(0.14)
    st = qy + I(0.84); sb = qy + qh - I(0.22)
    lo2, hi2 = 0.9930, 0.9983
    N = 22
    xs = [sl + (sr - sl) * k / (N - 1) for k in range(N)]
    raw = [0.9940 + 0.0030 * min(1.0, k / 13.0) + rnd.uniform(-0.0004, 0.0004) for k in range(N)]
    raw[6] = 0.99795; raw[5] = 0.9946; raw[7] = 0.9944
    for k in range(14, N):
        raw[k] = 0.99720 + rnd.uniform(-0.0003, 0.0003)
    def med3(a):
        b = sorted(a); m = len(b)
        return b[m // 2] if m % 2 else (b[m // 2 - 1] + b[m // 2]) / 2.0
    sm = [med3(raw[k:min(N, k + 3)]) for k in range(N - 2)]
    def ymap(v): return sb - (max(lo2, min(hi2, v)) - lo2) / (hi2 - lo2) * (sb - st)
    polyline([(xs[k], ymap(raw[k])) for k in range(N)], GR, 1.5)
    polyline([(xs[k], ymap(sm[k])) for k in range(N - 2)], BL, 2.8)
    rp = max(range(N), key=lambda k: raw[k])
    sp = max(range(N - 2), key=lambda k: sm[k])
    dot(xs[rp], ymap(raw[rp]), I(0.11), RD)
    dot(xs[sp], ymap(sm[sp]), I(0.15), NV, shape=MSO_SHAPE.DIAMOND)
    line(sl, sb, sr, sb, RGBColor(0xC7, 0xCD, 0xD6), 1.0)
    txt(qx + I(0.1), qy + qh - I(0.205), qw - I(0.2), I(0.19),
        [[("● 단일 최고(hunting)    ", dict(size=12, color=RD)), ("◆ median best → test 0.9987", dict(size=12, bold=True, color=NV))]],
        align=PP_ALIGN.CENTER)

    # Q4 color before/after (native pictures + native labels)
    qx, qy = quads[3]
    txt(qx + I(0.1), qy + I(0.06), qw - I(0.2), I(0.28), [[("Color 변경 전후 — target 색 대비↑ 분리도 향상", dict(size=13, bold=True, color=NV))]])
    imgs = [("p3r_color_baseline.png", "Before (파랑)"), ("p3r_color_c01.png", "After (빨강)")]
    isz = min(I(1.5), qh - I(0.95)); igap = I(0.5)
    total = isz * 2 + igap; ix0 = qx + (qw - total) // 2; iy = qy + I(0.42)
    for j, (fn_img, lab) in enumerate(imgs):
        ipath = _os.path.join(FIG_DIR, fn_img)
        ix = ix0 + j * (isz + igap)
        if _os.path.exists(ipath):
            slide.shapes.add_picture(ipath, Emu(int(ix)), Emu(int(iy)), Emu(int(isz)), Emu(int(isz)))
            rect(ix, iy, isz, isz, None, line=RGBColor(0xC7, 0xCD, 0xD6), lw=Pt(1.0)) if False else None
        txt(ix - I(0.15), iy + isz + I(0.04), isz + I(0.3), I(0.24), [[(lab, dict(size=12, bold=True, color=NV))]])



def _caption_runs(text, sz):
    """이미지 캡션을 '굵은 라벨 + 얇은 부연' 위계로 분리한다(첫 '·' 기준).
    예: '평균이동·치우침' → [평균이동(굵은 NAVY), '  치우침'(얇은 슬레이트, 소폭 작게)].
    '·' 가 없으면 통째로 굵은 라벨로 처리. 캡션 한 줄 안에서 임원이 핵심 라벨을 먼저 읽게 한다."""
    if "·" in text:
        lab, rest = text.split("·", 1)
        return [(lab.strip(), dict(size=sz, bold=True, color=NAVY)),
                ("  " + rest.strip(), dict(size=sz - 1, color=RGBColor(0x45, 0x52, 0x66)))]
    return [(text, dict(size=sz, bold=True, color=NAVY))]


def s_image_grid(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    imgs = d["images"]; cols = d.get("grid_cols", 3)
    rows = (len(imgs)+cols-1)//cols
    top = Inches(2.0)
    if d.get("bullets"):
        # 본문 불릿 폰트를 키워(3불릿 13.5→14) figure형 슬라이드 본문을 프리미엄 슬라이드
        # (slide16 2불릿 14.5)와 같은 톤으로 맞춘다 → 덱 전체 본문 밀도/스케일 통일(디자인 high).
        nb = len(d["bullets"]); bsz = d.get("bullet_size", 14.5 if nb <= 2 else 14)
        # 불릿 수에 비례해 본문 블록 높이를 잡아 아래 figure와 세로로 분리(겹침 방지).
        # body_h 가 spec 에 있으면 그 값을 우선(불릿이 많아도 figure 영역을 충분히 확보).
        bh = Inches(d["body_h"]) if d.get("body_h") else Inches(0.9 + 0.46 * max(0, nb - 2))
        _bullets_tf(slide, Inches(0.7), top, Inches(12.0), bh, d["bullets"], size=bsz, gap=5)
        top = Emu(int(top) + int(bh) + int(Inches(0.14)))
    # 결과 수치 배지. 제목이 짧으면 우측 상단(제목 줄)에 얹어 본문↔그리드 수직 슬롯을 통째로
    # 이미지에 양보(그림 확대). 제목이 길어 배지 영역(8.1")을 침범할 위험이면, 본문 아래 얇은
    # 슬롯(0.4")에 배치해 제목과의 충돌을 피한다(slide6 처럼 긴 제목 보호).
    if d.get("img_badge"):
        bw = Inches(4.6); bh2 = Inches(0.5)
        bx = Emu(int(Inches(12.7)) - int(bw))
        # 제목 폭 추정: 한글/전각 ~0.33", 영문/숫자/기호 ~0.18" (size 31 기준). 제목 우측 끝이
        # 배지 좌측을 넘으면 긴 제목으로 판단(헤더 size29→31 상향에 맞춰 폭 추정도 키움).
        ttl = d.get("title", "")
        est = sum(0.33 if ord(ch) > 0x2000 else 0.18 for ch in ttl)
        title_right = int(Inches(0.98)) + int(Inches(est))   # 제목 시작 x(0.98") + 추정 폭
        # 배지를 제목 줄(우상단)에 얹으려면 제목 우측 끝이 배지 좌측(8.1")과 충분한 간격(>0.9")을
        # 둬야 함. 제목이 중간 길이여도 배지와 빠듯하게 붙는 충돌(slide8)을 피하려고 안전선을
        # 7.95"→7.2"로 당겨, 애매하면 배지를 본문 아래 별도 슬롯으로 내린다.
        long_title = title_right > int(Inches(6.3))
        if not long_title:
            bay = Inches(0.62)
            _rect(slide, bx, bay, bw, bh2, RGBColor(0xF4,0xF6,0xF8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _rect(slide, bx, bay, Inches(0.09), bh2, ACCENT)
            _text(slide, bx+Inches(0.2), bay, Emu(int(bw)-int(Inches(0.36))), bh2,
                  [[(d["img_badge"], dict(size=13.5, bold=True, color=NAVY))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            bh2 = Inches(0.42)
            # 본문 마지막 줄(특히 회색 각주 불릿)의 descender 가 배지 상단 teal 경계에 닿지 않도록
            # 배지를 본문 블록 아래로 충분히 띄운다(0.04→0.16" — slide8 회색 각주↔녹색 칩 겹침 제거).
            bay = Emu(int(top) + int(Inches(0.16)))
            # 배지 텍스트 길이를 추정해 박스 폭을 맞춰(좌우 패딩 0.6") 한 줄에 들어가게 함 →
            # 긴 배지(예: 채택 구성 + '운영과 다른 데이터' 단서)가 박스 안에서 2줄로 잘리던
            # 충돌을 제거(slide8). 캔버스 폭(12.0") 안에서만 키운다.
            bdg = d["img_badge"]
            est_bw = sum(0.20 if ord(c) > 0x2000 else 0.115 for c in bdg) + 0.6
            bw = Emu(min(int(Inches(12.0)), max(int(Inches(4.6)), int(Inches(est_bw)))))
            # 4열 이상 그리드(식별 지도 4컷 등)에서는 배지를 우측 2컷에만 걸쳐 비대칭으로 보이지
            # 않도록 캔버스 가로 중앙(0.7~12.7")에 정렬해 4컷 전체 상단에 대칭으로 얹는다.
            if cols >= 4:
                bx = Emu(int(Inches(0.7)) + (int(Inches(12.0)) - int(bw)) // 2)
            else:
                bx = Emu(int(Inches(12.7)) - int(bw))
            _rect(slide, bx, bay, bw, bh2, RGBColor(0xF4,0xF6,0xF8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _rect(slide, bx, bay, Inches(0.09), bh2, ACCENT)
            _text(slide, bx+Inches(0.2), bay, Emu(int(bw)-int(Inches(0.36))), bh2,
                  [[(d["img_badge"], dict(size=13.5, bold=True, color=NAVY))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            top = Emu(int(top) + int(bh2) + int(Inches(0.22)))
    # 하단 note(전 패널 공통 설명)를 위한 여백 확보. 그리드 하단을 푸터 직전까지 끌어내려
    # 캔버스를 더 채우고(빈 하단 띠 제거), 이미지 셀을 가능한 한 크게 만든다. note 가 있는
    # figure 슬라이드(7·8·12·15·16)가 큰제목 슬라이드(10·13)보다 작아 보이던 'compact vs
    # full' 스케일 불일치(디자인 high)를 줄이려, 그리드 하단을 6.22→6.50", note 없는 경우
    # 6.62→6.74" 로 내려 그리드 영역 자체를 키운다(note 는 그만큼 아래 한 줄로 압축).
    note_h = Inches(0.30) if d.get("img_note") else Inches(0)
    bottom = Emu(int(Inches(6.80)) - int(note_h))
    # 패널 간 가로 여백(거터). 어둡고 빽빽한 패널(스펙트로그램 4컷 등)은 spec 의 grid_gap 으로
    # 약간 넓혀 '4장이 한 덩어리로 붙어 비슷해 보이는' 인상을 줄이고 분리감을 준다.
    gx, gy = Inches(d.get("grid_gap", 0.25)), Inches(d.get("grid_row_gap", 0.3))
    # 그림 2컷만 중앙에 떠 좌우가 비는 슬라이드(slide15)는, 우측 빈 공간에 '판정 기준' 요약 카드를
    # 두어 가로 폭을 채운다(figure row 허전함 제거). side_card 가 있으면 그리드 가용폭을 좌측으로
    # 한정하고, 비운 우측에 카드를 그린다(이미지 가공 없이 배치만 조정 — slide15 medium).
    sidecard = d.get("side_card")
    side_w = Inches(3.05) if sidecard else Inches(0)
    area_w = Emu(int(Inches(12.0)) - int(side_w) - (int(Inches(0.3)) if sidecard else 0))
    area_h = Emu(int(bottom) - int(top))
    cw = Emu(int((int(area_w) - int(gx)*(cols-1))/cols))
    chh = Emu(int((int(area_h) - int(gy)*(rows-1))/rows))
    # 캡션 높이를 열 수에 맞춰 조정(열이 많을수록 캡션이 좁아 줄바꿈되므로 약간 키움)
    # 열이 많을수록 타일 폭이 좁아 캡션이 단어 중간에서 끊김 → 4열 이상은 캡션 높이를 키우고
    # 폰트를 줄여(영문 결합어 fork+scratch 등이 한 줄에 들어가도록) 중간 줄바꿈을 줄인다.
    # 이미지 타일과 캡션 사이에 작은 숨통(cap_gap)을 둬 캡션이 타일 경계에 붙어 답답해 보이던
    # 인상을 완화(slide 7·8·9·12·15 공통 medium). 타일/이미지 실높이를 cap_gap 만큼 줄이고
    # 캡션 박스를 그만큼 아래로 내려, 이미지 가공 없이 간격만 확보한다.
    has_cap = any(im.get("caption") for im in imgs)
    has_diagram = any(im.get("diagram") for im in imgs)
    cap_gap = int(Inches(0.10)) if has_cap else 0
    cap_h = (Inches(0.46) if cols >= 5 else (Inches(0.44) if cols == 4 else Inches(0.36))) if has_cap else Inches(0)
    # 캡션 폰트를 한 단계 더 키워(11.5/12.5/13.5 → 12/13/14) figure형 슬라이드 캡션이 프리미엄
    # 슬라이드 본문 톤에 더 가깝게 읽히게 함(디자인 medium — 캡션이 작고 저대비해 보이던 인상 제거).
    cap_sz = 13 if cols >= 5 else (14 if cols == 4 else 14.5)
    cap_sz = d.get("cap_size", cap_sz)   # spec 에서 캡션 폰트를 한 단계 키울 수 있게(slide16 low)
    tile = d.get("img_tile")  # 패널 배경 틴트(연한 회색)로 흰 패널 경계를 살림
    img_box_h = int(chh) - int(cap_h) - cap_gap
    # --- 빈 하단/좌우 공백 제거: 셀 안에서 비율유지로 실제 렌더되는 이미지가 셀을 채우지
    #     못하면(가로폭 한정) 남는 세로 슬랙만큼 셀 높이를 줄이고 그리드 전체를 수직 중앙으로
    #     끌어내려, 이미지를 가능한 한 키우고 하단 빈 띠를 흡수한다. ---
    try:
        from PIL import Image as _PILImage
        ars = []
        for im in imgs:
            if im.get("diagram"):
                continue  # 네이티브 도식 셀은 셀을 그대로 채우므로 AR 산정에서 제외
            p = im["src"]
            if not os.path.isabs(p):
                cand = os.path.join(FIG_DIR, p)
                p = cand if os.path.exists(cand) else os.path.join(HERE, p)
            if os.path.exists(p):
                iw, ih = _PILImage.open(p).size
                ars.append(iw / ih)
        max_ar = max(ars) if ars else 1.0
    except Exception:
        max_ar = 1.0
    box_ar = int(cw) / max(1, img_box_h)
    col_off = 0  # 그리드 가로 중앙정렬용 좌측 오프셋
    # 단일 행(rows==1) 가로형 이미지가 셀보다 작아(box_ar>max_ar) 좌우에 큰 거터가 남으면,
    # 셀 높이를 키워 이미지를 가로폭에 맞게 확대한다(slide15 두 그래프가 중앙 55%에만 몰려
    # 좌우가 비던 문제 해결 — 이미지 가공 없이 배치 영역만 확대). 가용 area_h 안에서만 키운다.
    if (not has_diagram) and rows == 1 and max_ar <= box_ar and max_ar > 0:
        target_h = int(int(cw) / max_ar)            # 이미지가 셀 폭을 꽉 채우는 데 필요한 높이
        new_box_h = min(int(area_h) - int(cap_h), target_h)
        if new_box_h > img_box_h:
            img_box_h = new_box_h
            chh = Emu(img_box_h + int(cap_h))
            used_h = int(chh)
            # 키운 그리드를 area 안에서 수직 중앙으로 정렬(상·하 균형)
            top = Emu(int(top) + max(0, (int(area_h) - used_h) // 2))
            box_ar = int(cw) / max(1, img_box_h)
    if has_diagram:
        pass  # 네이티브 도식 그리드는 셀 폭/높이를 그대로 채움(AR 리사이즈 건너뜀)
    elif max_ar > box_ar:
        # 가로폭 한정 → 실제 이미지 높이는 cw/max_ar. 셀 세로 슬랙을 회수해 그리드를 수직 중앙 정렬.
        real_img_h = int(int(cw) / max_ar)
        slack_per_row = max(0, img_box_h - real_img_h)
        chh = Emu(int(chh) - slack_per_row)
        used_h = int(chh) * rows + int(gy) * (rows - 1)
        # 가로형(추세그래프 등) 다행 그리드는 슬랙을 완전 중앙정렬하면 본문↔그리드 사이 빈 띠가
        # 커진다 → 슬랙의 55%만 위에서 흡수해 그리드를 살짝 위로 올려 상단 빈 띠를 줄임(균형 유지).
        top = Emu(int(top) + max(0, int((int(area_h) - used_h) * 0.55)))
        img_box_h = int(chh) - int(cap_h) - cap_gap
    else:
        # 세로 한정(정사각·세로형 이미지가 넓은 셀에 들어가는 경우) → 셀 가로 슬랙(좌우 거터)을
        # 회수해 셀 폭을 이미지 실폭에 맞게 줄이고, 줄어든 만큼 그리드 전체를 가로 중앙으로 밀어
        # 이미지를 키우고 넓은 거터를 없앤다(slide10·11·15·22 그림 확대).
        real_img_w = int(img_box_h * max_ar)
        slack_per_col = max(0, int(cw) - real_img_w)
        # 셀 폭을 실제 이미지 폭 + 캡션 가독을 위한 소폭 패딩(0.12")만 남기고 축소
        keep_pad = int(Inches(0.12))
        new_cw = real_img_w + keep_pad
        if new_cw < int(cw):
            cw = Emu(new_cw)
            if cols >= 2:
                # 가로형 이미지(2열 추세그래프 포함): 셀을 중앙에 작게 몰지 않고, 남는 가로 슬랙을
                # 열 간격에 분배해 그리드가 폭을 꽉 채우게 한다(좌우 균형·하단 빈 사분면 인상 제거).
                # 좌우는 한 칸 거터(0.55")만 들이고, 나머지는 열 사이 간격으로 흡수(slide15 2그래프
                # 가 중앙 55%에만 몰려 좌측 하단이 비던 문제 해결 — 이미지 가공 없이 배치만 조정).
                side = int(Inches(0.55)) if cols == 2 else int(Inches(0.45))
                span = int(area_w) - 2 * side
                gx = Emu(max(int(Inches(0.25)), (span - int(cw) * cols) // (cols - 1)))
                col_off = side
            else:
                used_w = int(cw) * cols + int(gx) * (cols - 1)
                col_off = max(0, (int(area_w) - used_w) // 2)
    for i, im in enumerate(imgs):
        r, c = divmod(i, cols)
        x = Emu(int(Inches(0.7)) + col_off + c*(int(cw)+int(gx)))
        y = Emu(int(top) + r*(int(chh)+int(gy)))
        img_h = Emu(img_box_h)
        if im.get("diagram"):
            # 코드 스크린샷 대신 네이티브 도식을 셀 안에 직접 그린다(흰 카드 + 또렷한 테두리).
            _rect(slide, x, y, cw, img_h, WHITE,
                  line=RGBColor(0xB9, 0xC6, 0xDB), line_w=Pt(1.25))
            if im["diagram"] == "bytes_compare":
                _bytes_compare_diagram(slide, x, y, cw, img_h)
            if im["diagram"] == "hex_to_grade":
                _hex_to_grade_diagram(slide, x, y, cw, img_h)
            if im["diagram"] == "numba_composite":
                _numba_composite_diagram(slide, x, y, cw, img_h)
            if im["diagram"] == "pyvips_stream":
                _pyvips_stream_diagram(slide, x, y, cw, img_h)
            if im["diagram"] == "episode_trend":
                _episode_trend_diagram(slide, x, y, cw, img_h)
            if im["diagram"] == "p3_quad":
                _p3_quad_diagram(slide, x, y, cw, img_h)
            if im.get("caption"):
                cap_pad = int(gx) // 2
                _text(slide, Emu(int(x) - cap_pad), Emu(int(y)+int(chh)-int(cap_h)),
                      Emu(int(cw) + cap_pad * 2), cap_h,
                      [_caption_runs(im["caption"], cap_sz)],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            continue
        rf = im.get("result_frame")
        if tile:
            # 저대비·sparse 한 패널(원본 wafer/chip map, '정상' 타일 등)이 '빈 칸/깨진 이미지'로
            # 읽히지 않도록 타일을 또렷한 mid-gray 테두리(1.25pt)로 감싸 '의도된 데이터 카드'임을
            # 명확히 한다(이미지 가공 없이 프레임만 강화 — slide 8·12 빈 패널 인상 제거).
            # result_frame 패널(톤이 다른 '검출 결과' 이미지)은 teal 강조 테두리로 감싸 '튀는 것'이
            # 아니라 '의도된 결과 대비'임을 명확히 한다(slide7 medium — 3번째 이미지 톤 불일치).
            _rect(slide, x, y, cw, img_h, tile,
                  line=(ACCENT if rf else RGBColor(0xB9, 0xC6, 0xDB)),
                  line_w=Pt(1.9 if rf else 1.25))
        img_frame = d.get("img_frame", True) is not False and im.get("frame", True) is not False
        pic = _img_fit(slide, im["src"], x, y, cw, img_h, frame=(not tile) and img_frame)
        if rf:
            # '검출 결과' 라벨 칩을 타일 좌상단에 얹어 색감이 다른 결과 이미지가 의도된 대비임을 명시.
            rlw = int(Inches(1.12)); rlh = int(Inches(0.32))
            rlx = int(x) + int(Inches(0.08)); rly = int(y) + int(Inches(0.08))
            rchip = _rect(slide, Emu(rlx), Emu(rly), Emu(rlw), Emu(rlh), ACCENT,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = rchip.text_frame; tf.word_wrap = False
            tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = str(rf) if not isinstance(rf, bool) else "Detection result"
            _set_font(rr, size=10.5, bold=True, color=WHITE)
        # trend 산점도(matplotlib)의 영문 축 라벨·작은 눈금 숫자(좌측 'Measurement Value (nm)' 세로,
        # 하단 'time_index')는 비즈니스 톤과 이질적이고 가독성 하한이다 → 이미지 픽셀은 가공하지 않고,
        # 배치된 이미지 사각형의 좌측·하단 축 크롬 영역만 타일색 사각면으로 덮어(passe-partout 마스크)
        # 가린 뒤, 그 자리에 깔끔한 한글 축 설명('계측값'·'시간 순서')을 얹는다(slide15·16 medium 해결).
        # axis_cover=[좌측비율,하단비율](가린 영역 0~1). 산점도 본문(점·기준선)은 그대로 보존된다.
        if im.get("axis_cover") and hasattr(pic, "_fit_rect"):
            px_, py_, pw_, ph_ = [int(v) for v in pic._fit_rect]
            lf, bf = im["axis_cover"][0], im["axis_cover"][1]
            mask_col = _col(tile) if tile else WHITE
            lw_m = int(lf * pw_); bh_m = int(bf * ph_)
            # 좌측 축 숫자 strip 마스크(세로 축 라벨/눈금 가림)
            if lw_m > 0:
                _rect(slide, Emu(px_), Emu(py_), Emu(lw_m), Emu(ph_), mask_col)
            # 하단 축 라벨 strip 마스크('time_index'·눈금 숫자 가림)
            if bh_m > 0:
                _rect(slide, Emu(px_), Emu(py_ + ph_ - bh_m), Emu(pw_), Emu(bh_m), mask_col)
            # 깔끔한 한글 축 설명(임원 가독) — 좌측 세로축='계측값', 하단 가로축='시간 순서'
            if im.get("axis_y") is not False:
                _text(slide, Emu(px_ - int(Inches(0.02))), Emu(py_ + ph_ // 2 - int(Inches(0.12))),
                      Emu(lw_m + int(Inches(0.5))), Emu(int(Inches(0.24))),
                      [[(im.get("axis_y", "Measured value"), dict(size=9.5, bold=True, color=MUTED))]],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if im.get("axis_x") is not False:
                _text(slide, Emu(px_), Emu(py_ + ph_ - bh_m), Emu(pw_), Emu(max(bh_m, int(Inches(0.2)))),
                      [[(im.get("axis_x", "Time order →"), dict(size=9.5, bold=True, color=MUTED))]],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 하단 strip 만큼만 산점도 본문 높이를 좁혀 band(이상 음영)가 가려진 축 영역 아래로
            # 흘러내리지 않게 한다(가로 위치/폭은 원본 그대로라 band b0/b1 비율이 유지됨).
            pic._fit_rect = (px_, py_, pw_, ph_ - bh_m)
        # 거의 흰 배경의 저대비 원본 map(raw 패널)은 시각적으로 '빈 박스'처럼 읽힌다 → 이미지는
        # 가공하지 않고, 배치된 이미지 사각형 위에 아주 옅은 음영(반투명)만 얹어 '이미지가 존재함'을
        # 명확히 한다(slide8 medium 해결 — 픽셀 불변, 색면만 오버레이).
        if im.get("annot") and hasattr(pic, "_fit_rect"):
            px_, py_, pw_, ph_ = pic._fit_rect
            sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(px_), Emu(py_), Emu(pw_), Emu(ph_))
            _set_alpha_fill(sh, RGBColor(0x8E, 0x9F, 0xBC), 8)
        # 저대비 원본 wafer map(거의 백지로 보이는 raw 패널)은 가공 대신 결함 군집 위에 고채도
        # 주황 가이드 원/화살표를 얹어(이미지 픽셀 불변) '빈 그림'이 아니라 '대비 메시지의 좌측'
        # 임을 발표 거리에서 보이게 한다(slide8 medium 해결).
        if im.get("annot") and hasattr(pic, "_fit_rect"):
            _annot_over(slide, pic._fit_rect, im["annot"])
        # trend 미니차트: 세로 점선 이후 '이상 구간'(우측 일부)에 옅은 따뜻한 색 음영 밴드를
        # 얹어, 미세한 축 라벨 없이도 색·위치만으로 '오른쪽=이상'을 즉시 읽히게 한다(이미지 가공
        # 없이 반투명 색면 오버레이만 — slide16 low). band=[시작비율, 끝비율](가로 0~1).
        if im.get("band") and hasattr(pic, "_fit_rect"):
            px_, py_, pw_, ph_ = pic._fit_rect
            b0, b1 = im["band"][0], im["band"][1]
            bx_ = int(px_) + int(b0 * int(pw_))
            bw_ = int((b1 - b0) * int(pw_))
            sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(bx_), Emu(int(py_)),
                                        Emu(bw_), Emu(int(ph_)))
            # 밴드 불투명도를 12→16 으로 한 단계 올려 작은 썸네일에서도 '오른쪽=이상' 영역이 더
            # 또렷하게 읽히게 하고(저대비 산점도 보강 — slide15·16 medium), 정상/이상 경계에 가는
            # 빨강 세로 룰을 얹어 파랑 정상군과 빨강 이상군의 구분을 강화한다(이미지 가공 없이 오버레이).
            _set_alpha_fill(sh, RGBColor(0xE5, 0x3A, 0x1F), 16)
            _rect(slide, Emu(bx_), Emu(int(py_)), Emu(int(Pt(1.4))), Emu(int(ph_)),
                  RGBColor(0xE5, 0x3A, 0x1F))
            # 밴드 상단에 작은 '이상' 라벨로 의도를 명시(축 텍스트 의존 제거)
            lh_ = int(Inches(0.26)); lw_ = int(Inches(0.62))
            lbx = bx_ + (bw_ - lw_) // 2
            lab = _rect(slide, Emu(lbx), Emu(int(py_) + int(Inches(0.04))), Emu(lw_), Emu(lh_),
                        WHITE, line=RGBColor(0xE5, 0x3A, 0x1F), line_w=Pt(1),
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = lab.text_frame; tf.word_wrap = False
            tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = "Anomaly"
            _set_font(rr, size=10, bold=True, color=RGBColor(0xE5, 0x3A, 0x1F))
        if im.get("outline") and hasattr(pic, "_fit_rect"):
            # 흰 바탕에 색점 몇 개만 흩어진 식별 지도(objid) 패널은 박스 내부가 텅 비어 보인다 →
            # 이미지 픽셀은 가공하지 않고, 원본 wafer 디스크와 같은 기준 영역(팔각 윤곽선)을 옅게
            # 얹어 빈 흰 공간에 '같은 wafer 위 분류 결과'라는 기준틀을 부여한다(slide8 medium 해결).
            px_, py_, pw_, ph_ = [int(v) for v in pic._fit_rect]
            side_ = min(pw_, ph_)
            ox = px_ + (pw_ - side_) // 2; oy = py_ + (ph_ - side_) // 2
            # 옅은 팔각 윤곽(원본 raw map 의 octagon 디스크와 동일한 기준틀) — 가는 회청색 선만.
            oct_sp = slide.shapes.add_shape(MSO_SHAPE.OCTAGON,
                                            Emu(ox + side_ // 24), Emu(oy + side_ // 24),
                                            Emu(side_ - side_ // 12), Emu(side_ - side_ // 12))
            oct_sp.fill.background()
            oct_sp.line.color.rgb = RGBColor(0xC4, 0xD0, 0xE2); oct_sp.line.width = Pt(1.4)
            oct_sp.shadow.inherit = False
            # 식별 점이 모인 영역을 가리키는 고채도 주황 가이드 원(원본 raw 의 결함 군집과 동일 위치).
            if im.get("annot"):
                _annot_over(slide, (px_, py_, pw_, ph_), im["annot"])
        if im.get("pair"):
            # 원본↔식별 쌍(A쌍/B쌍)을 즉시 읽히게 타일 좌상단에 작은 라벨 칩을 얹어 4컷이 2x2
            # 비교임을 명확히 한다(같은 pair 끼리 같은 라벨). 이미지 가공 없이 오버레이만.
            plw = int(Inches(0.62)); plh = int(Inches(0.30))
            plx = int(x) + int(Inches(0.08)); ply = int(y) + int(Inches(0.08))
            pchip = _rect(slide, Emu(plx), Emu(ply), Emu(plw), Emu(plh), NAVY,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = pchip.text_frame; tf.word_wrap = False
            tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = "pair " + im["pair"]
            _set_font(rr, size=10.5, bold=True, color=WHITE)
        if im.get("watermark") and hasattr(pic, "_fit_rect"):
            # 의도적으로 거의 빈 '정상 — 패턴 없음' 패널이 렌더 누락/빈 이미지로 오인되지 않게,
            # 셀 가로 중앙에 옅은 기준선 + 흐린 워터마크 라벨을 얹어 '의도된 빈 상태(정상)'임을
            # 시각적으로 분명히 한다(이미지 가공 없이 오버레이만 — slide12 low).
            px_, py_, pw_, ph_ = [int(v) for v in pic._fit_rect]
            # 다른 3개 패널(짙은 텍스처)과의 밀도 불균형을 줄이려, '정상=균질' 패널 위에 옅은 회청색
            # 점 격자(homogeneous texture)를 규칙적으로 깔아 '깨진 빈 이미지'가 아니라 '고르게 분포한
            # 정상 텍스처'로 읽히게 한다(이미지 픽셀 불변 — 도형 오버레이만, slide12 medium 해결).
            nx, ny = 11, 11
            dot = int(Inches(0.045))
            mx0 = px_ + int(Inches(0.18)); my0 = py_ + int(Inches(0.18))
            gw = pw_ - int(Inches(0.36)); gh = ph_ - int(Inches(0.36))
            for gi in range(ny):
                for gj in range(nx):
                    # 주의: 그리드 거터 변수(gx/gy)와 충돌하지 않게 별도 지역변수(dgx/dgy) 사용.
                    # 과거 gx/gy 를 여기서 덮어써 watermark 다음 타일(4번째)의 x가 화면 밖으로
                    # 튕겨나가 렌더 누락되던 버그를 막는다.
                    dgx = mx0 + (gw * gj) // (nx - 1)
                    dgy = my0 + (gh * gi) // (ny - 1)
                    _rect(slide, Emu(dgx), Emu(dgy), Emu(dot), Emu(dot),
                          RGBColor(0xCD, 0xD6, 0xE3), shape=MSO_SHAPE.OVAL)
            base_y = py_ + ph_ // 2
            _rect(slide, Emu(px_ + int(Inches(0.12))), Emu(base_y),
                  Emu(pw_ - int(Inches(0.24))), Emu(int(Pt(1.4))), RGBColor(0xB6, 0xC1, 0xD2))
            # 워터마크 라벨에 옅은 흰 받침을 깔아 점 격자 위에서도 또렷이 읽히게 한다.
            _lw = int(Inches(2.1)); _lh = int(Inches(0.32))
            _rect(slide, Emu(px_ + (pw_ - _lw) // 2), Emu(base_y - int(Inches(0.42))),
                  Emu(_lw), Emu(_lh), WHITE, line=RGBColor(0xC2, 0xCC, 0xDC), line_w=Pt(0.75),
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _text(slide, Emu(px_), Emu(base_y - int(Inches(0.42))), Emu(pw_), Emu(_lh),
                  [[(im["watermark"], dict(size=11, bold=True, color=RGBColor(0x8A, 0x93, 0xA3)))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if im.get("badge"):
            # 거의 비어 보이는 '정상' 패널이 빈 칸/누락으로 읽히지 않게, 타일 안에 작은 라벨 칩을
            # 얹어 '의도적으로 깨끗함(정상=균질)'임을 명시(이미지 가공 없이 오버레이만).
            bw_ = int(cw) - int(Inches(0.36)); bh_ = int(Inches(0.34))
            bxx = int(x) + (int(cw) - bw_) // 2
            byy = int(y) + int(img_h) - bh_ - int(Inches(0.12))
            chip = _rect(slide, Emu(bxx), Emu(byy), Emu(bw_), Emu(bh_),
                         RGBColor(0xF4, 0xF7, 0xFB), line=RGBColor(0xB9, 0xC6, 0xDB),
                         line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = chip.text_frame; tf.word_wrap = False
            tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = im["badge"]
            _set_font(rr, size=10.5, bold=True, color=RGBColor(0x55, 0x60, 0x75))
        if im.get("caption"):
            # 캡션 박스를 셀 폭보다 좌우로 약간 넓혀(인접 거터의 절반씩) 한 줄 캡션이 단어 중간에서
            # 줄바꿈되지 않게 함(이미지 타일 폭은 그대로, 텍스트 박스만 확장). 셀 경계를 넘지 않는
            # 안전선 안에서만 확장해 이웃 캡션과 겹치지 않게 한다.
            cap_pad = int(gx) // 2
            cap_x = Emu(int(x) - cap_pad)
            cap_w = Emu(int(cw) + cap_pad * 2)
            _text(slide, cap_x, Emu(int(y)+int(chh)-int(cap_h)), cap_w, cap_h,
                  [_caption_runs(im["caption"], cap_sz)],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sidecard:
        # 우측 빈 공간을 '정상 vs 이상 판정 기준' 요약 카드로 채워 가로 폭 균형(slide15 medium).
        sx = Emu(int(Inches(0.7)) + int(area_w) + int(Inches(0.34)))
        sy = top; sh = Emu(int(bottom) - int(top))
        _rect(slide, sx, sy, side_w, sh, RGBColor(0xF4, 0xF7, 0xFB),
              line=RGBColor(0xB9, 0xC6, 0xDB), line_w=Pt(1.1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        hh = Inches(0.46)
        _rect(slide, sx, sy, side_w, hh, ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(slide, sx, Emu(int(sy) + int(hh) - int(Inches(0.06))), side_w, Inches(0.06), ACCENT)
        _text(slide, sx, sy, side_w, hh,
              [[(sidecard.get("head", ""), dict(size=12.5, bold=True, color=WHITE))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        srows = sidecard.get("rows", [])
        ry0 = int(sy) + int(hh) + int(Inches(0.16))
        rh = (int(sh) - int(hh) - int(Inches(0.28))) // max(1, len(srows))
        for k, (lab, txt) in enumerate(srows):
            ry = ry0 + k * rh
            badge_col = ACCENT if lab in ("정상", "Normal") else RGBColor(0xE5, 0x3A, 0x1F)
            _rect(slide, Emu(int(sx) + int(Inches(0.18))), Emu(ry),
                  Emu(int(Inches(0.78))), Emu(int(Inches(0.34))), badge_col,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _text(slide, Emu(int(sx) + int(Inches(0.18))), Emu(ry),
                  Emu(int(Inches(0.78))), Emu(int(Inches(0.34))),
                  [[(lab, dict(size=11, bold=True, color=WHITE))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _text(slide, Emu(int(sx) + int(Inches(0.18))), Emu(ry + int(Inches(0.42))),
                  Emu(int(side_w) - int(Inches(0.36))), Emu(rh - int(Inches(0.46))),
                  [[(txt, dict(size=11.5, color=INK))]], anchor=MSO_ANCHOR.TOP)
    if d.get("img_note"):
        # 하단 캡션 가독성: 연회색(MUTED) → 다크 슬레이트(#3A485C)로 대비 상향 + 폰트 확대(12→12.5).
        # 풀해상도 슬라이드 본문(최소 ~13)에 더 가깝게 키워 figure형 슬라이드와 톤(폰트 크기) 일관성
        # 을 맞춘다(짧게 다듬은 note 와 함께 한눈에 읽히게 — 디자인 medium 해결).
        _text(slide, Inches(0.7), Emu(int(Inches(6.80))-int(note_h)+int(Inches(0.02))), Inches(12.0), note_h,
              [[(d["img_note"], dict(size=13, color=RGBColor(0x3A,0x48,0x5C)))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    _footer(slide, idx)


def s_table(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    headers = d["headers"]; rows = d["rows"]
    nr, nc = len(rows)+1, len(headers)
    tw = Inches(12.0)
    # 표를 제목 바로 아래(1.78")에서 시작해 제목↔표 사이 공백을 줄이고, 본문 영역을 더 채우도록
    # 행 높이를 키운다. 남는 세로 공간은 표 전체를 수직 중앙으로 내려 상·하 균형(하단 공백 제거).
    # 각주(note)가 여러 줄(\n 포함)이면 표 하단(bottom0)을 끌어올려 2~3줄 각주가 footer 와
    # 겹치지 않게 공간을 확보한다(한 줄 각주는 기존대로 6.4").
    note_lines = (str(d.get("note", "")).count("\n") + 1) if d.get("note") else 0
    top0 = int(Inches(1.78))
    # 2줄 각주가 있어도 표 하단을 6.22"까지 내려(이전 6.05") 행 높이를 키우고 표가 본문 영역을
    # 꽉 채우게 한다(slide3 하단 옅은 여백 흡수). 각주는 표 바로 아래(min 캡)에 붙어 footer 와
    # 겹치지 않는다.
    bottom0 = int(Inches(6.22 if note_lines >= 2 else 6.4))
    avail = bottom0 - top0
    # 행 높이 상한을 1.2→1.35"로 키워 행 적은 표(3행+헤더)가 본문을 꽉 채우게 함(하단 공백 제거)
    row_h = min(int(Inches(1.35)), avail // nr)
    th = row_h * nr
    # 남는 세로 슬랙은 50%를 위에서 흡수해 표를 본문 세로 중앙에 균형 있게 배치(상·하 여백 균등)
    ty = top0 + max(0, int((avail - th) * 0.5))
    gtbl = slide.shapes.add_table(nr, nc, Inches(0.7), Emu(ty), tw, Emu(th)).table
    for ri in range(nr):
        gtbl.rows[ri].height = Emu(row_h)
    hl = d.get("highlight_row")
    for j, htxt in enumerate(headers):
        cell = gtbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(htxt); _set_font(r, size=13, bold=True, color=WHITE)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    # 설명형(여러 줄·서술) 셀은 좌측, 수치/짧은 라벨 셀은 가운데로 정렬해 읽기 흐름을 살림.
    def _is_descriptive(s):
        s = str(s)
        if "\n" in s:
            return True
        # 한글이 다수 포함되고 길이가 길면 서술형으로 간주(순수 수치/% 는 가운데 유지)
        han = sum(1 for ch in s if 0xAC00 <= ord(ch) <= 0xD7A3)
        return han >= 2 and len(s) >= 8
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtbl.cell(i+1, j)
            if hl is not None and i == hl:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF4,0xF6,0xF8)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else PANEL
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            if j == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.LEFT if _is_descriptive(val) else PP_ALIGN.CENTER
            # 서술형 좌측정렬 셀은 좌측 안쪽 여백을 줘 헤더·셀 경계에 붙지 않게 함
            cell.margin_left = Pt(7) if p.alignment == PP_ALIGN.LEFT else Pt(2)
            cell.margin_right = Pt(4)
            r = p.add_run(); r.text = str(val)
            bold = (hl is not None and i == hl)
            # 셀 폰트를 키워(11.5→13) 표 가독성을 프리미엄 표(slide13)와 동일 톤으로 맞춘다.
            # 단 5열 이상(slide13 P2 성과표)은 셀 폭이 좁아 12.5 로만 키워 수치가 줄바꿈되지 않게 함.
            cell_sz = 13 if nc <= 4 else 12.5
            _set_font(r, size=cell_sz, bold=bold, color=NAVY if bold else INK)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    if d.get("note"):
        # 각주가 footer(7.04")와 붙지 않도록 상한을 둔다. 표 하단을 6.22"로 내린 만큼 2줄 각주
        # 상한도 6.28"로 함께 내려(표와 겹치지 않게) footer 위에 2줄이 수렴하게 하고, 한 줄
        # 각주는 기존(6.5") 상한을 유지.
        note_cap = int(Inches(6.28 if note_lines >= 2 else 6.5))
        note_y = min(note_cap, ty + th + int(Inches(0.08)))
        # 여러 줄 각주는 폰트를 소폭 줄여(12.5→11.5) 줄 수가 늘어도 footer 와 겹치지 않게 한다.
        note_sz = 11.5 if note_lines >= 2 else 12.5
        # 좌우 여백을 본문 표와 동일(0.7~12.7")하게 두되 줄간격을 소폭 확대해 2줄 각주의 답답함 완화
        tb = slide.shapes.add_textbox(Inches(0.7), Emu(note_y), Inches(12.0), Inches(0.7))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        # note 에 명시적 줄바꿈(\n)이 있으면 문장별 문단으로 분리해 가독성 확보(한 줄에 4정보 압축 방지).
        # 여러 줄 각주는 항목별 '•' 마커를 앞에 붙여 위계를 분명히 하고, **...** 굵게 마크업을 지원해
        # 가장 중요한 단서(예: 동급 비교 아님)를 강조할 수 있게 한다.
        import re as _re_note
        segs = str(d["note"]).split("\n")
        for li, seg in enumerate(segs):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.line_spacing = 1.12; p.space_after = Pt(1)
            if len(segs) >= 2:
                # 행잉 인덴트: 각주가 한 줄을 넘겨 wrap 되어도 둘째 줄이 '•' 마커가 아니라 본문
                # 시작선에 정렬되게 한다(slide3 — wrap 된 둘째 줄이 좌측 끝에서 시작해 첫 줄과
                # 어긋나 보이던 문제 제거). marL=마커폭, indent=음수로 첫 줄만 마커부터 시작.
                hang = int(Inches(0.22))
                pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
                pPr.set('marL', str(hang)); pPr.set('indent', str(-hang))
                rm = p.add_run(); rm.text = "•  "
                _set_font(rm, size=note_sz, bold=True, color=ACCENT)
            for part in _re_note.split(r"(\*\*.+?\*\*)", seg):
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**"):
                    r.text = part[2:-2]
                    _set_font(r, size=note_sz, bold=True, color=NAVY)
                else:
                    r.text = part
                    _set_font(r, size=note_sz, color=RGBColor(0x3A, 0x48, 0x5C))
    _footer(slide, idx)


def s_closing(slide, d, idx):
    _bg(slide, WHITE)
    has_kpi = bool(d.get("kpis"))
    # 강조 규칙 통일: 본문 슬라이드와 동일한 '제목 왼쪽 세로 teal 바' + eyebrow 라벨(헤더 일관성)
    if d.get("kicker"):
        _rect(slide, Inches(0.7), Inches(0.72), Inches(0.11), Inches(0.28), ACCENT)
        _text(slide, Inches(1.0), Inches(0.7), Inches(11.5), Inches(0.3),
              [[(d["kicker"], dict(size=12.5, bold=True, color=ACCENT))]])
        ty = Inches(1.12)
    else:
        ty = Inches(1.0)
    # 본문 슬라이드 헤더(_title_block)와 완전히 동일한 좌측 teal 바(0.135w×0.84h)·제목 위계(31)로
    # 통일 → 개요/마무리/상세 슬라이드 헤더 톤을 한 스펙으로 맞춘다(템플릿 혼재 medium 완화).
    _rect(slide, Inches(0.7), ty, Inches(0.135), Inches(0.84), ACCENT)
    _text(slide, Inches(1.02), Emu(int(ty)+int(Inches(0.04))), Inches(11.5), Inches(1.0),
          [[(d["title"], dict(size=31, bold=True, color=NAVY))]])
    if d.get("bullets"):
        # 본문 불릿 블록을 약간 내려(2.3→2.62) 헤더↔카드 사이 수직 중앙으로 정렬해, 본문이 위로
        # 몰리고 본문↔카드 사이에 큰 빈 띠가 생기던 인상을 제거(줄간격도 소폭 키워 균형 보강).
        # 본문 4불릿(첫 불릿의 보조설명을 회색 lvl=1 한 줄로 분리) + 회색 각주 1줄이 아래 KPI
        # 카드(5.02") 위에서 끝나도록, 블록을 약간 올리고(2.62→2.46) 가용 높이를 카드 직전
        # (2.46→4.86", 2.40")까지로 넓히며 불릿 간격을 8→6 으로 조여 한 줄 늘어난 본문이 카드
        # 영역으로 흘러내려 겹치던 충돌을 제거한다(slide17 — 첫 불릿 2줄 압축 + 보조설명 분리).
        _bullets_tf(slide, Inches(0.97), Inches(2.40), Inches(11.4), Inches(2.42),
                    d["bullets"], size=15.6, gap=11, line_spacing=1.08)
    if has_kpi:
        kpis = d["kpis"]; n = len(kpis)
        gap = Inches(0.3); total = Inches(11.4)
        cw = Emu(int((total - gap*(n-1)) / n))
        # 카드 내부 상하 패딩이 과해 숫자-라벨-부제가 떠 보이던 인상(slide18 low)을 줄이고자
        # 카드 높이를 1.86"로 줄여 세 줄을 더 단단히 묶고, 줄어든 만큼 카드를 약간 내려(4.92")
        # 본문 블록과의 간격을 유지한다(footer 7.04" 와도 안전 간격 확보).
        # 보조설명(회색 '–' 각주) 줄과 KPI 카드 사이 수직 간격을 넓혀(4.84→5.02) 카드 영역을
        # 시각적으로 분리한다(slide18 low — 카드가 각주 바로 아래 붙어 답답하던 인상 제거).
        # 늘린 간격만큼 카드 높이를 소폭 줄여(1.86→1.80) footer(7.04") 위 안전선을 유지한다.
        x0 = Inches(0.97); y = Inches(5.02); ch = Inches(1.80)
        # 4칩 큰 숫자 폰트 통일: 가장 긴 big 텍스트(예: '오경보 0.00%')가 카드 폭 안에 들어가는
        # 크기를 한 번 계산해 모든 카드에 동일 적용 → '오경보 0.00%'만 커 보이던 불일치 제거.
        big_pt = 27
        card_inner = int(cw) - int(Inches(0.2))
        for st in kpis:
            est_w = sum((0.62 if ord(c) > 0x2000 else 0.40) for c in str(st["big"]))  # pt27 글자폭 추정(in)
            while big_pt > 18 and est_w * (big_pt / 27.0) * int(Inches(1)) > card_inner:
                big_pt -= 1
        for i, st in enumerate(kpis):
            x = Emu(int(x0) + i*(int(cw)+int(gap)))
            # 마무리 KPI 4칩은 색·높이를 완전히 통일 — 모두 동일 회색(PANEL) 배경 + teal 상단 바.
            # (이전엔 사업가치 카드만 연한 teal 틴트로 차별화했으나, 마무리 슬라이드는 칩 간
            #  색·높이 일관성이 더 중요해 강조 틴트를 제거하고 회색으로 통일한다.)
            _rect(slide, x, y, cw, ch, PANEL)
            _rect(slide, x, y, cw, Inches(0.1), ACCENT)
            # 카드 내부(big+label+sub)를 한 텍스트 블록으로 묶어 카드 전체 높이에 수직 중앙정렬 →
            # 상단 teal 바 아래 가용 영역에 균등 배치(이전엔 big/label 이 위로 몰리고 하단이 비었음).
            blk = [[(st["big"], dict(size=big_pt, bold=True, color=NAVY))],
                   [(st["label"], dict(size=13.2, bold=True, color=INK))]]
            if st.get("sub"):
                blk.append([(st.get("sub"), dict(size=11.4, color=RGBColor(0x55,0x60,0x75)))])
            tbk = _text(slide, x+Inches(0.1), Emu(int(y)+int(Inches(0.1))),
                        Emu(int(cw)-int(Inches(0.2))), Emu(int(ch)-int(Inches(0.1))), blk,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # big 줄과 label 줄 사이 숨통(큰 숫자 아래 여백) — 한 블록이라 단락 간격으로 분리
            tfk = tbk.text_frame
            tfk.paragraphs[0].space_after = Pt(8)
            if st.get("sub"):
                tfk.paragraphs[1].space_after = Pt(2)
    _footer(slide, idx)


def _flow_box(slide, x, y, w, h, b):
    color = b.get("color", "panel")
    _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    dark = _is_dark(color)
    # Accent 배경 위에서는 gold가 저대비로 보일 수 있어 흰색 볼드로 통일하고,
    # gold는 navy/navy2 배경 강조에만 사용한다.
    is_accent = str(_col(color)).upper() == "96A0AD"
    tcol = WHITE if dark else NAVY
    if is_accent:
        scol = WHITE
        tagcol = WHITE
    elif dark:
        scol = RGBColor(0xD2, 0xDE, 0xF3)    # navy 위 본문
        tagcol = GOLD                        # navy 위 강조만 앰버
    else:
        scol = MUTED
        tagcol = ACCENT
    pad = Inches(0.1)
    lines = [[(b.get("title", ""), dict(size=12.5, bold=True, color=tcol))]]
    if b.get("sub"):
        lines.append([(b["sub"], dict(size=9.5, color=scol))])
    if b.get("tags"):
        tagtxt = "  ".join(b["tags"]) if isinstance(b["tags"], list) else str(b["tags"])
        lines.append([(tagtxt, dict(size=9, bold=True, color=tagcol))])
    _text(slide, Emu(int(x) + int(pad)), y, Emu(int(w) - 2 * int(pad)), h, lines,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def s_flow(slide, d, idx):
    """아키텍처/파이프라인 흐름도: 단계(stage)를 좌→우로, 단계 안 박스는 세로로 쌓고 단계 사이 화살표."""
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    stages = d["stages"]; n = len(stages)
    area_x, area_w = int(Inches(0.6)), int(Inches(12.13))
    # 박스가 세로로 과하게 길어 상·하단 dead space 가 컸음 → 밴드 높이를 정보량에 맞게 축소(4.25→3.4")
    # 하고, 단계 라벨(상단)~caption(하단) 사이 수직 중앙에 배치해 공백을 균형 있게 흡수.
    # 한 단계가 박스 2개로 쪼개지는 경우(Known/Unknown 등)는 세로 공간이 더 필요해 밴드를 약간 키운다.
    # 밴드 하단은 caption(6.55") 위 안전선(6.42") 을 넘지 않도록 top 을 함께 제한(겹침 방지).
    max_boxes = max(len(st["boxes"]) for st in stages)
    band_bottom = int(Inches(6.42))
    area_h = int(Inches(3.75 if max_boxes >= 2 else 3.4))
    top = int(Inches(2.5)) + max(0, (int(Inches(4.25)) - area_h)) // 2
    if top + area_h > band_bottom:
        top = band_bottom - area_h
    arrow_w = int(Inches(0.4))
    col_w = (area_w - arrow_w * (n - 1)) // n
    bgap = int(Inches(0.2))

    def box_centers(m):
        """단계 안 m개 박스의 세로 중심 y좌표 리스트(박스 높이 동일 분할 기준)."""
        bh = (area_h - bgap * (m - 1)) // m
        return [top + bi * (bh + bgap) + bh // 2 for bi in range(m)]

    for si, st in enumerate(stages):
        cx = area_x + si * (col_w + arrow_w)
        if st.get("label"):
            _text(slide, Emu(cx), Emu(top - int(Inches(0.36))), Emu(col_w), Inches(0.32),
                  [[(st["label"], dict(size=11, bold=True, color=ACCENT))]], align=PP_ALIGN.CENTER)
        boxes = st["boxes"]; m = len(boxes)
        bh = (area_h - bgap * (m - 1)) // m
        for bi, b in enumerate(boxes):
            by = top + bi * (bh + bgap)
            _flow_box(slide, Emu(cx), Emu(by), Emu(col_w), Emu(bh), b)
        if si < n - 1:
            # 다음 단계 박스 수에 맞춰 화살표를 각 박스 세로 중심으로 분기 연결(공백·끊김 방지).
            nxt_m = len(stages[si + 1]["boxes"])
            ah = int(Inches(0.3))
            ax = cx + col_w + int(Emu(15000))
            aw = arrow_w - int(Emu(60000))
            if nxt_m <= 1:
                ay = top + area_h // 2 - ah // 2
                _rect(slide, Emu(ax), Emu(ay), Emu(aw), Emu(ah), ACCENT,
                      shape=MSO_SHAPE.RIGHT_ARROW)
            else:
                # 현재 단계 중심 → 세로 분배 막대 → 각 다음 박스로 향하는 화살표
                # 커넥터(스텁·분배 막대)를 굵게(약 0.05") 키워 분기 흐름이 한눈에 이어지게 함
                ys = box_centers(nxt_m)
                src_y = top + area_h // 2
                cw_bar = int(Emu(46000))   # 분배 막대/스텁 두께(기존 22000→46000)
                bar_x = cx + col_w + arrow_w // 2 - cw_bar // 2
                bar_w = cw_bar
                y_lo, y_hi = min(ys), max(ys)
                # 중심에서 분배 막대까지 짧은 수평 스텁
                _rect(slide, Emu(cx + col_w), Emu(src_y - cw_bar // 2),
                      Emu(bar_x - (cx + col_w) + bar_w), Emu(cw_bar), ACCENT)
                # 수직 분배 막대
                _rect(slide, Emu(bar_x), Emu(y_lo), Emu(bar_w), Emu(y_hi - y_lo), ACCENT)
                # 각 다음 박스로 화살표
                for yc in ys:
                    _rect(slide, Emu(bar_x + bar_w), Emu(yc - ah // 2),
                          Emu((cx + col_w + arrow_w) - (bar_x + bar_w) + int(Emu(8000))),
                          Emu(ah), ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
    if d.get("caption"):
        _text(slide, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.38),
              [[(d["caption"], dict(size=11, color=MUTED))]])
    _footer(slide, idx)


def s_timeline(slide, d, idx):
    """로드맵 타임라인: 가로 축 + 마일스톤(완료=채움/예정=외곽선), 위/아래 번갈아 라벨."""
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    items = d["items"]; n = len(items)
    # 타임라인 라인의 좌우 시작/끝 여백을 대칭으로 맞춤(1.25"씩) — 마지막 노드가 우측 마진 안쪽에
    # 오도록 끝점을 당겨, 마지막 빈 원 노드/카드가 우측 끝에 붙던 비대칭을 제거(slide17 low).
    ax0, ax1 = int(Inches(1.25)), int(Inches(12.08))
    # 축을 약간 위로 올려(4.3→4.08) 타임라인이 하단부에 몰리지 않게 하고 상하 균형을 맞춤.
    axy = int(Inches(4.08))
    # 상단 요약 배지(완료/예정)를 제목 줄 우측(0.66")에 정렬해 고립감을 줄이고 제목과 한 라인에 둠.
    n_done = sum(1 for it in items if it.get("done"))
    _chip(slide, Inches(9.45), Inches(0.66), f"완료 {n_done}", fill=ACCENT, fg=NAVY,
          w=Inches(1.45), h=Inches(0.44), size=13)
    _chip(slide, Inches(11.05), Inches(0.66), f"예정 {n - n_done}", fill=PANEL, fg=MUTED,
          w=Inches(1.45), h=Inches(0.44), size=13)
    _rect(slide, Emu(ax0), Emu(axy), Emu(ax1 - ax0), Pt(3), LINE)
    step = (ax1 - ax0) // max(1, n - 1) if n > 1 else 0
    rr = int(Inches(0.16))
    for i, it in enumerate(items):
        cx = ax0 + i * step if n > 1 else ax0 + (ax1 - ax0) // 2
        done = it.get("done", False)
        above = (i % 2 == 0)
        # 노드↔라벨 사이를 가는 회색 스템으로 연결해 위/아래 라벨이 떠 보이지 않게 함
        stem_h = int(Inches(0.5))
        if above:
            _rect(slide, Emu(cx - int(Pt(0.9))), Emu(axy - rr - stem_h), Emu(int(Pt(1.8))), Emu(stem_h),
                  RGBColor(0xC7, 0xCF, 0xDA))
        else:
            _rect(slide, Emu(cx - int(Pt(0.9))), Emu(axy + rr), Emu(int(Pt(1.8))), Emu(stem_h),
                  RGBColor(0xC7, 0xCF, 0xDA))
        if done:
            _rect(slide, Emu(cx - rr), Emu(axy - rr + int(Pt(1.5))), Emu(2 * rr), Emu(2 * rr),
                  ACCENT, shape=MSO_SHAPE.OVAL)
        else:
            _rect(slide, Emu(cx - rr), Emu(axy - rr + int(Pt(1.5))), Emu(2 * rr), Emu(2 * rr),
                  WHITE, line=ACCENT, line_w=Pt(2.25), shape=MSO_SHAPE.OVAL)
        # 박스 폭을 약간 키워(2.3→2.5") 마지막 노드 라벨이 우측 가장자리에서 빠듯하게 줄바꿈되는
        # 것을 완화. 단 우측 끝 노드는 박스가 캔버스(12.7")를 넘지 않게 좌측으로 당겨 배치.
        box_w = int(Inches(2.5))
        bx = cx - box_w // 2
        right_edge = int(Inches(12.78))
        if bx + box_w > right_edge:
            bx = right_edge - box_w
        if bx < int(Inches(0.3)):
            bx = int(Inches(0.3))
        # 축을 4.08"로 올려 위 라벨(BOTTOM)·아래 라벨(TOP)을 노드 스템에 맞춰 재배치(상하 균형).
        if above:
            ty = int(Inches(2.28)); bh = Inches(1.32)
        else:
            ty = int(Inches(4.44)); bh = Inches(1.32)
        # 노드 라벨 뒤에 옅은 카드 박스를 깔아 타임라인 상·하단 빈 띠를 채우고 시각 밀도를 높임
        # (완료=연회색 강조, 예정=연회색 패널). 카드 상단에 가는 강조 바로 단계 위계를 줌.
        card_h = int(Inches(1.16))
        cy = ty if above else ty + (int(bh) - card_h)
        card_fill = RGBColor(0xF4, 0xF6, 0xF8) if done else PANEL
        _rect(slide, Emu(bx), Emu(cy), Emu(box_w), Emu(card_h), card_fill,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(slide, Emu(bx), Emu(cy), Emu(box_w), Inches(0.07),
              ACCENT if done else RGBColor(0xC2, 0xCC, 0xDC))
        # 예정(미완료) 카드는 본문 라벨을 한 단계 진한 슬레이트로 처리해 완료(NAVY) 대비 위계는
        # 유지하되 회색 배경 위 저대비를 해소(slide17 low — 예정 카드 본문이 회색 위 회색으로 흐리게
        # 보이던 인상 제거). 카드 상단 바·테두리만 회색으로 두고 텍스트는 가독성을 확보한다.
        lbl_col = NAVY if done else RGBColor(0x5A, 0x64, 0x74)
        sub_col = MUTED if done else RGBColor(0x73, 0x7D, 0x8C)
        lines = [[(it.get("date", ""), dict(size=12.5, bold=True, color=(ACCENT if done else MUTED)))],
                 [(it.get("label", ""), dict(size=11.5, bold=True, color=lbl_col))]]
        if it.get("sub"):
            lines.append([(it["sub"], dict(size=9.5, color=sub_col))])
        # 카드 안쪽에 맞춰 라벨을 수직 중앙 정렬(카드 영역에 균형 있게)
        _text(slide, Emu(bx + int(Inches(0.1))), Emu(cy), Emu(box_w - int(Inches(0.2))),
              Emu(card_h), lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(0.6), Inches(6.55), Inches(8), Inches(0.34),
          [[("● 완료 / 진행 중", dict(size=10.5, bold=True, color=ACCENT)),
            ("       ○ 예정", dict(size=10.5, bold=True, color=MUTED))]])
    if d.get("caption"):
        _text(slide, Inches(7.0), Inches(6.55), Inches(5.7), Inches(0.34),
              [[(d["caption"], dict(size=10.5, color=MUTED))]], align=PP_ALIGN.RIGHT)
    _footer(slide, idx)


def s_cards(slide, d, idx):
    """과제별 카드(컬럼) 시각화: 상단 배지(P1/P2/P3)+과제명 / 대표 썸네일(figures 원본) /
    아래 3줄(문제·접근·결과, 라벨은 틸·결과 줄 강조). 글(표)만 빽빽하던 'P1·P2·P3 한눈에'를
    스캔 가능한 3카드로 대체한다(이미지 가공 없이 원본 썸네일을 카드 안에 비율유지로 배치)."""
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    cards = d["cards"][:3]
    n = len(cards)
    # 3카드를 캔버스 가로(0.7~12.7")에 균등 배치. 카드 사이 거터 0.34".
    x0 = int(Inches(0.7)); total_w = int(Inches(12.0)); gap = int(Inches(0.34))
    cw = (total_w - gap * (n - 1)) // n
    # 카드 세로 슬롯: 제목 아래(1.92")~note 위. note 가 있으면 하단을 끌어올린다.
    cy0 = int(Inches(1.92))
    cbottom = int(Inches(6.46)) if d.get("note") else int(Inches(6.48))
    ch = cbottom - cy0
    badge_h = int(Inches(0.62))           # 상단 배지(P번호+과제명) 띠
    thumb_h = int(Inches(1.88))           # 대표 썸네일 영역
    rows_top = cy0 + badge_h + thumb_h    # 3줄(문제·접근·결과) 시작 y
    rows_h = (cy0 + ch) - rows_top
    for i, cd in enumerate(cards):
        x = x0 + i * (cw + gap)
        # 카드 외곽: 옅은 패널 + 또렷한 테두리(영역 확정)
        _rect(slide, Emu(x), Emu(cy0), Emu(cw), Emu(ch), WHITE,
              line=RGBColor(0xC4, 0xD0, 0xE2), line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # ── 상단 배지 띠(네이비): P번호(틸 강조) + 과제명(흰색) ──
        _rect(slide, Emu(x), Emu(cy0), Emu(cw), Emu(badge_h), NAVY,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 라운드 하단 모서리를 사각으로 덮어 띠 아래가 카드 본문과 자연스럽게 이어지게 함
        _rect(slide, Emu(x), Emu(cy0 + badge_h // 2), Emu(cw), Emu(badge_h // 2), NAVY)
        _rect(slide, Emu(x + int(Inches(0.22))), Emu(cy0 + (badge_h - int(Inches(0.42))) // 2),
              Emu(int(Inches(0.7))), Emu(int(Inches(0.42))), ACCENT,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, Emu(x + int(Inches(0.22))), Emu(cy0),
              Emu(int(Inches(0.7))), Emu(badge_h),
              [[(cd["no"], dict(size=17, bold=True, color=WHITE))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Emu(x + int(Inches(1.04))), Emu(cy0),
              Emu(cw - int(Inches(1.2))), Emu(badge_h),
              [[(cd["task"], dict(size=12.5, bold=True, color=WHITE))]],
              anchor=MSO_ANCHOR.MIDDLE)
        # ── 대표 썸네일(figures 원본, 비율유지) ──
        ty = cy0 + badge_h
        # 옅은 틴트 타일 위에 비율유지로 배치(저대비 원본이 '빈 칸'으로 읽히지 않게 프레임 부여)
        tpad = int(Inches(0.14))
        _rect(slide, Emu(x + tpad), Emu(ty + tpad // 2), Emu(cw - 2 * tpad),
              Emu(thumb_h - tpad), RGBColor(0xE4, 0xEA, 0xF4),
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        inset = int(Inches(0.20))
        _img_fit(slide, cd["thumb"], Emu(x + tpad + inset), Emu(ty + tpad // 2 + inset // 2),
                 Emu(cw - 2 * tpad - 2 * inset), Emu(thumb_h - tpad - inset), frame=False)
        # ── 아래 3줄(문제·접근·결과) — 라벨은 틸, 결과 줄 강조 ──
        labels = [("Problem", "problem"), ("Approach", "approach"), ("Impact", "result")]
        rgap = int(Inches(0.10))
        bottom_pad = int(Inches(0.18))
        rh = min(int(Inches(0.52)), (rows_h - bottom_pad - rgap * (len(labels) - 1)) // len(labels))
        rx = x + int(Inches(0.2)); rw = cw - int(Inches(0.4))
        for k, (lab, key) in enumerate(labels):
            ry = rows_top + k * (rh + rgap)
            is_result = (key == "result")
            if is_result:
                # 결과 줄은 옅은 틸 틴트 배경 + 좌측 틸 바로 강조
                _rect(slide, Emu(rx - int(Inches(0.08))), Emu(ry), Emu(rw + int(Inches(0.16))),
                      Emu(rh), RGBColor(0xF4, 0xF6, 0xF8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                _rect(slide, Emu(rx - int(Inches(0.08))), Emu(ry), Emu(int(Inches(0.08))),
                      Emu(rh), ACCENT)
            # 라벨 칩(틸 텍스트) + 본문
            label_w = int(Inches(0.96))
            row_pad = int(Inches(0.08))
            _text(slide, Emu(rx), Emu(ry + row_pad), Emu(label_w), Emu(rh - 2 * row_pad),
                  [[(lab, dict(size=11, bold=True, color=ACCENT))]],
                  anchor=MSO_ANCHOR.TOP)
            body_col = NAVY if is_result else INK
            _text(slide, Emu(rx + label_w + int(Inches(0.04))), Emu(ry + row_pad),
                  Emu(rw - label_w - int(Inches(0.04))), Emu(rh - 2 * row_pad),
                  [[(cd[key], dict(size=11.5, bold=is_result, color=body_col))]],
                  anchor=MSO_ANCHOR.TOP, wrap=True)
    if d.get("note"):
        _text(slide, Inches(0.7), Emu(int(Inches(6.52))), Inches(12.0), Inches(0.5),
              [[(d["note"], dict(size=11.5, color=RGBColor(0x3A, 0x48, 0x5C)))]],
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    _footer(slide, idx)


def _ba_row(slide, x, y, w, h, label, before, after):
    """before → after 한 행 도식: [라벨] [회색 before 박스] →(틸 화살표) [틸 after 박스].
    '왜 좋은가'를 글 도식으로 보여준다(이미지 가공 없이 PowerPoint 도형만)."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    lblw = int(Inches(0.92))
    # 좌측 카테고리 라벨(틸 굵게)
    _text(slide, Emu(x), Emu(y), Emu(lblw), Emu(h),
          [[(label, dict(size=12, bold=True, color=ACCENT))]], anchor=MSO_ANCHOR.MIDDLE)
    arr_w = int(Inches(0.46))
    box_w = (w - lblw - arr_w) // 2
    # before 박스(회색)
    bx = x + lblw
    _rect(slide, Emu(bx), Emu(y), Emu(box_w), Emu(h), PANEL,
          line=RGBColor(0xC4, 0xD0, 0xE2), line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(slide, Emu(bx + int(Inches(0.12))), Emu(y), Emu(box_w - int(Inches(0.24))), Emu(h),
          [[(before, dict(size=11.5, color=RGBColor(0x55, 0x60, 0x75)))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 화살표(틸)
    ax = bx + box_w
    _rect(slide, Emu(ax + int(Inches(0.06))), Emu(y + (h - int(Inches(0.3))) // 2),
          Emu(arr_w - int(Inches(0.12))), Emu(int(Inches(0.3))), ACCENT,
          shape=MSO_SHAPE.RIGHT_ARROW)
    # after 박스(틸 틴트 + 좌측 틸 바)
    ax2 = ax + arr_w
    _rect(slide, Emu(ax2), Emu(y), Emu(box_w), Emu(h), RGBColor(0xF4, 0xF6, 0xF8),
          line=ACCENT, line_w=Pt(1.1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _rect(slide, Emu(ax2), Emu(y), Emu(int(Inches(0.08))), Emu(h), ACCENT)
    _text(slide, Emu(ax2 + int(Inches(0.16))), Emu(y), Emu(box_w - int(Inches(0.26))), Emu(h),
          [[(after, dict(size=11.5, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _tech_card(slide, x, y, w, h, d):
    """기법 1개를 '무엇을 무엇으로 바꾸는지' 보이는 카드로 표현."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    _rect(slide, Emu(x), Emu(y), Emu(w), Emu(h), WHITE,
          line=RGBColor(0xC4, 0xD0, 0xE2), line_w=Pt(1.15), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _rect(slide, Emu(x), Emu(y), Emu(int(Inches(0.08))), Emu(h), ACCENT)
    _text(slide, Emu(x + int(Inches(0.18))), Emu(y + int(Inches(0.10))),
          Emu(w - int(Inches(0.36))), Emu(int(Inches(0.28))),
          [[(d["name"], dict(size=13, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Emu(x + int(Inches(0.18))), Emu(y + int(Inches(0.39))),
          Emu(w - int(Inches(0.36))), Emu(int(Inches(0.24))),
          [[(d["role"], dict(size=9.5, bold=True, color=ACCENT))]], anchor=MSO_ANCHOR.MIDDLE)

    fy = y + int(Inches(0.72)); fh = int(Inches(0.36))
    node_gap = int(Inches(0.16)); arr_w = int(Inches(0.26))
    node_w = (w - int(Inches(0.42)) - node_gap * 2 - arr_w) // 2
    bx = x + int(Inches(0.20))
    _rect(slide, Emu(bx), Emu(fy), Emu(node_w), Emu(fh), PANEL,
          line=LINE, line_w=Pt(0.9), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(slide, Emu(bx + int(Inches(0.06))), Emu(fy), Emu(node_w - int(Inches(0.12))), Emu(fh),
          [[(d["before"], dict(size=8.8, bold=True, color=RGBColor(0x55, 0x60, 0x75)))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ax = bx + node_w + node_gap
    _rect(slide, Emu(ax), Emu(fy + int(Inches(0.08))), Emu(arr_w), Emu(int(Inches(0.20))),
          ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
    bx2 = ax + arr_w + node_gap
    _rect(slide, Emu(bx2), Emu(fy), Emu(node_w), Emu(fh), RGBColor(0xF4, 0xF6, 0xF8),
          line=ACCENT, line_w=Pt(0.9), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(slide, Emu(bx2 + int(Inches(0.06))), Emu(fy), Emu(node_w - int(Inches(0.12))), Emu(fh),
          [[(d["after"], dict(size=8.8, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _text(slide, Emu(x + int(Inches(0.22))), Emu(y + int(Inches(1.16))),
          Emu(w - int(Inches(0.44))), Emu(h - int(Inches(1.24))),
          [[(d["desc"], dict(size=9.2, color=INK))]], anchor=MSO_ANCHOR.TOP)


def s_pipeline(slide, d, idx):
    """데이터 파이프라인: WHY(필요성) + 기술 스택(생성/뷰어 구분) + '왜 좋은가' before→after 도식.
    좌측 = WHY 박스 + 두 스택 카드(생성 파이프라인 / 운영 뷰어), 우측 = before→after 비교 박스
    (변환·저장·뷰어 3행) + 논문 수식 이미지 2장(작게 배치). 모두 실제 코드 사실만 사용."""
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    top = int(Inches(1.82))
    bottom = int(Inches(6.92))
    # ── WHY 박스(상단 전폭) ──
    why = d.get("why")
    wy_h = int(Inches(0.86)) if why else 0
    if why:
        wy = top
        _rect(slide, Emu(int(Inches(0.7))), Emu(wy), Emu(int(Inches(12.0))), Emu(wy_h),
              RGBColor(0xF2, 0xF5, 0xFA), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(slide, Emu(int(Inches(0.7))), Emu(wy), Emu(int(Inches(0.1))), Emu(wy_h), ACCENT)
        _text(slide, Emu(int(Inches(0.92))), Emu(wy), Emu(int(Inches(1.35))), Emu(wy_h),
              [[("Rationale", dict(size=12, bold=True, color=ACCENT))],
               [("Need", dict(size=9.5, bold=True, color=MUTED))]],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Emu(int(Inches(2.26))), Emu(wy), Emu(int(Inches(10.18))), Emu(wy_h),
              [[(why, dict(size=12.5, color=INK))]], anchor=MSO_ANCHOR.MIDDLE)
    if d.get("tech_cards"):
        col_top = top + wy_h + (int(Inches(0.18)) if why else 0)
        lx = int(Inches(0.7)); lw = int(Inches(7.45)); gap = int(Inches(0.22))
        rx = lx + lw + int(Inches(0.36)); rw = int(Inches(12.7)) - rx
        _text(slide, Emu(lx), Emu(col_top), Emu(lw), Emu(int(Inches(0.30))),
              [[("기법별 역할", dict(size=13, bold=True, color=NAVY))]],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Emu(rx), Emu(col_top), Emu(rw), Emu(int(Inches(0.30))),
              [[("Before → After", dict(size=13, bold=True, color=NAVY))]],
              anchor=MSO_ANCHOR.MIDDLE)

        grid_top = col_top + int(Inches(0.42))
        card_w = (lw - gap) // 2
        card_h = int(Inches(1.58))
        row_gap = int(Inches(0.20))
        for ti, tc in enumerate(d["tech_cards"][:4]):
            r, c = divmod(ti, 2)
            tx = lx + c * (card_w + gap)
            ty = grid_top + r * (card_h + row_gap)
            _tech_card(slide, tx, ty, card_w, card_h, tc)

        bas = d.get("ba_rows", [])
        ba_top = grid_top
        ba_gap = int(Inches(0.18))
        ba_h = (bottom - ba_top - ba_gap * (len(bas) - 1)) // max(1, len(bas))
        for bi, ba in enumerate(bas):
            by = ba_top + bi * (ba_h + ba_gap)
            _ba_row(slide, Emu(rx), Emu(by), Emu(rw), Emu(ba_h),
                    ba["label"], ba["before"], ba["after"])
        _footer(slide, idx)
        return
    # ── 좌/우 2단: 좌=기술 스택 카드 2개, 우=before→after 도식 ──
    col_top = top + wy_h + (int(Inches(0.18)) if why else 0)
    col_h = bottom - col_top
    lx = int(Inches(0.7)); lw = int(Inches(5.85)); colgap = int(Inches(0.3))
    rx = lx + lw + colgap; rw = int(Inches(12.7)) - rx
    # 좌측 컬럼 헤더
    _text(slide, Emu(lx), Emu(col_top), Emu(lw), Emu(int(Inches(0.32))),
          [[("Actual stack", dict(size=13, bold=True, color=NAVY))]],
          anchor=MSO_ANCHOR.MIDDLE)
    stacks = d.get("stacks", [])
    sy0 = col_top + int(Inches(0.42))
    sgap = int(Inches(0.16))
    sh = (col_h - int(Inches(0.42)) - sgap * (len(stacks) - 1)) // max(1, len(stacks))
    for si, st in enumerate(stacks):
        sy = sy0 + si * (sh + sgap)
        _rect(slide, Emu(lx), Emu(sy), Emu(lw), Emu(sh), WHITE,
              line=RGBColor(0xC4, 0xD0, 0xE2), line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 헤더 띠(틸) — 스택 이름
        hh = int(Inches(0.42))
        _rect(slide, Emu(lx), Emu(sy), Emu(lw), Emu(hh), ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(slide, Emu(lx), Emu(sy + hh // 2), Emu(lw), Emu(hh // 2), ACCENT)
        _text(slide, Emu(lx + int(Inches(0.2))), Emu(sy), Emu(lw - int(Inches(0.4))), Emu(hh),
              [[(st["head"], dict(size=12.5, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
        # 본문 항목(① ② …) — 핵심어 굵게(**...**) 지원
        _bullets_tf(slide, Emu(lx + int(Inches(0.22))), Emu(sy + hh + int(Inches(0.10))),
                    Emu(lw - int(Inches(0.44))), Emu(sh - hh - int(Inches(0.18))),
                    st["items"], size=st.get("size", 11), gap=3, line_spacing=1.06)
    # 우측 컬럼 헤더
    _text(slide, Emu(rx), Emu(col_top), Emu(rw), Emu(int(Inches(0.32))),
          [[("Before → After", dict(size=13, bold=True, color=NAVY))]],
          anchor=MSO_ANCHOR.MIDDLE)
    bas = d.get("ba_rows", [])
    # before→after 3행 도식
    ba_top = col_top + int(Inches(0.42))
    # 하단에 논문 수식 이미지 2장(작게) 배치 → before→after 행 영역과 이미지 영역을 나눔
    imgs = d.get("formula_imgs", [])
    # 논문 수식 썸네일 밴드를 1.42→1.62" 로 키워(약 +14% 크게 렌더) 작아서 단독 판독이 어렵던
    # 두 보조 도식(hex→Grade·palette 저장식)을 '보조 도식'으로 분명히 보이게 한다(slide5 low,
    # 이미지 가공 없이 배치 영역만 확대). before→after 3행은 줄어든 만큼 줄간격으로 흡수.
    img_band_h = int(Inches(1.62)) if imgs else 0
    ba_area_h = (bottom - ba_top) - (img_band_h + int(Inches(0.16)) if imgs else 0)
    ba_gap = int(Inches(0.16))
    ba_h = (ba_area_h - ba_gap * (len(bas) - 1)) // max(1, len(bas))
    for bi, ba in enumerate(bas):
        by = ba_top + bi * (ba_h + ba_gap)
        _ba_row(slide, Emu(rx), Emu(by), Emu(rw), Emu(ba_h),
                ba["label"], ba["before"], ba["after"])
    if imgs:
        iy = bottom - img_band_h
        igap = int(Inches(0.22))
        iw = (rw - igap * (len(imgs) - 1)) // len(imgs)
        for ii, im in enumerate(imgs):
            ix = rx + ii * (iw + igap)
            # 논문 수식 이미지: 흰 카드 + 또렷한 테두리, 캡션 한 줄
            cap_h = int(Inches(0.28))
            _rect(slide, Emu(ix), Emu(iy), Emu(iw), Emu(img_band_h - cap_h),
                  WHITE, line=RGBColor(0xC4, 0xD0, 0xE2), line_w=Pt(1))
            _img_fit(slide, im["src"], Emu(ix + int(Inches(0.06))), Emu(iy + int(Inches(0.06))),
                     Emu(iw - int(Inches(0.12))), Emu(img_band_h - cap_h - int(Inches(0.12))),
                     frame=False)
            _text(slide, Emu(ix), Emu(iy + img_band_h - cap_h), Emu(iw), Emu(cap_h),
                  [[("Reference / ", dict(size=10, bold=True, color=ACCENT)),
                    (im.get("caption", ""), dict(size=10.5, bold=True, color=RGBColor(0x3A,0x48,0x5C)))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def _draw_hardneg(slide, X, Y, W, H):
    """Paper-style hard-negative filtering diagram.
    q 주변의 inner exclusion zone, hard-negative band, easy-negative outer region을
    embedding geometry로 표현한다."""
    def PX(f): return int(X + f * W)
    def PY(f): return int(Y + f * H)
    NAVY = RGBColor(0x0F, 0x1E, 0x3D); INK = RGBColor(0x22, 0x26, 0x2E)
    RED = RGBColor(0xCC, 0x33, 0x28); BLUE = RGBColor(0x2B, 0x66, 0xD9)
    ORANGE = RGBColor(0xE0, 0x8A, 0x1E); GRAYL = RGBColor(0xCF, 0xD4, 0xDB)
    SOFT = RGBColor(0xE2, 0xE7, 0xEF); GRID = RGBColor(0xEC, 0xF0, 0xF5)
    GREEN = RGBColor(0x2B, 0xA6, 0x6B); MUTED = RGBColor(0x6F, 0x78, 0x86)
    ARRC = RGBColor(0x94, 0x9E, 0xAD)
    M = int(Inches(0.13))

    def rect(fx, fy, fw, fh, fill, line=None, shape=MSO_SHAPE.RECTANGLE, lw=0.8):
        return _rect(slide, Emu(PX(fx)), Emu(PY(fy)), Emu(int(fw * W)), Emu(int(fh * H)),
                     fill, line=line, shape=shape, line_w=Pt(lw))

    def txt(fx, fy, fw, fh, text, size=8.5, bold=False, color=INK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, name=FONT):
        _text(slide, Emu(PX(fx)), Emu(PY(fy)), Emu(int(fw * W)), Emu(int(fh * H)),
              [[(text, dict(size=size, bold=bold, color=color, name=name))]],
              align=align, anchor=anchor)

    def dashed(sp):
        ln = sp.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))

    def line(fx1, fy1, fx2, fy2, color=ARRC, width=0.9, dash=False):
        cn = slide.shapes.add_connector(1, Emu(PX(fx1)), Emu(PY(fy1)), Emu(PX(fx2)), Emu(PY(fy2)))
        cn.line.color.rgb = color
        cn.line.width = Pt(width)
        cn.shadow.inherit = False
        if dash:
            dashed(cn)
        return cn

    def mk(fx, fy, shp, fill, line_c=None, size=1.0, lw=1.2):
        mm = int(M * size)
        return _rect(slide, Emu(PX(fx) - mm // 2), Emu(PY(fy) - mm // 2), Emu(mm), Emu(mm),
                     fill, line=line_c, shape=shp, line_w=Pt(lw))

    # Panel and light embedding axes.
    rect(0.02, 0.02, 0.96, 0.95, RGBColor(0xFA, 0xFB, 0xFD), line=SOFT, lw=0.9)
    txt(0.055, 0.055, 0.40, 0.06, "embedding space around anchor a", size=8.5,
        bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    txt(0.625, 0.055, 0.28, 0.055, "s(a,x)=cos(a,x)", size=7.8, bold=True,
        color=MUTED, align=PP_ALIGN.RIGHT, name="Consolas")

    # Coordinate frame.
    line(0.11, 0.69, 0.75, 0.69, color=GRID, width=0.8)
    line(0.18, 0.18, 0.18, 0.76, color=GRID, width=0.8)
    txt(0.835, 0.705, 0.08, 0.04, "z₁", size=7.5, color=MUTED)
    txt(0.125, 0.165, 0.08, 0.04, "z₂", size=7.5, color=MUTED)

    qx, qy = 0.36, 0.50
    # Rings centered on the anchor. Outer ring = retrieved hard-neighbor candidates, including p+ and false negatives.
    # Inner ring = positive-neighbor margin where same-pattern false negatives are excluded.
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(PX(qx - 0.165)), Emu(PY(qy - 0.185)),
                                   Emu(int(0.33 * W)), Emu(int(0.37 * H)))
    inner.fill.background(); inner.line.color.rgb = RED; inner.line.width = Pt(1.1)
    inner.shadow.inherit = False; dashed(inner)
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(PX(qx - 0.300)), Emu(PY(qy - 0.330)),
                                   Emu(int(0.60 * W)), Emu(int(0.66 * H)))
    outer.fill.background(); outer.line.color.rgb = ARRC; outer.line.width = Pt(1.25)
    outer.shadow.inherit = False; dashed(outer)
    txt(qx - 0.16, qy - 0.23, 0.30, 0.04, "positive-neighbor margin", size=7.3,
        bold=True, color=RED, align=PP_ALIGN.LEFT)
    txt(qx + 0.13, qy + 0.205, 0.34, 0.04, "semi-hard annulus: keep", size=7.3,
        bold=True, color=BLUE, align=PP_ALIGN.LEFT)

    # Similarity spokes; too-hard points are deliberately near the anchor/p+.
    for fx, fy, col, dash in [
        (0.47, 0.37, ARRC, False),   # positive
        (0.43, 0.42, RED, True),     # too hard
        (0.51, 0.47, RED, True),     # too hard
        (0.15, 0.33, BLUE, False),   # ambiguous
        (0.13, 0.64, BLUE, False),   # ambiguous
        (0.72, 0.64, ORANGE, False), # easy
    ]:
        line(qx, qy, fx, fy, color=col if dash else ARRC, width=0.85, dash=dash)

    mk(qx, qy, MSO_SHAPE.OVAL, NAVY, size=1.25)
    txt(qx - 0.075, qy + 0.025, 0.06, 0.04, "a", size=11, bold=True, color=NAVY)
    mk(0.47, 0.37, MSO_SHAPE.OVAL, WHITE, line_c=INK, size=1.15, lw=1.4)
    txt(0.495, 0.325, 0.08, 0.04, "p⁺", size=11, bold=True, color=NAVY)

    # Too-hard false negatives placed closest to the anchor/positive.
    mk(0.43, 0.42, MSO_SHAPE.DIAMOND, RED, size=1.05)
    mk(0.51, 0.47, MSO_SHAPE.DIAMOND, RED, size=1.05)
    txt(0.535, 0.455, 0.12, 0.035, "false neg", size=7.2, bold=True, color=RED)
    # Useful hard negatives in the middle band.
    mk(0.15, 0.33, MSO_SHAPE.ISOSCELES_TRIANGLE, BLUE, size=1.10)
    mk(0.13, 0.64, MSO_SHAPE.ISOSCELES_TRIANGLE, BLUE, size=1.10)
    mk(0.61, 0.53, MSO_SHAPE.ISOSCELES_TRIANGLE, BLUE, size=1.10)
    # Too easy far negatives.
    mk(0.72, 0.64, MSO_SHAPE.RECTANGLE, ORANGE, size=1.05)
    mk(0.82, 0.50, MSO_SHAPE.RECTANGLE, ORANGE, size=1.05)

    # Direct legend as compact paper-style labels.
    legend = [
        (MSO_SHAPE.OVAL, WHITE, INK, "positive p⁺"),
        (MSO_SHAPE.DIAMOND, RED, None, "too hard / false neg"),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, BLUE, None, "semi-hard / keep"),
        (MSO_SHAPE.RECTANGLE, ORANGE, None, "too easy"),
    ]
    lx, ly = 0.70, 0.13
    for i, (shp, fill, line_c, label) in enumerate(legend):
        yy = ly + i * 0.075
        mk(lx, yy, shp, fill, line_c=line_c, size=0.80, lw=1.0)
        txt(lx + 0.035, yy - 0.020, 0.23, 0.04, label, size=7.5, color=INK,
            align=PP_ALIGN.LEFT)

    # Loss role strip: what positive/negative samples do mathematically.
    role_boxes = [
        (0.065, 0.275, GREEN, "positive p⁺", "numerator term\nattractive pull"),
        (0.365, 0.275, RED, "negative n", "denominator term\nrepulsive term"),
        (0.665, 0.280, BLUE, "semi-hard n", "outside margin\nboundary pressure"),
    ]
    for x0, w0, col, head, body in role_boxes:
        rect(x0, 0.790, w0, 0.140, WHITE, line=RGBColor(0xD3, 0xDA, 0xE3), lw=0.9)
        rect(x0, 0.790, 0.012, 0.140, col)
        txt(x0 + 0.025, 0.800, w0 - 0.045, 0.030, head, size=6.9, bold=True,
            color=col, align=PP_ALIGN.LEFT)
        txt(x0 + 0.025, 0.835, w0 - 0.045, 0.078, body, size=5.8, bold=True,
            color=INK, align=PP_ALIGN.LEFT)


def _draw_numba(slide, X, Y, W, H):
    """학술 발표용 텍스트 시각화: Numba가 Python loop overhead를 줄이는 원리."""
    BG = RGBColor(0xF7, 0xF7, 0xF7)
    CARD = RGBColor(0xFF, 0xFF, 0xFF)
    BD = RGBColor(0x99, 0x99, 0x99)
    SOFT = RGBColor(0xD4, 0xD4, 0xD4)
    PY_BG = RGBColor(0xF1, 0xF1, 0xF1)
    PY_LINE = RGBColor(0x8A, 0x8A, 0x8A)
    JIT_BG = RGBColor(0xE9, 0xE9, 0xE9)
    CORE_BG = RGBColor(0xEF, 0xEF, 0xEF)
    INK = RGBColor(0x1C, 0x1C, 0x1C)
    MUTED_INK = RGBColor(0x5F, 0x5F, 0x5F)
    DARK = RGBColor(0x5A, 0x5A, 0x5A)
    ARROW = RGBColor(0x9A, 0x9A, 0x9A)
    def PX(f): return int(X + f * W)
    def PY(f): return int(Y + f * H)

    def box(cx, cy, cw, ch, fill=CARD, line=SOFT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, lw=0.8):
        return _rect(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
                     fill, line=line, shape=shape, line_w=Pt(lw))

    def txt(cx, cy, cw, ch, text, size=9, bold=False, color=INK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, name=FONT):
        _text(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
              [[(text, dict(size=size, bold=bold, color=color, name=name))]],
              align=align, anchor=anchor)

    def arrow(cx, cy, cw, ch, color=ARROW):
        _rect(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
              color, shape=MSO_SHAPE.RIGHT_ARROW)

    def down_arrow(cx, cy, cw, ch, color=ARROW):
        _rect(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
              color, shape=MSO_SHAPE.DOWN_ARROW)

    box(0.02, 0.02, 0.96, 0.95, BG, line=BD, shape=MSO_SHAPE.RECTANGLE, lw=0.9)

    # Before: Python runtime appears inside the per-pixel loop.
    box(0.055, 0.07, 0.89, 0.34, CARD, line=SOFT)
    box(0.08, 0.10, 0.16, 0.08, PY_BG, line=PY_LINE)
    txt(0.08, 0.102, 0.16, 0.076, "BEFORE", size=12.0, bold=True, color=INK)
    txt(0.27, 0.095, 0.40, 0.08, "Interpreter overhead inside loop", size=12.0, bold=True, color=INK)
    txt(0.70, 0.10, 0.18, 0.08, "N times", size=15, bold=True, color=PY_LINE)

    nodes = [
        (0.085, 0.245, 0.085, 0.085, "pixel", CARD, SOFT),
        (0.205, 0.225, 0.105, 0.125, "PY", PY_BG, PY_LINE),
        (0.350, 0.245, 0.085, 0.085, "pixel", CARD, SOFT),
        (0.470, 0.225, 0.105, 0.125, "PY", PY_BG, PY_LINE),
        (0.615, 0.245, 0.085, 0.085, "...", CARD, SOFT),
        (0.735, 0.225, 0.105, 0.125, "PY", PY_BG, PY_LINE),
    ]
    for nx, ny, nw, nh, label, fill, line in nodes:
        box(nx, ny, nw, nh, fill, line=line, shape=MSO_SHAPE.RECTANGLE)
        txt(nx, ny, nw, nh, label, size=12.0, bold=(label == "PY"),
            color=(PY_LINE if label == "PY" else INK), name="Consolas")
    for ax in [0.172, 0.312, 0.438, 0.577, 0.703]:
        arrow(ax, 0.268, 0.030, 0.035, color=ARROW)
    txt(0.20, 0.352, 0.64, 0.035, "per-pixel interpreter dispatch", size=12.0, bold=True, color=MUTED_INK)

    # After: the first call builds native code, then that code executes the loop.
    box(0.055, 0.47, 0.89, 0.43, CARD, line=SOFT)
    box(0.08, 0.50, 0.16, 0.08, JIT_BG, line=BD)
    txt(0.08, 0.502, 0.16, 0.076, "AFTER", size=12.0, bold=True, color=INK)
    txt(0.265, 0.495, 0.42, 0.08, "First call JIT-compiles typed loop", size=12.0, bold=True, color=INK)

    box(0.095, 0.590, 0.185, 0.070, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT, shape=MSO_SHAPE.RECTANGLE)
    txt(0.100, 0.594, 0.175, 0.062, "input dtype", size=12.0, bold=True, color=INK)
    arrow(0.292, 0.612, 0.040, 0.032)
    box(0.345, 0.580, 0.175, 0.090, JIT_BG, line=BD, shape=MSO_SHAPE.RECTANGLE)
    txt(0.350, 0.584, 0.165, 0.080, "JIT\ncompile", size=12.0, bold=True, color=INK)
    arrow(0.535, 0.612, 0.040, 0.032)
    box(0.575, 0.580, 0.290, 0.095, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT, shape=MSO_SHAPE.RECTANGLE)
    txt(0.585, 0.584, 0.270, 0.087, "native loop\nreused", size=11.5, bold=True, color=INK)

    box(0.075, 0.735, 0.205, 0.075, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT, shape=MSO_SHAPE.RECTANGLE)
    txt(0.080, 0.739, 0.195, 0.067, "input data", size=12.0, bold=True, color=INK)
    arrow(0.292, 0.756, 0.035, 0.032)
    box(0.315, 0.700, 0.430, 0.145, DARK, line=DARK, shape=MSO_SHAPE.RECTANGLE)
    txt(0.330, 0.715, 0.400, 0.115, "native loop", size=16.0, bold=True, color=WHITE)
    arrow(0.760, 0.756, 0.040, 0.032)
    box(0.815, 0.735, 0.100, 0.075, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT, shape=MSO_SHAPE.RECTANGLE)
    txt(0.815, 0.739, 0.100, 0.067, "result", size=12.0, bold=True, color=INK)


def _draw_pyvips(slide, X, Y, W, H):
    """학술 발표용 텍스트 시각화: libvips streaming/lazy evaluation 설명."""
    BG = RGBColor(0xF7, 0xF7, 0xF7)
    CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
    LINE_C = RGBColor(0x99, 0x99, 0x99)
    SOFT_LINE = RGBColor(0xD4, 0xD4, 0xD4)
    HOT = RGBColor(0xF1, 0xF1, 0xF1)
    HOT_LINE = RGBColor(0x8A, 0x8A, 0x8A)
    COOL = RGBColor(0xE9, 0xE9, 0xE9)
    COOL_LINE = RGBColor(0xA8, 0xA8, 0xA8)
    STRIPE = RGBColor(0x6A, 0x6A, 0x6A)
    INK = RGBColor(0x1C, 0x1C, 0x1C)
    MUTED_INK = RGBColor(0x5F, 0x5F, 0x5F)
    ARROW = RGBColor(0x9A, 0x9A, 0x9A)
    def PX(f): return int(X + f * W)
    def PY(f): return int(Y + f * H)

    def box(cx, cy, cw, ch, fill=CARD_BG, line=SOFT_LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, lw=0.8):
        return _rect(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
                     fill, line=line, shape=shape, line_w=Pt(lw))

    def txt(cx, cy, cw, ch, text, size=9, bold=False, color=INK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, name=FONT):
        _text(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
              [[(text, dict(size=size, bold=bold, color=color, name=name))]],
              align=align, anchor=anchor)

    def arrow(cx, cy, cw, ch, color=ARROW):
        _rect(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
              color, shape=MSO_SHAPE.RIGHT_ARROW)

    def down_arrow(cx, cy, cw, ch, color=ARROW):
        _rect(slide, Emu(PX(cx)), Emu(PY(cy)), Emu(int(cw * W)), Emu(int(ch * H)),
              color, shape=MSO_SHAPE.DOWN_ARROW)

    box(0.02, 0.02, 0.96, 0.95, BG, line=LINE_C, shape=MSO_SHAPE.RECTANGLE, lw=0.9)

    # Full-buffer path: one wide row avoids cramped nested boxes.
    box(0.055, 0.085, 0.89, 0.285, CARD_BG, line=SOFT_LINE)
    box(0.080, 0.115, 0.145, 0.075, HOT, line=HOT_LINE)
    txt(0.080, 0.118, 0.145, 0.068, "PILLOW", size=12.0, bold=True, color=INK)
    txt(0.255, 0.120, 0.225, 0.060, "full-frame buffer", size=12.0, bold=True, color=MUTED_INK)

    box(0.095, 0.222, 0.090, 0.072, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT_LINE, shape=MSO_SHAPE.RECTANGLE)
    for i in range(3):
        y = 0.236 + i * 0.014
        _rect(slide, Emu(PX(0.108)), Emu(PY(y)), Emu(int(0.064 * W)), Emu(int(0.006 * H)),
              RGBColor(0xB8, 0xB8, 0xB8), shape=MSO_SHAPE.RECTANGLE)
    txt(0.075, 0.298, 0.130, 0.040, "FILE", size=12.0, bold=True, color=INK)
    arrow(0.205, 0.258, 0.035, 0.032)
    box(0.260, 0.210, 0.145, 0.100, HOT, line=HOT_LINE, shape=MSO_SHAPE.RECTANGLE)
    for i in range(5):
        y = 0.225 + i * 0.014
        _rect(slide, Emu(PX(0.275)), Emu(PY(y)), Emu(int(0.115 * W)), Emu(int(0.007 * H)),
              STRIPE, shape=MSO_SHAPE.RECTANGLE)
    txt(0.230, 0.314, 0.205, 0.035, "full RGB frame", size=10.2, bold=True, color=INK)
    arrow(0.425, 0.258, 0.035, 0.032)
    box(0.485, 0.207, 0.210, 0.128, HOT, line=HOT_LINE, shape=MSO_SHAPE.RECTANGLE)
    txt(0.490, 0.214, 0.200, 0.034, "full-image", size=12.0, bold=True, color=INK)
    txt(0.500, 0.258, 0.180, 0.066, "decode\nencode", size=12.0, bold=True, color=MUTED_INK)
    arrow(0.710, 0.258, 0.035, 0.032)
    box(0.790, 0.222, 0.090, 0.072, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT_LINE, shape=MSO_SHAPE.RECTANGLE)
    for i in range(4):
        y = 0.235 + i * 0.012
        _rect(slide, Emu(PX(0.803)), Emu(PY(y)), Emu(int(0.064 * W)), Emu(int(0.006 * H)),
              STRIPE, shape=MSO_SHAPE.RECTANGLE)
    txt(0.765, 0.298, 0.140, 0.040, "output", size=12.0, bold=True, color=INK)

    # libvips path: one wide row with a single C pipeline block.
    box(0.055, 0.455, 0.89, 0.315, CARD_BG, line=SOFT_LINE)
    box(0.080, 0.485, 0.160, 0.075, COOL, line=COOL_LINE)
    txt(0.080, 0.488, 0.160, 0.068, "LIBVIPS", size=12.0, bold=True, color=INK)
    txt(0.255, 0.490, 0.225, 0.060, "active-stripe buffer", size=12.0, bold=True, color=MUTED_INK)

    box(0.095, 0.606, 0.090, 0.072, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT_LINE, shape=MSO_SHAPE.RECTANGLE)
    for i in range(3):
        y = 0.620 + i * 0.014
        _rect(slide, Emu(PX(0.108)), Emu(PY(y)), Emu(int(0.064 * W)), Emu(int(0.006 * H)),
              RGBColor(0xB8, 0xB8, 0xB8), shape=MSO_SHAPE.RECTANGLE)
    txt(0.075, 0.690, 0.130, 0.040, "FILE", size=12.0, bold=True, color=INK)
    arrow(0.205, 0.642, 0.035, 0.035)
    box(0.260, 0.592, 0.145, 0.100, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT_LINE, shape=MSO_SHAPE.RECTANGLE)
    for i in range(5):
        y = 0.607 + i * 0.014
        fill = STRIPE if i == 2 else RGBColor(0xD8, 0xD8, 0xD8)
        _rect(slide, Emu(PX(0.275)), Emu(PY(y)), Emu(int(0.115 * W)), Emu(int(0.007 * H)),
              fill, shape=MSO_SHAPE.RECTANGLE)
    txt(0.245, 0.696, 0.175, 0.035, "RGB stripe", size=10.8, bold=True, color=INK)
    arrow(0.425, 0.642, 0.035, 0.035)
    box(0.485, 0.585, 0.210, 0.150, COOL, line=COOL_LINE, shape=MSO_SHAPE.RECTANGLE)
    txt(0.490, 0.588, 0.200, 0.060, "streaming\nC pipeline", size=9.4, bold=True, color=INK)
    txt(0.500, 0.655, 0.180, 0.050, "decode\nencode", size=10.4, bold=True, color=MUTED_INK)
    arrow(0.710, 0.642, 0.035, 0.035)
    box(0.790, 0.606, 0.090, 0.072, RGBColor(0xF4, 0xF4, 0xF4), line=SOFT_LINE, shape=MSO_SHAPE.RECTANGLE)
    for i in range(4):
        y = 0.619 + i * 0.012
        _rect(slide, Emu(PX(0.803)), Emu(PY(y)), Emu(int(0.064 * W)), Emu(int(0.006 * H)),
              STRIPE, shape=MSO_SHAPE.RECTANGLE)
    txt(0.765, 0.690, 0.140, 0.040, "output", size=12.0, bold=True, color=INK)

    box(0.145, 0.840, 0.710, 0.070, RGBColor(0xEA, 0xEA, 0xEA), line=SOFT_LINE, shape=MSO_SHAPE.RECTANGLE)
    txt(0.145, 0.846, 0.710, 0.058, "same output with bounded active rows", size=12.0, bold=True, color=INK)


def s_papertext(slide, d, idx):
    """논문 figure(Fig.1/2)의 텍스트 내용을 네이티브 텍스트 박스로 렌더(이미지 캡쳐 X).
    각 fig = 얇은 테두리 박스 + (선택 head) + (라벨 줄 + 들여쓴 monospace 내용) + 하단 Fig 캡션."""
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    top = int(Inches(2.0))
    if d.get("bullets"):
        bh = int(Inches(d.get("body_h", 0.7)))
        marker_color = ACCENT
        _bullets_tf(slide, Inches(0.7), Emu(top), Inches(12.0), Emu(bh),
                    d["bullets"], size=d.get("bullet_size", 14), gap=5, marker=marker_color)
        top = top + bh + int(Inches(0.14))
    figs = d["figs"]
    n = len(figs)
    cols = d.get("cols", 2 if n >= 2 else 1)
    rows = (n + cols - 1) // cols
    bottom = int(Inches(6.95))
    area_h = bottom - top
    cap_h = 0 if d.get("hide_fig_captions") else int(Inches(0.36))
    gx = int(Inches(0.4)); gy = int(Inches(0.28))
    x0 = int(Inches(0.7)); total_w = int(Inches(12.0))
    cw = (total_w - gx * (cols - 1)) // cols
    chh = (area_h - gy * (rows - 1)) // rows
    box_h = chh - cap_h
    MONO = "Consolas"
    def _kr(t): return any('가' <= ch <= '힣' for ch in t)
    for i, fg in enumerate(figs):
        r, c = divmod(i, cols)
        fx = x0 + c * (cw + gx)
        fy = top + r * (chh + gy)
        if fg.get("diagram") == "hardneg":
            # 반장 레이아웃: wafer 후보, loss 포함/제외, embedding 결과를 한 셀 안에 요약
            _draw_hardneg_half(slide, Emu(fx), Emu(fy), Emu(cw), Emu(box_h))
        elif fg.get("diagram") == "simclr_big":
            _draw_simclr_big(slide, Emu(fx), Emu(fy), Emu(cw), Emu(box_h))
        elif fg.get("diagram") == "pos_neg_hard":
            _draw_pos_neg_hard(slide, Emu(fx), Emu(fy), Emu(cw), Emu(box_h))
        elif fg.get("diagram") == "numba":
            _draw_numba(slide, fx, fy, cw, box_h)
        elif fg.get("diagram") == "pyvips":
            _draw_pyvips(slide, fx, fy, cw, box_h)
        elif fg.get("src"):
            # 이미지 figure 는 테두리 박스 없이 셀을 거의 가득 채워 최대 크기로 렌더한다
            # (테두리+padding 이 이미지를 줄이던 문제 제거 — 사용자 요청). 텍스트 박스만 테두리 유지.
            _img_fit(slide, fg["src"], Emu(fx + int(Inches(0.04))), Emu(fy + int(Inches(0.04))),
                     Emu(cw - int(Inches(0.08))), Emu(box_h - int(Inches(0.08))), frame=False)
        else:
            _rect(slide, Emu(fx), Emu(fy), Emu(cw), Emu(box_h), WHITE,
                  line=RGBColor(0x33, 0x37, 0x40), line_w=Pt(1.25))
            lines = []
            if fg.get("head"):
                lines.append([(fg["head"], dict(size=14.5, bold=True, color=NAVY))])
                lines.append([(" ", dict(size=5, color=WHITE))])
            for si, sec in enumerate(fg["sections"]):
                lab, content = sec[0], sec[1]
                lines.append([(lab, dict(size=13.5, bold=True, italic=(not _kr(lab)),
                                         color=NAVY, name=(FONT if _kr(lab) else MONO)))])
                for cc in content:
                    lines.append([("    " + cc, dict(size=13, color=INK,
                                                     name=(FONT if _kr(cc) else MONO)))])
                if si < len(fg["sections"]) - 1:
                    lines.append([(" ", dict(size=5.5, color=WHITE))])
            _text(slide, Emu(fx + int(Inches(0.30))), Emu(fy + int(Inches(0.14))),
                  Emu(cw - int(Inches(0.54))), Emu(box_h - int(Inches(0.28))),
                  lines, anchor=MSO_ANCHOR.MIDDLE)
        if cap_h:
            caption_color = RGBColor(0x1C, 0x1C, 0x1C) if d.get("grayscale") else NAVY
            _text(slide, Emu(fx), Emu(fy + box_h + int(Inches(0.03))), Emu(cw), Emu(cap_h),
                  [[(fg["caption"], dict(size=12.5, bold=True, color=caption_color))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx, size=d.get("footer_size", 9))


def s_archflow(slide, d, idx):
    """전체 아키텍쳐 흐름을 네이티브 도형(박스+화살표)으로 그린다(이미지 캡쳐 X).
    stage = 전폭 박스 또는 split(나란한 2개 박스). 사이에 down-arrow. 테두리 검정, 화살표 네이비."""
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    def I(v): return int(Inches(v))
    def _kr(t): return any('가' <= ch <= '힣' for ch in t)
    MONO = "Consolas"
    BD = RGBColor(0x33, 0x37, 0x40)
    ARR = NAVY
    TINT = RGBColor(0xF2, 0xF4, 0xF8)
    stages = d["stages"]
    n = len(stages)
    cx = int(EMU_W) // 2; full_w = I(9.8); full_x = cx - full_w // 2
    top = I(1.95); bottom = I(6.95)
    arrow_h = I(0.32)

    def box_h(b):
        if b.get("lines"):
            return I(0.36) + I(0.28) * len(b["lines"])
        if b.get("sub"):
            return I(0.74)
        return I(0.52)

    def stage_h(st):
        if st.get("split"):
            return max(box_h(b) for b in st["split"])
        return box_h(st)

    def draw_box(bx, by, bw, bh, b, tint=False):
        _rect(slide, Emu(bx), Emu(by), Emu(bw), Emu(bh),
              (TINT if tint else WHITE), line=BD, line_w=Pt(1.4), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        if b.get("lines"):
            _text(slide, Emu(bx + I(0.2)), Emu(by + I(0.05)), Emu(bw - I(0.4)), Emu(I(0.28)),
                  [[(b["text"], dict(size=12, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)
            ly = by + I(0.36)
            for ln in b["lines"]:
                _text(slide, Emu(bx + I(0.32)), Emu(ly), Emu(bw - I(0.46)), Emu(I(0.26)),
                      [[(ln, dict(size=10.5, color=INK, name=(FONT if _kr(ln) else MONO)))]],
                      anchor=MSO_ANCHOR.MIDDLE)
                ly += I(0.28)
        elif b.get("sub"):
            _text(slide, Emu(bx), Emu(by + I(0.09)), Emu(bw), Emu(I(0.30)),
                  [[(b["text"], dict(size=12.5, bold=True, color=NAVY))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _text(slide, Emu(bx), Emu(by + I(0.40)), Emu(bw), Emu(I(0.28)),
                  [[(b["sub"], dict(size=10.5, color=MUTED))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            _text(slide, Emu(bx), Emu(by), Emu(bw), Emu(bh),
                  [[(b["text"], dict(size=12.5, bold=True, color=NAVY))]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    heights = [stage_h(s) for s in stages]
    total = sum(heights) + arrow_h * (n - 1)
    y = top + max(0, (bottom - top - total) // 2)
    for i, st in enumerate(stages):
        bh = heights[i]
        if st.get("split"):
            boxes = st["split"]; m = len(boxes); gap = I(0.5)
            sw = (full_w - gap * (m - 1)) // m
            for j, b in enumerate(boxes):
                draw_box(full_x + j * (sw + gap), y, sw, bh, b, tint=True)
        else:
            draw_box(full_x, y, full_w, bh, st)
        if i < n - 1:
            ay = y + bh
            _rect(slide, Emu(cx - I(0.12)), Emu(ay + I(0.03)), Emu(I(0.24)), Emu(arrow_h - I(0.06)),
                  ARR, shape=MSO_SHAPE.DOWN_ARROW)
            if st.get("arrow"):
                _text(slide, Emu(cx + I(0.28)), Emu(ay), Emu(I(4.4)), Emu(arrow_h),
                      [[(st["arrow"], dict(size=10, color=MUTED))]], anchor=MSO_ANCHOR.MIDDLE)
        y = y + bh + arrow_h
    _footer(slide, idx)


def _metric_card_compact(slide, x, y, w, h, big, label, sub="", color=COVER_BAR):
    _rect(slide, x, y, w, h, PANEL, line=LINE)
    _rect(slide, x, y, w, Inches(0.08), color)
    _text(slide, x+Inches(0.12), y+Inches(0.14), Emu(int(w)-int(Inches(0.24))), Inches(0.32),
          [[(big, dict(size=20, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER)
    _text(slide, x+Inches(0.12), y+Inches(0.48), Emu(int(w)-int(Inches(0.24))), Inches(0.24),
          [[(label, dict(size=10.5, bold=True, color=INK))]], align=PP_ALIGN.CENTER)
    if sub:
        _text(slide, x+Inches(0.12), y+Inches(0.72), Emu(int(w)-int(Inches(0.24))), Inches(0.18),
              [[(sub, dict(size=8.6, color=MUTED))]], align=PP_ALIGN.CENTER)


def _mini_title(slide, x, y, w, text, color=COVER_BAR):
    _rect(slide, x, y, w, Inches(0.36), PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _rect(slide, x, y, Inches(0.08), Inches(0.36), color)
    _text(slide, x+Inches(0.16), y, Emu(int(w)-int(Inches(0.24))), Inches(0.36),
          [[(text, dict(size=12.5, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _image_pair(slide, x, y, w, h, img1, img2, lab1, lab2):
    gap = Inches(0.10)
    iw = Emu((int(w) - int(gap)) // 2)
    _img_fit(slide, img1, x, y, iw, h-Inches(0.32), frame=True)
    _img_fit(slide, img2, Emu(int(x)+int(iw)+int(gap)), y, iw, h-Inches(0.32), frame=True)
    _text(slide, x, Emu(int(y)+int(h)-int(Inches(0.28))), iw, Inches(0.26),
          [[(lab1, dict(size=10.5, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER)
    _text(slide, Emu(int(x)+int(iw)+int(gap)), Emu(int(y)+int(h)-int(Inches(0.28))), iw, Inches(0.26),
          [[(lab2, dict(size=10.5, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER)


def _evidence_table(slide, x, y, w, h, headers, rows, col_fracs,
                    highlight_rows=None, font_size=9.0, recipe_col=1,
                    header_size=9.2, row_gap_pt=0.6):
    """Compact evidence table for portfolio-metric slides."""
    highlight_rows = set(highlight_rows or [])
    total_w = int(w)
    total_h = int(h)
    fr_sum = float(sum(col_fracs))
    widths = [int(total_w * f / fr_sum) for f in col_fracs]
    widths[-1] += total_w - sum(widths)
    header_h = int(Inches(0.42))
    row_h = int((total_h - header_h) / max(1, len(rows)))

    cx = int(x)
    for j, head in enumerate(headers):
        cw = widths[j]
        _rect(slide, Emu(cx), y, Emu(cw), Emu(header_h), NAVY, line=WHITE)
        _text(slide, Emu(cx + int(Inches(0.04))), y, Emu(cw - int(Inches(0.08))), Emu(header_h),
              [[(str(head), dict(size=header_size, bold=True, color=WHITE))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw

    for i, row in enumerate(rows):
        cy = int(y) + header_h + i * row_h
        fill = RGBColor(0xE9, 0xF3, 0xED) if i in highlight_rows else (PANEL if i % 2 == 0 else WHITE)
        line = RGBColor(0xC8, 0xD8, 0xCE) if i in highlight_rows else RGBColor(0xE0, 0xE6, 0xEF)
        cx = int(x)
        for j, val in enumerate(row):
            cw = widths[j]
            _rect(slide, Emu(cx), Emu(cy), Emu(cw), Emu(row_h - int(Pt(row_gap_pt))), fill, line=line)
            align = PP_ALIGN.LEFT if j == recipe_col else PP_ALIGN.CENTER
            pad = int(Inches(0.06 if j == recipe_col else 0.035))
            _text(slide, Emu(cx + pad), Emu(cy + int(Inches(0.01))),
                  Emu(cw - 2 * pad), Emu(row_h - int(Inches(0.02))),
                  [[(str(val), dict(size=font_size, bold=(i in highlight_rows), color=NAVY if i in highlight_rows else INK))]],
                  align=align, anchor=MSO_ANCHOR.MIDDLE)
            cx += cw


def s_p1_known_perf(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P1 | Known evidence", "Known 16-class backbone scan and staged improvement")
    _text(slide, Inches(0.80), Inches(1.42), Inches(11.8), Inches(0.34),
          [[("portfolio.md 기준 성능값: 실전 현업 데이터 16 class / 1,500 labeled / 4:1 stratified split.",
             dict(size=13.2, bold=True, color=NAVY))]])

    headers = ["Stage", "Configuration", "Weighted F1", "Decision / note"]
    rows = [
        ("Baseline", "ImageNet pretrained CNN", "0.78", "starting point"),
        ("Backbone", "ViT", "0.81", "same split"),
        ("Backbone", "Swin", "0.84", "same split"),
        ("Backbone", "EffNetV2", "0.85", "same split"),
        ("Backbone", "MaxViT", "0.87", "119.5M params / 74.2G FLOPs"),
        ("Backbone", "ConvNeXtV2", "0.87", "selected: 88.6M params / 45.1G FLOPs"),
        ("HPO", "ConvNeXtV2 + Optuna", "0.92", "warmup→cosine schedule, focal loss, class-imbalance weighting, augmentation strength"),
        ("Cascade", "ConvNeXtV2 + Optuna + ROI-YOLO", "0.95", "low-confidence wafer only"),
    ]
    _evidence_table(slide, Inches(0.70), Inches(1.88), Inches(12.0), Inches(4.10),
                    headers, rows, [1.15, 2.90, 1.05, 4.05], highlight_rows={7},
                    font_size=10.0, header_size=10.0, recipe_col=1)

    card_y = Inches(6.18)
    cards = [
        ("0.78 → 0.87", "backbone gain"),
        ("0.87 → 0.92", "Optuna HPO"),
        ("0.92 → 0.95", "ROI-YOLO cascade"),
    ]
    cw = Inches(3.72); gap = Inches(0.36); x0 = Inches(0.92)
    for i, (big, lab) in enumerate(cards):
        x = Emu(int(x0) + i * (int(cw) + int(gap)))
        _metric_card_compact(slide, x, card_y, cw, Inches(0.78), big, lab, "", color=COVER_BAR)
    _footer(slide, idx)


def s_unknown_ablation(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P1 | Unknown evidence", "Contrastive recipe ablation on generated evaluation set")
    _text(slide, Inches(0.80), Inches(1.40), Inches(11.8), Inches(0.34),
          [[("portfolio.md 기준: 생성 데이터 개발 지표입니다. 실전 운영 검증(2,000장 review, 13 후보 group 중 7건 확인)과 분리해 제시합니다.",
             dict(size=12.6, bold=True, color=NAVY))]])

    headers = ["#", "Recipe (per class 500, normal 2000)", "M1 capture", "M2 noise", "M3 Completeness", "Sil"]
    rows = [
        ("1", "Global InfoNCE only (baseline)", "0.9337", "15.78%", "0.9468", "0.582"),
        ("2", "+ Local DenseCL (LW=0.5)", "0.9361", "13.87%", "0.9502", "0.514"),
        ("3", "+ MoCo Queue 4096", "0.9356", "9.45%", "0.9474", "0.573"),
        ("4", "+ NV-Retriever NEG 0.72", "0.9250", "8.23%", "0.9485", "0.611"),
        ("5", "+ NeCo 0.2 (5-tool full)", "0.9559", "6.66%", "0.9660", "0.6104"),
        ("6", "Final recipe (Local DenseCL 제외 4-tool)", "0.9559", "6.66%", "0.9660", "0.781"),
        ("7", "Final recipe + noise threshold τ=0.5", "0.9619", "0.00%", "0.9679", "0.781"),
    ]
    _evidence_table(slide, Inches(0.62), Inches(1.90), Inches(12.10), Inches(4.30),
                    headers, rows, [0.42, 4.80, 1.20, 1.05, 1.45, 0.82],
                    highlight_rows={6}, font_size=11.0, header_size=10.5, recipe_col=1)

    _rect(slide, Inches(0.82), Inches(6.38), Inches(11.75), Inches(0.34), PANEL, line=LINE)
    _rect(slide, Inches(0.82), Inches(6.38), Inches(0.10), Inches(0.34), COVER_BAR)
    _text(slide, Inches(1.00), Inches(6.38), Inches(11.35), Inches(0.34),
          [[("Final generated-set metric: capture 0.9619 / noise 0.00% / Completeness 0.9679 / Silhouette 0.781.",
             dict(size=11.0, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def s_p2_intro(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P2 | Chip multi-label", "Single-failure source to 2-combo evaluation setting")

    # Problem definition only: source와 eval을 분리해서 보여준다.
    x = Inches(0.65); y = Inches(1.82); w = Inches(3.45); h = Inches(4.88)
    _rect(slide, x, y, w, h, PANEL, line=LINE)
    _rect(slide, x, y, Inches(0.10), h, COVER_BAR)
    lines = [
        [("Problem", dict(size=15, bold=True, color=NAVY))],
        [("real 2-combo GT 부족\nlimited supervision / validation", dict(size=13.8, bold=True, color=NAVY))],
        [(" ", dict(size=4, color=WHITE))],
        [("Source", dict(size=15, bold=True, color=NAVY))],
        [("single-failure source만 보유\n직접적인 2-combo supervision 없음", dict(size=12.8, color=INK))],
        [(" ", dict(size=4, color=WHITE))],
        [("Eval", dict(size=15, bold=True, color=NAVY))],
        [("3,850-chip benchmark\n2-combo/negative tail 포함", dict(size=12.8, color=INK))],
    ]
    _text(slide, x+Inches(0.28), y+Inches(0.28), w-Inches(0.55), h-Inches(0.45), lines)

    def chip_tile(img, label, tx, ty, tw, th):
        _rect(slide, tx, ty, tw, th, WHITE, line=RGBColor(0xD6,0xDE,0xEA))
        _img_fit(slide, img, tx+Inches(0.06), ty+Inches(0.06), tw-Inches(0.12), th-Inches(0.34), frame=False)
        _text(slide, tx+Inches(0.04), ty+th-Inches(0.25), tw-Inches(0.08), Inches(0.18),
              [[(label, dict(size=9.4, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rx = Inches(4.42); rw = Inches(8.10)
    src_y = Inches(1.82); src_h = Inches(2.08)
    eval_y = Inches(4.34); eval_h = Inches(2.36)

    _rect(slide, rx, src_y, rw, src_h, WHITE, line=LINE)
    _rect(slide, rx, src_y, Inches(0.10), src_h, COVER_BAR)
    _text(slide, rx+Inches(0.24), src_y+Inches(0.16), rw-Inches(0.48), Inches(0.28),
          [[("Source bank: single-failure samples", dict(size=14.4, bold=True, color=NAVY))]])
    _text(slide, rx+Inches(0.24), src_y+Inches(0.48), rw-Inches(0.48), Inches(0.20),
          [[("train split 이후 train source 안에서만 synthetic combination을 생성합니다.",
             dict(size=10.4, color=MUTED))]])
    source_imgs = [
        ("chip_eval_bank_boundary_selected.png", "bb"),
        ("chip_eval_fork_selected.png", "fk"),
        ("chip_eval_scratch_selected.png", "sc"),
        ("chip_eval_scratch_rot_selected.png", "sr"),
    ]
    tile_w = Inches(1.40); tile_h = Inches(1.18); gap = Inches(0.30)
    tx0 = int(rx) + int(Inches(0.48)); ty = src_y + Inches(0.76)
    for i, (img, lab) in enumerate(source_imgs):
        chip_tile(img, lab, Emu(tx0 + i*(int(tile_w)+int(gap))), ty, tile_w, tile_h)

    _rect(slide, rx+Inches(2.95), Inches(3.96), Inches(1.60), Inches(0.28),
          COVER_BAR, shape=MSO_SHAPE.DOWN_ARROW)
    _text(slide, rx+Inches(4.72), Inches(3.95), Inches(2.7), Inches(0.28),
          [[("evaluation requires combo classes and negative tail", dict(size=10.0, bold=True, color=MUTED))]],
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    _rect(slide, rx, eval_y, rw, eval_h, WHITE, line=LINE)
    _rect(slide, rx, eval_y, Inches(0.10), eval_h, RGBColor(0x2B,0xA6,0x6B))
    _text(slide, rx+Inches(0.24), eval_y+Inches(0.16), rw-Inches(0.48), Inches(0.28),
          [[("Eval target: 3,850-chip controlled benchmark", dict(size=14.4, bold=True, color=NAVY))]])
    _text(slide, rx+Inches(0.24), eval_y+Inches(0.48), rw-Inches(0.48), Inches(0.20),
          [[("single 4 + 2-combo 6 + Normal/Invalid/OOD negative tail.",
             dict(size=10.4, color=MUTED))]])
    eval_imgs = [
        ("chip_combo_bb_fork_selected.png", "bb+fk"),
        ("chip_combo_bb_scratch_selected.png", "bb+sc"),
        ("chip_combo_fork_scratch_rot_selected.png", "fk+sr"),
        ("chip_eval_normal_selected.png", "normal"),
        ("chip_eval_invalid_selected.png", "invalid"),
    ]
    tile_w2 = Inches(1.26); tile_h2 = Inches(1.32); gap2 = Inches(0.17)
    tx0 = int(rx) + int(Inches(0.38)); ty = eval_y + Inches(0.82)
    for i, (img, lab) in enumerate(eval_imgs):
        chip_tile(img, lab, Emu(tx0 + i*(int(tile_w2)+int(gap2))), ty, tile_w2, tile_h2)

    _rect(slide, x+Inches(0.28), y+h-Inches(0.88), w-Inches(0.56), Inches(0.58),
          WHITE, line=RGBColor(0xD6,0xDE,0xEA))
    _text(slide, x+Inches(0.45), y+h-Inches(0.82), w-Inches(0.90), Inches(0.46),
          [[("Design implication: distribution gap requires probability modeling and reject logic.",
             dict(size=10.2, bold=True, color=NAVY))]],
          anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx, size=d.get("footer_size", 9))


def s_p2_fcmpm(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P2 | FCM-PM 원리", "FCM-PM: positive-probability gain with FAR-tail control")
    _text(slide, Inches(0.85), Inches(1.52), Inches(11.75), Inches(0.34),
          [[("FCM은 full-cover 합성으로 약한 2-combo positive probability를 끌어올리고, Pair Mask는 합성 배경 loss를 차단해 Normal/OOD FAR tail을 낮춥니다.",
             dict(size=13.2, color=INK))]])

    fig_x, fig_y, fig_w, fig_h = Inches(0.78), Inches(1.98), Inches(11.78), Inches(2.30)
    _rect(slide, fig_x, fig_y, fig_w, fig_h, WHITE, line=LINE)
    _img_fit(slide, "fcm_pm_panel.png", fig_x+Inches(0.10), fig_y+Inches(0.12),
             fig_w-Inches(0.20), fig_h-Inches(0.24), frame=False)

    card_y = Inches(4.52); card_h = Inches(1.05); card_w = Inches(3.70); gap = Inches(0.22); x0 = Inches(0.78)
    cards = [
        ("Normal baseline", "2-combo pos. prob 낮음\nFAR 낮지만 recall 약함"),
        ("FCM", "pos. prob 상승\nFAR tail 동반 상승 가능"),
        ("FCM-PM", "pos. prob 유지\nFAR tail 0.00%로 억제"),
    ]
    for i, (head, body) in enumerate(cards):
        x = Emu(int(x0) + i*(int(card_w)+int(gap)))
        fill = PANEL if i != 2 else RGBColor(0xED,0xF7,0xF2)
        line = LINE if i != 2 else RGBColor(0x8D,0xC8,0xA5)
        bar = COVER_BAR if i != 2 else RGBColor(0x2B,0xA6,0x6B)
        _rect(slide, x, card_y, card_w, card_h, fill, line=line)
        _rect(slide, x, card_y, Inches(0.09), card_h, bar)
        _text(slide, x+Inches(0.22), card_y+Inches(0.12), card_w-Inches(0.40), Inches(0.24),
              [[(head, dict(size=12.2, bold=True, color=NAVY))]])
        _text(slide, x+Inches(0.22), card_y+Inches(0.43), card_w-Inches(0.40), Inches(0.48),
              [[(body, dict(size=11.4, bold=(i==2), color=NAVY if i==2 else INK,
                            name=FONT))]],
              anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            ax = Emu(int(x) + int(card_w) + int(gap)//2 - int(Inches(0.10)))
            _rect(slide, ax, card_y+Inches(0.42), Inches(0.20), Inches(0.20),
                  RGBColor(0xC7,0xCF,0xDA), shape=MSO_SHAPE.RIGHT_ARROW)

    _rect(slide, Inches(1.02), Inches(5.82), Inches(11.22), Inches(0.62), WHITE, line=LINE)
    _rect(slide, Inches(1.02), Inches(5.82), Inches(0.10), Inches(0.62), COVER_BAR)
    _text(slide, Inches(1.22), Inches(5.88), Inches(1.56), Inches(0.18),
          [[("Algorithmic view", dict(size=9.8, bold=True, color=MUTED))]],
          anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(2.78), Inches(5.87), Inches(9.20), Inches(0.44),
          [[("FCM: x_mix = M * x_A + (1-M) * x_B,   y_mix = y_A ∪ y_B",
             dict(size=10.4, bold=True, color=NAVY, name="Cambria Math"))],
           [("Pair Mask: L_total = BCE(x_mix, y_mix) + λ BCE(x_mask, y_A / y_B)",
             dict(size=10.4, bold=True, color=NAVY, name="Cambria Math"))]],
          anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def s_p2_validation(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P2 | 검증 결과", "Multi-label recipe performance table (per class 2000)")
    _text(slide, Inches(0.72), Inches(1.40), Inches(12.0), Inches(0.32),
          [[("portfolio.md 기준: 현업 failure chip 원천 + 도메인 확률분포 기반 생성/검증. 단일 대표 모델은 FCM-PM + val_margin selection입니다.",
             dict(size=12.4, bold=True, color=NAVY))]])

    main_headers = ["#", "Recipe", "bit_F1", "single", "2combo", "FAR"]
    main_rows = [
        ("1", "BCE + Label Smoothing", "0.1093", "0.1896", "0.0668", "99.47%"),
        ("2", "Sigmoid Focal Loss", "0.7980", "0.8724", "0.7050", "45.72%"),
        ("3", "Asymmetric Loss (ASL)", "0.6435", "0.5379", "0.7320", "100%"),
        ("4", "CutMix (random rectangle)", "0.9359", "0.9566", "0.9070", "42.05%"),
        ("5", "CutMix + Pair Mask", "0.9491", "0.9728", "0.9281", "24.62%"),
        ("6", "FCM-PM + val_f1 selection", "0.9652", "1.0000", "0.9517", "0.15%"),
        ("7", "FCM-PM + val_margin selection", "0.9927", "0.9996", "0.9871", "0.00%"),
        ("8", "vote_majority_bits Ensemble", "0.9956", "1.0000", "0.9921", "0.00%"),
        ("9", "Knowledge Distillation", "0.9799", "1.0000", "0.9638", "0.00%"),
    ]
    _evidence_table(slide, Inches(0.42), Inches(1.82), Inches(8.05), Inches(4.58),
                    main_headers, main_rows,
                    [0.35, 2.85, 0.76, 0.72, 0.82, 0.68],
                    highlight_rows={6, 7}, font_size=9.2, header_size=9.2, recipe_col=1)

    side_headers = ["Candidate", "NI-FAR", "OOD-FAR", "Cost"]
    side_rows = [
        ("FCM-PM + val_margin", "0.00%", "0.00%", "1x / 1x / 1x"),
        ("vote_majority_bits", "0.00%", "0.00%", "5x / 1/5x / 5x"),
        ("Knowledge Distillation", "0.00%", "0.00%", "1x / 1x / 1x"),
    ]
    _rect(slide, Inches(8.74), Inches(1.82), Inches(4.08), Inches(0.34), PANEL, line=LINE)
    _text(slide, Inches(8.84), Inches(1.82), Inches(3.88), Inches(0.34),
          [[("Negative-tail and deployment cost", dict(size=11.0, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _evidence_table(slide, Inches(8.74), Inches(2.26), Inches(4.08), Inches(1.76),
                    side_headers, side_rows,
                    [1.65, 0.62, 0.70, 1.00],
                    highlight_rows={0, 1}, font_size=8.8, header_size=8.5, recipe_col=0)

    _rect(slide, Inches(8.74), Inches(4.30), Inches(4.08), Inches(1.58), PANEL, line=LINE)
    _rect(slide, Inches(8.74), Inches(4.30), Inches(0.10), Inches(1.58), COVER_BAR)
    _text(slide, Inches(8.96), Inches(4.48), Inches(3.62), Inches(0.28),
          [[("Portfolio interpretation", dict(size=12.0, bold=True, color=NAVY))]])
    lines = [
        [("FCM-PM + val_margin", dict(size=10.3, bold=True, color=NAVY))],
        [("대표 single model: bit_F1 0.9927 / Total FAR 0.00%", dict(size=9.8, color=INK))],
        [("Ensemble은 bit_F1 0.9956이지만 cost 5x", dict(size=9.8, color=INK))],
        [("KD는 1x cost 후보: bit_F1 0.9799 / FAR 0.00%", dict(size=9.8, color=INK))],
    ]
    _text(slide, Inches(8.96), Inches(4.82), Inches(3.62), Inches(0.82), lines)

    strip_y = Inches(6.52); strip_h = Inches(0.46); strip_x = Inches(0.64); strip_w = Inches(12.05)
    _rect(slide, strip_x, strip_y, strip_w, strip_h, PANEL, line=LINE)
    _rect(slide, strip_x, strip_y, Inches(0.10), strip_h, COVER_BAR)
    _text(slide, strip_x+Inches(0.18), strip_y, strip_w-Inches(0.36), strip_h,
          [[("대표 single model: bit_F1 0.9927 / Total FAR 0.00%. Champion ensemble: bit_F1 0.9956 / Total FAR 0.00%지만 cost는 5x / 1/5x / 5x.",
             dict(size=10.6, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx, size=d.get("footer_size", 9))


def s_p2_selection(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P2 | selection & reject", "val_margin selection and negative-tail control")

    # Left: val-margin selection
    lx = Inches(0.75); ly = Inches(1.78); lw = Inches(5.70); lh = Inches(4.95)
    _rect(slide, lx, ly, lw, lh, WHITE, line=LINE)
    _rect(slide, lx, ly, Inches(0.10), lh, COVER_BAR)
    _text(slide, lx+Inches(0.28), ly+Inches(0.22), lw-Inches(0.56), Inches(0.32),
          [[("val-margin checkpoint selection", dict(size=16.0, bold=True, color=NAVY))]])
    _text(slide, lx+Inches(0.28), ly+Inches(0.62), lw-Inches(0.56), Inches(0.25),
          [[("val_margin = mean(p_positive) - max(p_negative)", dict(size=12.0, bold=True, color=INK, name="Consolas"))]])

    chart_x = int(lx) + int(Inches(0.62)); chart_y = int(ly) + int(Inches(1.46))
    base_y = chart_y + int(Inches(1.55)); bar_w = int(Inches(0.16)); step = int(Inches(0.78))
    vals_f1 = [0.42, 0.52, 0.64, 0.86, 0.72]
    vals_margin = [0.18, 0.27, 0.39, 0.54, 0.82]
    vals_eval = [0.20, 0.30, 0.44, 0.57, 0.84]
    colors = [RGBColor(0xB8,0xC2,0xD0), RGBColor(0x2F,0x3A,0x4F), RGBColor(0x78,0xB8,0x8D)]
    SEL_F1 = RGBColor(0xD9,0x8C,0x1F)
    SEL_MARGIN = RGBColor(0x2B,0xA6,0x6B)
    labels = ["val-F1", "val-margin", "eval-F1"]
    for li, (lab, col) in enumerate(zip(labels, colors)):
        xx = int(lx) + int(Inches(0.60)) + li * int(Inches(1.36))
        _rect(slide, Emu(xx), ly+Inches(0.99), Inches(0.20), Inches(0.20), col)
        _text(slide, Emu(xx+int(Inches(0.26))), ly+Inches(0.955), Inches(1.02), Inches(0.24),
              [[(lab, dict(size=12.0, bold=True, color=INK))]], anchor=MSO_ANCHOR.MIDDLE)
    for i in range(5):
        bx = chart_x + i*step
        hf1 = int(Inches(vals_f1[i] * 1.12))
        hm = int(Inches(vals_margin[i] * 1.12))
        he = int(Inches(vals_eval[i] * 1.12))
        _rect(slide, Emu(bx), Emu(base_y-hf1), Emu(bar_w), Emu(hf1),
              SEL_F1 if i == 3 else colors[0])
        _rect(slide, Emu(bx+bar_w+int(Inches(0.06))), Emu(base_y-hm), Emu(bar_w), Emu(hm),
              SEL_MARGIN if i == 4 else colors[1])
        _rect(slide, Emu(bx+2*bar_w+int(Inches(0.12))), Emu(base_y-he), Emu(bar_w), Emu(he),
              SEL_MARGIN if i == 4 else (SEL_F1 if i == 3 else colors[2]))
        _text(slide, Emu(bx-int(Inches(0.02))), Emu(base_y+int(Inches(0.07))), Inches(0.58), Inches(0.16),
              [[(f"epoch {i+1}", dict(size=8.8, color=MUTED))]], align=PP_ALIGN.CENTER)
    _rect(slide, Emu(chart_x-int(Inches(0.08))), Emu(base_y), Inches(4.08), Pt(1), LINE)
    _text(slide, lx+Inches(3.28), ly+Inches(1.22), Inches(0.98), Inches(0.22),
          [[("val-F1 max", dict(size=10.2, bold=True, color=SEL_F1))]], align=PP_ALIGN.CENTER)
    _text(slide, lx+Inches(4.10), ly+Inches(1.26), Inches(1.06), Inches(0.22),
          [[("margin max", dict(size=10.2, bold=True, color=SEL_MARGIN))]], align=PP_ALIGN.CENTER)
    _text(slide, lx+Inches(0.34), ly+Inches(3.72), lw-Inches(0.68), Inches(0.20),
          [[("portfolio evidence: correlation with test bit_F1", dict(size=9.6, bold=True, color=MUTED))]],
          align=PP_ALIGN.CENTER)
    card_y = ly+Inches(3.92); card_h = Inches(0.66); card_w = Inches(2.42)
    impacts = [
        ("val_f1", "Spearman -0.10", "weak"),
        ("val_margin", "Spearman +0.56", "selected"),
    ]
    for ii, (head, val, delta) in enumerate(impacts):
        cx = Emu(int(lx)+int(Inches(0.30))+ii*(int(card_w)+int(Inches(0.25))))
        side_col = SEL_F1 if ii == 0 else SEL_MARGIN
        _rect(slide, cx, card_y, card_w, card_h, RGBColor(0xF4,0xF7,0xFB), line=RGBColor(0xD7,0xDE,0xE8))
        _rect(slide, cx, card_y, Inches(0.07), card_h, side_col)
        _text(slide, cx+Inches(0.14), card_y+Inches(0.06), Inches(0.70), Inches(0.18),
              [[(head, dict(size=9.6, bold=True, color=MUTED))]])
        _text(slide, cx+Inches(0.14), card_y+Inches(0.24), Inches(1.26), Inches(0.24),
              [[(val, dict(size=10.8, bold=True, color=NAVY))]])
        _text(slide, cx+Inches(1.44), card_y+Inches(0.19), Inches(0.82), Inches(0.25),
              [[(delta, dict(size=9.8, bold=True, color=side_col))]],
              align=PP_ALIGN.RIGHT)

    # Right: negative-tail control
    rx = Inches(6.85); ry = Inches(1.78); rw = Inches(5.70); rh = Inches(4.95)
    _rect(slide, rx, ry, rw, rh, WHITE, line=LINE)
    _rect(slide, rx, ry, Inches(0.10), rh, RGBColor(0x2B,0xA6,0x6B))
    _text(slide, rx+Inches(0.28), ry+Inches(0.22), rw-Inches(0.56), Inches(0.32),
          [[("Reject gate and deployment candidates", dict(size=15.2, bold=True, color=NAVY))]])
    _text(slide, rx+Inches(0.28), ry+Inches(0.60), rw-Inches(0.56), Inches(0.28),
          [[("max-prob < 0.55 → Normal gate; ensemble/KD는 배포 후보로 비교",
             dict(size=11.4, bold=True, color=INK))]])

    bit_labels = ["bb", "fk", "sc", "sr"]
    bit_colors = [RGBColor(0x2F,0x67,0xF6), RGBColor(0x18,0xA9,0x7D),
                  RGBColor(0xF5,0xA6,0x23), RGBColor(0xD9,0x3A,0x5A)]
    for bi, (lab, col) in enumerate(zip(bit_labels, bit_colors)):
        bx = int(rx) + int(Inches(0.34)) + bi * int(Inches(0.60))
        _rect(slide, Emu(bx), ry+Inches(0.99), Inches(0.20), Inches(0.20), col)
        _text(slide, Emu(bx+int(Inches(0.26))), ry+Inches(0.955), Inches(0.34), Inches(0.24),
              [[(lab, dict(size=12.0, bold=True, color=INK))]], anchor=MSO_ANCHOR.MIDDLE)

    panels = [
        ("single bb", [0.84, 0.12, 0.13, 0.11], "accept"),
        ("2-combo bb+sc", [0.62, 0.13, 0.60, 0.12], "accept"),
        ("tail / OOD", [0.55, 0.33, 0.31, 0.29], "reject"),
    ]
    px0 = int(rx) + int(Inches(0.42)); py0 = int(ry) + int(Inches(1.42))
    pw = int(Inches(1.42)); ph = int(Inches(1.46)); pg = int(Inches(0.24))
    for pi, (name, probs, tag) in enumerate(panels):
        px = px0 + pi * (pw + pg)
        _rect(slide, Emu(px), Emu(py0), Emu(pw), Emu(ph), RGBColor(0xF7,0xF9,0xFC), line=RGBColor(0xD7,0xDE,0xE8))
        _text(slide, Emu(px+int(Inches(0.08))), Emu(py0+int(Inches(0.06))), Emu(pw-int(Inches(0.16))), Inches(0.16),
              [[(name, dict(size=9.2, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER)
        base = py0 + int(Inches(1.08)); max_h = int(Inches(0.72)); bw = int(Inches(0.17)); gap = int(Inches(0.10))
        bars_w = 4 * bw + 3 * gap
        bx0 = px + (pw - bars_w)//2
        _rect(slide, Emu(bx0-int(Inches(0.04))), Emu(base), Emu(bars_w+int(Inches(0.08))), Pt(1), LINE)
        for bj, val in enumerate(probs):
            bh = int(max_h * val)
            bx = bx0 + bj * (bw + gap)
            _rect(slide, Emu(bx), Emu(base-bh), Emu(bw), Emu(bh), bit_colors[bj])
        _text(slide, Emu(px+int(Inches(0.08))), Emu(py0+int(Inches(1.16))), Emu(pw-int(Inches(0.16))), Inches(0.18),
              [[(tag, dict(size=9.2, bold=True, color=RGBColor(0x2B,0xA6,0x6B) if tag == "accept" else RGBColor(0xCC,0x33,0x28)))]],
              align=PP_ALIGN.CENTER)
    _text(slide, rx+Inches(0.34), ry+Inches(3.34), rw-Inches(0.68), Inches(0.20),
          [[("portfolio evidence: inference-stage robustness candidates", dict(size=9.6, bold=True, color=MUTED))]],
          align=PP_ALIGN.CENTER)
    card_y = ry+Inches(3.56); card_h = Inches(0.66); card_w = Inches(2.42)
    impacts = [
        ("max-prob gate", "< 0.55 → Normal", "tail control"),
        ("ensemble", "bit_F1 0.9956", "FAR 0.00%"),
    ]
    for ii, (head, val, delta) in enumerate(impacts):
        cx = Emu(int(rx)+int(Inches(0.30))+ii*(int(card_w)+int(Inches(0.25))))
        _rect(slide, cx, card_y, card_w, card_h, RGBColor(0xF4,0xF7,0xFB), line=RGBColor(0xD7,0xDE,0xE8))
        _rect(slide, cx, card_y, Inches(0.07), card_h, RGBColor(0x2B,0xA6,0x6B))
        _text(slide, cx+Inches(0.14), card_y+Inches(0.06), Inches(1.02), Inches(0.18),
              [[(head, dict(size=9.4, bold=True, color=MUTED))]])
        _text(slide, cx+Inches(0.14), card_y+Inches(0.24), Inches(1.30), Inches(0.24),
              [[(val, dict(size=10.6, bold=True, color=NAVY))]])
        _text(slide, cx+Inches(1.50), card_y+Inches(0.19), Inches(0.76), Inches(0.25),
              [[(delta, dict(size=9.6, bold=True, color=RGBColor(0x2B,0xA6,0x6B)))]],
              align=PP_ALIGN.RIGHT)
    _text(slide, rx+Inches(0.52), ry+Inches(4.34), rw-Inches(1.04), Inches(0.25),
          [[("Knowledge Distillation: bit_F1 0.9799 / Total FAR 0.00% / 1x cost",
             dict(size=9.2, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER)

    _footer(slide, idx)


def s_p3_intro(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P3 | Trend anomaly generator", "Real-label scarcity addressed by generator benchmark")
    _text(slide, Inches(0.85), Inches(1.86), Inches(6.2), Inches(0.52),
          [[("Core contribution — 10년 trend 판정 기준을 generator rule로 정식화해 validation benchmark를 구축했습니다.",
             dict(size=16, bold=True, color=NAVY))]])
    # left flow
    steps = [
        ("현장 기준", "BBD / Overlay / CD 10년 trend 판정 경험"),
        ("Generator rule", "Region 5종 + Noise 3종 + Anomaly 5종"),
        ("Synthetic benchmark", "normal 3,500 + abnormal 3,500 = 7,000개"),
        ("Gate validation", "Binary F1 0.9967 / abnormal recall 0.9987"),
    ]
    x = Inches(0.90); y = Inches(2.75); w = Inches(5.45); bh = Inches(0.68); gap = Inches(0.18)
    for i, (head, body) in enumerate(steps):
        yy = Emu(int(y) + i*(int(bh)+int(gap)))
        _rect(slide, x, yy, w, bh, RGBColor(0xE8,0xEC,0xF2) if i==len(steps)-1 else PANEL, line=LINE)
        _rect(slide, x, yy, Inches(0.09), bh, COVER_BAR)
        _text(slide, x+Inches(0.24), yy+Inches(0.08), Inches(1.45), Inches(0.25),
              [[(head, dict(size=11.2, bold=True, color=NAVY))]])
        _text(slide, x+Inches(1.72), yy, w-Inches(1.92), bh,
              [[(body, dict(size=12.4, bold=(i==len(steps)-1), color=NAVY if i==len(steps)-1 else INK))]],
              anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps)-1:
            _rect(slide, x+Inches(2.58), Emu(int(yy)+int(bh)-int(Inches(0.02))), Inches(0.24), Inches(0.22),
                  COVER_BAR, shape=MSO_SHAPE.DOWN_ARROW)

    # right metrics + image
    rx = Inches(6.85); ry = Inches(2.05)
    _metric_card_compact(slide, rx, ry, Inches(1.75), Inches(1.05), "7,000", "trend samples", "3,500+3,500")
    _metric_card_compact(slide, rx+Inches(1.95), ry, Inches(1.75), Inches(1.05), "5×3×5", "rule space", "region/noise/anomaly")
    _metric_card_compact(slide, rx+Inches(3.90), ry, Inches(1.75), Inches(1.05), "0.9967", "Binary F1", "PoC")
    _rect(slide, rx, Inches(3.34), Inches(5.65), Inches(3.18), WHITE, line=LINE)
    _rect(slide, rx, Inches(3.34), Inches(5.65), Inches(0.40), NAVY)
    _text(slide, rx+Inches(0.18), Inches(3.40), Inches(5.30), Inches(0.28),
          [[("Stabilization and shortcut-control stack", dict(size=12.6, bold=True, color=WHITE))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tune_cards = [
        ("LR sweep", "3e-5/3e-4 stable\n1e-4/1e-3 collapse"),
        ("Label smoothing", "0.02 sweet spot\nFN 4.6 -> 0.8"),
        ("Stochastic depth", "0.05 regularization\nFP 3.8 -> 1.4"),
        ("Focal + EMA", "gamma=2.0, EMA 0.95\nhard-anomaly stability"),
        ("Normal ratio", "hard-normal 3,300\nFP/FN balance"),
        ("Color + median + NT", "color/smoothing check\nval-F1 median, NT 0.9"),
    ]
    tx0, ty0 = int(rx) + int(Inches(0.22)), int(Inches(3.88))
    tw, th = int(Inches(2.54)), int(Inches(0.62))
    hgap, vgap = int(Inches(0.16)), int(Inches(0.13))
    for i, (head, body) in enumerate(tune_cards):
        cx = tx0 + (i % 2) * (tw + hgap)
        cy = ty0 + (i // 2) * (th + vgap)
        _rect(slide, Emu(cx), Emu(cy), Emu(tw), Emu(th), PANEL, line=LINE)
        _rect(slide, Emu(cx), Emu(cy), Emu(int(Inches(0.06))), Emu(th), COVER_BAR)
        _text(slide, Emu(cx+int(Inches(0.13))), Emu(cy+int(Inches(0.07))),
              Emu(tw-int(Inches(0.20))), Emu(int(Inches(0.17))),
              [[(head, dict(size=10.4, bold=True, color=NAVY))]])
        _text(slide, Emu(cx+int(Inches(0.13))), Emu(cy+int(Inches(0.26))),
              Emu(tw-int(Inches(0.20))), Emu(int(Inches(0.29))),
              [[(body, dict(size=8.8, color=MUTED))]])
    _text(slide, rx+Inches(0.15), Inches(6.24), Inches(5.35), Inches(0.22),
          [[("Interpretation: generator validity, training stability, and shortcut checks are evaluated together.",
             dict(size=9.6, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER)
    _footer(slide, idx)


def s_p3_generator(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P3 | generator rule", "Baseline generator: episode + noise sampling")
    _text(slide, Inches(0.85), Inches(1.47), Inches(11.6), Inches(0.32),
          [[("정상 trend는 region별 baseline 위에 episode density와 noise family를 분리해 샘플링합니다. anomaly injection은 다음 장에서 분리합니다.",
             dict(size=12.8, color=INK))]])

    # Normal baseline only: episode/density/noise sampling.
    fig_x, fig_y, fig_w, fig_h = Inches(0.78), Inches(1.82), Inches(11.75), Inches(3.82)
    _rect(slide, fig_x, fig_y, fig_w, fig_h, WHITE, line=LINE)
    _episode_trend_diagram(slide, fig_x+Inches(0.10), fig_y+Inches(0.08),
                           fig_w-Inches(0.20), fig_h-Inches(0.16))

    # Bottom rule cards: equations for how one synthetic normal trend is assembled.
    card_y = Inches(5.72); card_h = Inches(0.92); card_w = Inches(2.88); gap = Inches(0.18); x0 = Inches(0.78)
    cards = [
        ("Episode", "K ~ U(Kmin,Kmax),  L_k ~ U(Lmin,Lmax)", "구간 개수와 길이를 먼저 결정"),
        ("Density", "m_k ~ Cat(dense,sparse,missing,thin)", "구간별 관측 point 수를 결정"),
        ("Noise", "ε_t ∈ {N(0,σ²), AR(1), Laplace(0,b)}", "측정 산포 family를 episode별 부여"),
        ("Output", "y_t = μ_region(t) + ε_t  if active", "missing 구간 제외 후 normal trend 생성"),
    ]
    for i, (head, formula, desc) in enumerate(cards):
        x = Emu(int(x0) + i*(int(card_w)+int(gap)))
        _rect(slide, x, card_y, card_w, card_h, WHITE, line=LINE)
        _rect(slide, x, card_y, card_w, Inches(0.26), NAVY)
        _text(slide, x+Inches(0.10), card_y, card_w-Inches(0.20), Inches(0.27),
              [[(head, dict(size=10.4, bold=True, color=WHITE))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, x+Inches(0.10), card_y+Inches(0.34), card_w-Inches(0.20), Inches(0.24),
              [[(formula, dict(size=9.1, bold=True, color=NAVY, name="Cambria Math"))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, x+Inches(0.14), card_y+Inches(0.62), card_w-Inches(0.28), Inches(0.22),
              [[(desc, dict(size=9.2, color=INK, name=FONT))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def s_p3_result(slide, d, idx):
    _bg(slide, WHITE)
    _title_block_compact(slide, "P3 | PoC validation", "Representative anomaly rules and synthetic gate validation")
    _text(slide, Inches(0.85), Inches(1.50), Inches(11.5), Inches(0.34),
          [[("Normal baseline에 5가지 anomaly injection rule을 더해 normal/abnormal 1차 gate를 검증했습니다.",
             dict(size=14.2, color=INK))]])

    # 5 anomaly types in a 2-row figure grid. Keep trend panels image-dominant:
    # no per-card text inside the trend panels; the 6th cell carries labels/equations.
    x0 = Inches(1.15); y0 = Inches(1.86); cw = Inches(3.55); ch = Inches(1.90)
    gapx = Inches(0.18); gapy = Inches(0.14)
    imgs = [
        ("trend_mean_shift_view.png", "Mean shift", "y_t = x_t + δ"),
        ("trend_standard_deviation_view.png", "Std change", "y_t = μ + A(x_t - μ)"),
        ("trend_spike_view.png", "Spike", "y_t = x_t + Σ A_i 1(t=t_i)"),
        ("trend_drift_view.png", "Drift", "y_t = x_t + r(t-s)"),
        ("trend_context_view.png", "Context", "legend-relative shift"),
    ]
    for i, (src, head, formula) in enumerate(imgs):
        row, col = divmod(i, 3)
        x = Emu(int(x0)+col*(int(cw)+int(gapx)))
        y = Emu(int(y0)+row*(int(ch)+int(gapy)))
        _rect(slide, x, y, cw, ch, WHITE, line=LINE)
        _img_cover(slide, src, x+Inches(0.06), y+Inches(0.06),
                   cw-Inches(0.12), ch-Inches(0.12), focus_y=0.50)

    fx = Emu(int(x0)+2*(int(cw)+int(gapx)))
    fy = Emu(int(y0)+1*(int(ch)+int(gapy)))
    _rect(slide, fx, fy, cw, ch, PANEL, line=LINE)
    _rect(slide, fx, fy, Inches(0.10), ch, COVER_BAR)
    _text(slide, fx+Inches(0.24), fy+Inches(0.16), cw-Inches(0.42), Inches(0.24),
          [[("Injection equations", dict(size=13.0, bold=True, color=NAVY))]])
    eq_lines = [
        [("mean:  y_t = x_t + δ", dict(size=9.0, color=INK, name="Consolas"))],
        [("scale: y_t = μ + A(x_t - μ)", dict(size=9.0, color=INK, name="Consolas"))],
        [("spike: y_t = x_t + Σ A_i 1(t=t_i)", dict(size=9.0, color=INK, name="Consolas"))],
        [("drift: y_t = x_t + r(t-s)", dict(size=9.0, color=INK, name="Consolas"))],
        [("context: fleet/reference-relative shift", dict(size=9.0, color=INK, name="Consolas"))],
    ]
    _text(slide, fx+Inches(0.26), fy+Inches(0.48), cw-Inches(0.52), ch-Inches(0.58),
          eq_lines, anchor=MSO_ANCHOR.MIDDLE)

    # result strip
    strip_x, strip_y, strip_w, strip_h = Inches(0.90), Inches(5.90), Inches(11.55), Inches(0.50)
    _rect(slide, strip_x, strip_y, strip_w, strip_h, PANEL, line=LINE)
    vals = [
        ("Test 1,500", "normal 750 / abnormal 5종×150"),
        ("F1 0.9967", "TN/FN/FP/TP 746/1/4/749"),
        ("Recall 0.9987", "abnormal class"),
        ("5-seed best 0.9987", "TN/FN/FP/TP 748/0/2/750"),
        ("threshold 0.9", "normal gate"),
    ]
    cell_w = Emu(int(strip_w) // len(vals))
    for i, (big, sub) in enumerate(vals):
        x = Emu(int(strip_x) + i*int(cell_w))
        if i > 0:
            _rect(slide, x, strip_y+Inches(0.09), Pt(1), strip_h-Inches(0.18), LINE)
        _text(slide, x+Inches(0.08), strip_y+Inches(0.08), cell_w-Inches(0.16), Inches(0.22),
              [[(big, dict(size=11.5, bold=True, color=NAVY if i not in (2, 3) else RGBColor(0x2B,0xA6,0x6B)))]],
              align=PP_ALIGN.CENTER)
        _text(slide, x+Inches(0.08), strip_y+Inches(0.32), cell_w-Inches(0.16), Inches(0.20),
              [[(sub, dict(size=9.4, color=MUTED))]], align=PP_ALIGN.CENTER)

    _rect(slide, Inches(0.90), Inches(6.46), Inches(11.5), Inches(0.30), WHITE, line=LINE)
    _rect(slide, Inches(0.90), Inches(6.46), Inches(0.10), Inches(0.30), COVER_BAR)
    _text(slide, Inches(1.08), Inches(6.46), Inches(11.10), Inches(0.30),
          [[("Validation controls: LR sweep, color/smoothing, label smoothing, stochastic depth, focal+EMA, val-F1 median, normal threshold 0.9",
             dict(size=10.2, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def _number_chip(slide, x, y, n, label, body, w=Inches(2.55), h=Inches(0.72), fill=None):
    fill = fill or PANEL
    _rect(slide, x, y, w, h, fill, line=LINE)
    _rect(slide, x, y, Inches(0.09), h, ACCENT)
    _rect(slide, x+Inches(0.18), y+Inches(0.17), Inches(0.34), Inches(0.34), ACCENT,
          shape=MSO_SHAPE.OVAL)
    _text(slide, x+Inches(0.18), y+Inches(0.17), Inches(0.34), Inches(0.34),
          [[(str(n), dict(size=11, bold=True, color=WHITE))]], align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, x+Inches(0.62), y+Inches(0.10), w-Inches(0.78), Inches(0.22),
          [[(label, dict(size=11.2, bold=True, color=NAVY))]])
    _text(slide, x+Inches(0.62), y+Inches(0.36), w-Inches(0.78), h-Inches(0.42),
          [[(body, dict(size=9.7, color=MUTED))]])


def _draw_simclr_big(slide, x, y, w, h):
    """SimCLR half: crop the reference figure whitespace and redraw loss large."""
    def PX(f): return Emu(int(x) + int(f * int(w)))
    def PY(f): return Emu(int(y) + int(f * int(h)))
    def WW(f): return Emu(int(f * int(w)))
    def HH(f): return Emu(int(f * int(h)))

    def label(fx, fy, fw, fh, text, size=8.0, bold=False, color=INK,
              align=PP_ALIGN.CENTER, name=FONT):
        _text(slide, PX(fx), PY(fy), WW(fw), HH(fh),
              [[(text, dict(size=size, bold=bold, color=color, name=name))]],
              align=align, anchor=MSO_ANCHOR.MIDDLE)

    _rect(slide, x, y, w, h, WHITE, line=RGBColor(0xEA, 0xEE, 0xF4), line_w=Pt(1.0))

    # Reference figure: use a whitespace-trimmed paper crop so the SimCLR mechanism is readable.
    _img_fit(slide, "ref_simclr_framework_crop.png", PX(0.015), PY(0.010), WW(0.97), HH(0.660), frame=False)

    # Loss is a crop from the paper figure, not a retyped approximation.
    _img_fit(slide, "ref_simclr_loss_crop.png", PX(0.040), PY(0.690), WW(0.92), HH(0.275), frame=False)


def _draw_pos_neg_hard_old(slide, x, y, w, h):
    """Three-column explanation: positive / negative / hard negative."""
    BLACK = RGBColor(0x16, 0x18, 0x1D)
    DARK = NAVY
    MID = RGBColor(0x55, 0x60, 0x70)
    LIGHT = RGBColor(0xEA, 0xEE, 0xF4)
    HARD_FILL = RGBColor(0xC5, 0xCB, 0xD6)
    SOFT = RGBColor(0xF8, 0xFA, 0xFC)

    def PX(f): return Emu(int(x) + int(f * int(w)))
    def PY(f): return Emu(int(y) + int(f * int(h)))
    def WW(f): return Emu(int(f * int(w)))
    def HH(f): return Emu(int(f * int(h)))

    def label(fx, fy, fw, fh, text, size=8.0, bold=False, color=INK,
              align=PP_ALIGN.CENTER, name=FONT):
        _text(slide, PX(fx), PY(fy), WW(fw), HH(fh),
              [[(text, dict(size=size, bold=bold, color=color, name=name))]],
              align=align, anchor=MSO_ANCHOR.MIDDLE)

    def connector(fx1, fy1, fx2, fy2, color, width=1.0):
        cn = slide.shapes.add_connector(1, PX(fx1), PY(fy1), PX(fx2), PY(fy2))
        cn.line.color.rgb = color
        cn.line.width = Pt(width)
        cn.shadow.inherit = False
        return cn

    def dot(fx, fy, color, shape=MSO_SHAPE.OVAL, size=0.026, line=None):
        sz = WW(size)
        _rect(slide, Emu(int(PX(fx)) - int(sz) // 2), Emu(int(PY(fy)) - int(sz) // 2),
              sz, sz, color, line=line, shape=shape, line_w=Pt(1.1))

    _rect(slide, x, y, w, h, WHITE, line=LIGHT, line_w=Pt(1.0))

    cols = [
        ("Positive", "NUMERATOR", "sim ↑", "pull\ncompact"),
        ("Negative", "DENOMINATOR", "sim low", "weak push"),
        ("Hard negative", "DENOMINATOR", "sim high", "boundary"),
    ]
    col_x = [0.030, 0.350, 0.670]
    col_w = 0.300
    for i, (head, role, loss_txt, emb_txt) in enumerate(cols):
        cx = col_x[i]
        _rect(slide, PX(cx), PY(0.045), WW(col_w), HH(0.875), SOFT, line=LIGHT)
        _rect(slide, PX(cx), PY(0.045), WW(col_w), HH(0.018), COVER_BAR)
        label(cx, 0.080, col_w, 0.075, head, size=15.2, bold=True, color=NAVY)

        # local embedding sketch
        qx, qy = cx + col_w * 0.34, 0.335
        dot(qx, qy, BLACK, size=0.044)
        label(qx - 0.055, qy + 0.052, 0.065, 0.055, "q", size=15.0, bold=True,
              color=BLACK, name="Cambria Math")
        if i == 0:
            sx, sy = cx + col_w * 0.68, 0.260
            dot(sx, sy, WHITE, line=BLACK, size=0.052)
            connector(sx, sy, qx, qy, DARK, width=1.4)
            label(sx + 0.030, sy - 0.070, 0.09, 0.055, "p⁺", size=15.0,
                  bold=True, color=BLACK, name="Cambria Math")
        elif i == 1:
            sx, sy = cx + col_w * 0.76, 0.385
            dot(sx, sy, DARK, MSO_SHAPE.ISOSCELES_TRIANGLE, size=0.050)
            connector(qx, qy, sx, sy, MID, width=1.0)
            label(sx - 0.025, sy + 0.058, 0.09, 0.055, "n⁻", size=14.8,
                  bold=True, color=DARK, name="Cambria Math")
        else:
            sx, sy = cx + col_w * 0.60, 0.300
            dot(sx, sy, DARK, MSO_SHAPE.DIAMOND, size=0.052)
            connector(qx, qy, sx, sy, DARK, width=1.6)
            label(sx + 0.030, sy - 0.070, 0.10, 0.055, "h⁻", size=14.8,
                  bold=True, color=DARK, name="Cambria Math")

        _rect(slide, PX(cx + 0.020), PY(0.515), WW(col_w - 0.040), HH(0.155),
              WHITE, line=LIGHT)
        label(cx + 0.035, 0.528, col_w - 0.070, 0.052, role,
              size=12.4, bold=True, color=NAVY, name="Consolas")
        label(cx + 0.035, 0.585, col_w - 0.070, 0.058, loss_txt,
              size=14.0, bold=True, color=BLACK)

        _rect(slide, PX(cx + 0.020), PY(0.720), WW(col_w - 0.040), HH(0.135),
              WHITE, line=LIGHT)
        label(cx + 0.035, 0.740, col_w - 0.070, 0.090, emb_txt,
              size=14.0, bold=True, color=BLACK)


def _draw_pos_neg_hard(slide, x, y, w, h):
    """Paper-style role panel: no nested text boxes; large labels and separators."""
    BLACK = RGBColor(0x16, 0x18, 0x1D)
    DARK = NAVY
    MID = RGBColor(0x55, 0x60, 0x70)
    LIGHT = RGBColor(0xEA, 0xEE, 0xF4)
    HARD_FILL = RGBColor(0xC5, 0xCB, 0xD6)

    def PX(f): return Emu(int(x) + int(f * int(w)))
    def PY(f): return Emu(int(y) + int(f * int(h)))
    def WW(f): return Emu(int(f * int(w)))
    def HH(f): return Emu(int(f * int(h)))

    def label(fx, fy, fw, fh, text, size=12.0, bold=False, color=INK,
              align=PP_ALIGN.CENTER, name=FONT):
        _text(slide, PX(fx), PY(fy), WW(fw), HH(fh),
              [[(text, dict(size=size, bold=bold, color=color, name=name))]],
              align=align, anchor=MSO_ANCHOR.MIDDLE)

    def connector(fx1, fy1, fx2, fy2, color, width=1.0):
        cn = slide.shapes.add_connector(1, PX(fx1), PY(fy1), PX(fx2), PY(fy2))
        cn.line.color.rgb = color
        cn.line.width = Pt(width)
        cn.shadow.inherit = False
        return cn

    def dot(fx, fy, color, shape=MSO_SHAPE.OVAL, size=0.026, line=None):
        sz = WW(size)
        _rect(slide, Emu(int(PX(fx)) - int(sz) // 2), Emu(int(PY(fy)) - int(sz) // 2),
              sz, sz, color, line=line, shape=shape, line_w=Pt(1.1))

    def vline(fx, fy1, fy2, width=1.0):
        cn = slide.shapes.add_connector(1, PX(fx), PY(fy1), PX(fx), PY(fy2))
        cn.line.color.rgb = LIGHT
        cn.line.width = Pt(width)
        cn.shadow.inherit = False

    def hline(fx1, fx2, fy, width=1.0):
        cn = slide.shapes.add_connector(1, PX(fx1), PY(fy), PX(fx2), PY(fy))
        cn.line.color.rgb = LIGHT
        cn.line.width = Pt(width)
        cn.shadow.inherit = False

    _rect(slide, x, y, w, h, WHITE, line=LIGHT, line_w=Pt(1.0))
    vline(0.335, 0.07, 0.93)
    vline(0.665, 0.07, 0.93)
    hline(0.055, 0.945, 0.535)

    cols = [
        ("Positive", "PULL TERM", "similar", "pull closer", 0.035),
        ("Negative", "PUSH TERM", "far away", "weak push", 0.365),
        ("Hard negative", "PUSH TERM", "very similar", "strong push", 0.695),
    ]
    col_w = 0.270
    for i, (head, role, loss_txt, emb_txt, cx) in enumerate(cols):
        label(cx, 0.070, col_w, 0.075, head, size=16.2, bold=True, color=NAVY)

        qx, qy = cx + col_w * 0.33, 0.325
        dot(qx, qy, BLACK, size=0.050)
        label(qx - 0.062, qy + 0.060, 0.070, 0.060, "q",
              size=16.0, bold=True, color=BLACK, name="Cambria Math")
        if i == 0:
            sx, sy = cx + col_w * 0.70, 0.245
            dot(sx, sy, WHITE, line=BLACK, size=0.060)
            connector(sx, sy, qx, qy, DARK, width=1.55)
            label(sx + 0.032, sy - 0.075, 0.095, 0.060, "p⁺",
                  size=16.0, bold=True, color=BLACK, name="Cambria Math")
        elif i == 1:
            sx, sy = cx + col_w * 0.76, 0.385
            dot(sx, sy, DARK, MSO_SHAPE.ISOSCELES_TRIANGLE, size=0.058)
            connector(qx, qy, sx, sy, MID, width=1.0)
            label(sx - 0.025, sy + 0.065, 0.095, 0.060, "n⁻",
                  size=15.8, bold=True, color=DARK, name="Cambria Math")
        else:
            sx, sy = cx + col_w * 0.62, 0.295
            dot(sx, sy, HARD_FILL, line=BLACK, size=0.060)
            connector(qx, qy, sx, sy, DARK, width=1.6)
            label(sx + 0.032, sy - 0.075, 0.10, 0.060, "h⁻",
                  size=15.8, bold=True, color=DARK, name="Cambria Math")

        label(cx, 0.590, col_w, 0.060, role, size=13.5,
              bold=True, color=NAVY, name="Consolas")
        label(cx, 0.665, col_w, 0.070, loss_txt, size=16.2,
              bold=True, color=BLACK)
        label(cx, 0.815, col_w, 0.090, emb_txt, size=15.0,
              bold=True, color=BLACK)


def _draw_hardneg_half_old(slide, x, y, w, h):
    """Legacy half-slide hard-negative explanation.
    Shows sample role -> loss inclusion -> embedding effect without turning it into a full slide."""
    RED = RGBColor(0xC9, 0x35, 0x2B)
    BLUE = RGBColor(0x2B, 0x66, 0xD9)
    GREEN = RGBColor(0x2B, 0xA6, 0x6B)
    ORANGE = RGBColor(0xE0, 0x8A, 0x1E)
    SOFT = RGBColor(0xF7, 0xF9, 0xFC)
    SOFT_RED = RGBColor(0xFF, 0xF1, 0xEE)
    SOFT_BLUE = RGBColor(0xF0, 0xF5, 0xFF)
    SOFT_GREEN = RGBColor(0xEC, 0xF8, 0xF2)
    GRID = RGBColor(0xE8, 0xEC, 0xF2)

    _rect(slide, x, y, w, h, WHITE, line=LINE)
    head_h = Inches(0.36)
    note_h = Inches(0.20)
    _rect(slide, x, y, w, head_h, NAVY)
    _text(slide, x, y, w, head_h,
          [[("B. hard-negative filtering: loss에 넣을 negative를 다시 고른다",
             dict(size=11.9, bold=True, color=WHITE))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bx, by = x + Inches(0.16), y + Inches(0.48)
    bw = w - Inches(0.32)
    col_gap = Inches(0.11)
    col_w = Emu((int(bw) - int(col_gap) * 2) // 3)
    img_h = Inches(0.55)
    cap_h = Inches(0.18)
    role_h = Inches(0.32)
    samples = [
        ("Positive", "wafer_center_scratch.png", 0.50, 0.50, GREEN, SOFT_GREEN, "numerator\nanchor로 당김"),
        ("False negative", "wafer_edge_top_scratch.png", 0.50, 0.00, RED, SOFT_RED, "exclude from Σ\n잘못된 반발 제거"),
        ("Semi-hard negative", "wafer_ringdots.png", 0.50, 0.50, BLUE, SOFT_BLUE, "denominator\n경계 학습 유지"),
    ]
    for i, (label, img, fx_focus, fy_focus, color, fill, role) in enumerate(samples):
        cx = Emu(int(bx) + i * (int(col_w) + int(col_gap)))
        _rect(slide, cx, by, col_w, img_h + cap_h + role_h, fill, line=LINE)
        _img_cover(slide, img, cx + Inches(0.08), by + Inches(0.06),
                   col_w - Inches(0.16), img_h - Inches(0.07),
                   focus_x=fx_focus, focus_y=fy_focus)
        _text(slide, cx + Inches(0.04), by + img_h + Inches(0.01),
              col_w - Inches(0.08), cap_h,
              [[(label, dict(size=6.9, bold=True, color=color))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _rect(slide, cx + Inches(0.09), by + img_h + cap_h + Inches(0.02),
              col_w - Inches(0.18), role_h - Inches(0.10), WHITE, line=RGBColor(0xD9,0xE0,0xEA))
        _text(slide, cx + Inches(0.12), by + img_h + cap_h + Inches(0.055),
              col_w - Inches(0.24), role_h - Inches(0.16),
              [[(role, dict(size=5.9, bold=True, color=NAVY))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    eq_y = by + Inches(1.16)
    eq_h = Inches(0.32)
    _rect(slide, bx, eq_y, bw, eq_h, RGBColor(0xFA,0xFB,0xFD), line=LINE)
    _text(slide, bx + Inches(0.08), eq_y + Inches(0.04), bw - Inches(0.16), eq_h - Inches(0.08),
          [[("L = -log exp(s(a,p+)/τ) / [exp(s(a,p+)/τ)+Σkeep exp(s(a,n)/τ)]   FN 제외",
             dict(size=7.7, bold=True, color=NAVY, name="Consolas"))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    plot_y = eq_y + eq_h + Inches(0.13)
    note_y = y + h - note_h - Inches(0.10)
    plot_h = Emu(max(int(Inches(0.72)), int(note_y - plot_y - Inches(0.08))))
    plot_w = Emu((int(bw) - int(Inches(0.22))) // 2)

    def dot(cx, cy, color, shape=MSO_SHAPE.OVAL, sz=0.09):
        _rect(slide, cx - Inches(sz/2), cy - Inches(sz/2), Inches(sz), Inches(sz),
              color, shape=shape)

    def connector(x1, y1, x2, y2, color):
        cn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        cn.line.color.rgb = color
        cn.line.width = Pt(1.0)
        cn.shadow.inherit = False
        return cn

    # Before: false negative remains in denominator, causing repulsion inside the same pattern.
    px1 = bx
    _rect(slide, px1, plot_y, plot_w, plot_h, WHITE, line=LINE)
    _text(slide, px1 + Inches(0.08), plot_y + Inches(0.05), plot_w - Inches(0.16), Inches(0.16),
          [[("Naive denominator", dict(size=7.8, bold=True, color=RED))]],
          align=PP_ALIGN.CENTER)
    base_x, base_y = px1 + Emu(int(plot_w * 0.42)), plot_y + Emu(int(plot_h * 0.58))
    for dx, dy in [(-0.10, -0.06), (-0.02, 0.03), (0.07, -0.02)]:
        dot(base_x + Inches(dx), base_y + Inches(dy), GREEN)
    fnx, fny = base_x + Inches(0.20), base_y - Inches(0.12)
    dot(fnx, fny, RED, MSO_SHAPE.DIAMOND, sz=0.105)
    connector(base_x + Inches(0.08), base_y - Inches(0.02), fnx + Inches(0.23), fny - Inches(0.13), RED)
    _text(slide, px1 + Inches(0.10), plot_y + plot_h - Inches(0.24), plot_w - Inches(0.20), Inches(0.19),
          [[("same-pattern sample에도 repulsive gradient",
             dict(size=6.2, bold=True, color=INK))]], align=PP_ALIGN.CENTER)

    # After: false negative is removed; semi-hard negatives define the boundary.
    px2 = Emu(int(bx) + int(plot_w) + int(Inches(0.22)))
    _rect(slide, px2, plot_y, plot_w, plot_h, WHITE, line=LINE)
    _text(slide, px2 + Inches(0.08), plot_y + Inches(0.05), plot_w - Inches(0.16), Inches(0.16),
          [[("Positive-aware filtering", dict(size=7.8, bold=True, color=BLUE))]],
          align=PP_ALIGN.CENTER)
    c2x, c2y = px2 + Emu(int(plot_w * 0.34)), plot_y + Emu(int(plot_h * 0.58))
    for dx, dy in [(-0.07, -0.04), (0.00, 0.04), (0.07, -0.02), (0.02, -0.08)]:
        dot(c2x + Inches(dx), c2y + Inches(dy), GREEN)
    dot(px2 + Emu(int(plot_w * 0.70)), plot_y + Emu(int(plot_h * 0.42)), BLUE, MSO_SHAPE.ISOSCELES_TRIANGLE, sz=0.11)
    dot(px2 + Emu(int(plot_w * 0.75)), plot_y + Emu(int(plot_h * 0.66)), BLUE, MSO_SHAPE.ISOSCELES_TRIANGLE, sz=0.11)
    connector(px2 + Emu(int(plot_w * 0.58)), plot_y + Emu(int(plot_h * 0.33)),
              px2 + Emu(int(plot_w * 0.58)), plot_y + Emu(int(plot_h * 0.82)), BLUE)
    _text(slide, px2 + Inches(0.10), plot_y + plot_h - Inches(0.24), plot_w - Inches(0.20), Inches(0.19),
          [[("compact cluster + sharper boundary",
             dict(size=6.2, bold=True, color=INK))]], align=PP_ALIGN.CENTER)

    _rect(slide, bx, note_y, bw, note_h, SOFT, line=GRID)
    _text(slide, bx + Inches(0.08), note_y, bw - Inches(0.16), note_h,
          [[("too-close 후보는 버리고, margin 밖 semi-hard negative만 남겨 Unknown cluster를 안정화",
             dict(size=7.0, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)


def _draw_hardneg_half(slide, x, y, w, h):
    """Clean half-slide figure: embedding geometry + InfoNCE role."""
    BLACK = RGBColor(0x12, 0x12, 0x12)
    GRAY = RGBColor(0x9A, 0xA3, 0xB1)
    LIGHT = RGBColor(0xEA, 0xEE, 0xF4)
    RED = RGBColor(0xC9, 0x2A, 0x2A)
    BLUE = RGBColor(0x2B, 0x66, 0xD9)
    GREEN = RGBColor(0x2B, 0xA6, 0x6B)

    def PX(f): return Emu(int(x) + int(f * int(w)))
    def PY(f): return Emu(int(y) + int(f * int(h)))
    def WW(f): return Emu(int(f * int(w)))
    def HH(f): return Emu(int(f * int(h)))

    def dashed(sp):
        ln = sp.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))

    def label(fx, fy, fw, fh, text, size=8.0, bold=False, color=BLACK,
              align=PP_ALIGN.CENTER, name=FONT):
        _text(slide, PX(fx), PY(fy), WW(fw), HH(fh),
              [[(text, dict(size=size, bold=bold, color=color, name=name))]],
              align=align, anchor=MSO_ANCHOR.MIDDLE)

    def connector(fx1, fy1, fx2, fy2, color=GRAY, width=0.9, dash=False):
        cn = slide.shapes.add_connector(1, PX(fx1), PY(fy1), PX(fx2), PY(fy2))
        cn.line.color.rgb = color
        cn.line.width = Pt(width)
        cn.shadow.inherit = False
        if dash:
            dashed(cn)
        return cn

    def point(fx, fy, shape, fill, line=None, size=0.030, lw=1.0):
        sz = WW(size)
        sp = _rect(slide, Emu(int(PX(fx)) - int(sz) // 2), Emu(int(PY(fy)) - int(sz) // 2),
                   sz, sz, fill, line=line, shape=shape, line_w=Pt(lw))
        return sp

    _rect(slide, x, y, w, h, WHITE, line=LIGHT, line_w=Pt(1.0))
    label(0.02, 0.02, 0.47, 0.065, "Positive-aware hard negative mining",
          size=9.4, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    label(0.62, 0.02, 0.34, 0.055, "s(q,x)=cos(q,x)",
          size=7.4, bold=True, color=MUTED, align=PP_ALIGN.RIGHT, name="Consolas")

    qx, qy = 0.42, 0.37
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, PX(qx - 0.19), PY(qy - 0.22),
                                   WW(0.38), HH(0.44))
    inner.fill.background()
    inner.line.color.rgb = RED
    inner.line.width = Pt(1.0)
    inner.shadow.inherit = False
    dashed(inner)
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, PX(qx - 0.32), PY(qy - 0.34),
                                   WW(0.64), HH(0.68))
    outer.fill.background()
    outer.line.color.rgb = GRAY
    outer.line.width = Pt(1.0)
    outer.shadow.inherit = False
    dashed(outer)

    p = (0.54, 0.25)
    fn1, fn2 = (0.48, 0.36), (0.56, 0.42)
    sh1, sh2, sh3 = (0.22, 0.24), (0.62, 0.54), (0.67, 0.27)

    connector(p[0], p[1], qx, qy, GREEN, width=1.05)
    connector(qx, qy, fn1[0], fn1[1], RED, width=1.0, dash=True)
    connector(qx, qy, fn2[0], fn2[1], RED, width=1.0, dash=True)
    connector(qx, qy, sh2[0], sh2[1], BLUE, width=1.0)
    connector(qx, qy, sh3[0], sh3[1], BLUE, width=1.0)

    point(qx, qy, MSO_SHAPE.OVAL, BLACK, size=0.036)
    label(qx - 0.045, qy + 0.035, 0.06, 0.05, "q", size=12.5,
          bold=True, color=BLACK, name="Cambria Math")
    point(p[0], p[1], MSO_SHAPE.OVAL, WHITE, line=BLACK, size=0.040, lw=1.35)
    label(p[0] + 0.035, p[1] - 0.045, 0.08, 0.05, "p⁺", size=12.2,
          bold=True, color=BLACK, name="Cambria Math")
    for fx, fy in [fn1, fn2]:
        point(fx, fy, MSO_SHAPE.DIAMOND, RED, size=0.035)
    for fx, fy in [sh1, sh2, sh3]:
        point(fx, fy, MSO_SHAPE.ISOSCELES_TRIANGLE, BLUE, size=0.037)

    label(0.18, 0.09, 0.34, 0.05, "false-negative exclusion margin",
          size=7.1, bold=True, color=RED, align=PP_ALIGN.LEFT)
    label(0.62, 0.58, 0.30, 0.05, "semi-hard negatives kept",
          size=7.2, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    label(0.62, 0.18, 0.28, 0.075, "denominator\nboundary gradient",
          size=7.0, bold=True, color=BLUE)
    label(0.36, 0.47, 0.30, 0.065, "removed:\nno repulsive gradient",
          size=6.9, bold=True, color=RED)
    label(0.47, 0.13, 0.28, 0.05, "numerator attraction",
          size=7.1, bold=True, color=GREEN)

    # Embedding effect, kept as a small paper-style inset.
    _rect(slide, PX(0.07), PY(0.65), WW(0.22), HH(0.18), WHITE, line=LIGHT)
    label(0.08, 0.665, 0.20, 0.035, "naive", size=7.0, bold=True, color=RED)
    for fx, fy in [(0.13, 0.74), (0.15, 0.70), (0.18, 0.76)]:
        point(fx, fy, MSO_SHAPE.OVAL, GREEN, size=0.022)
    point(0.23, 0.68, MSO_SHAPE.DIAMOND, RED, size=0.025)
    connector(0.18, 0.73, 0.24, 0.67, RED, width=0.9)

    _rect(slide, PX(0.33), PY(0.65), WW(0.27), HH(0.18), WHITE, line=LIGHT)
    label(0.345, 0.665, 0.24, 0.035, "filtered", size=7.0, bold=True, color=BLUE)
    for fx, fy in [(0.41, 0.74), (0.425, 0.705), (0.445, 0.735), (0.435, 0.77)]:
        point(fx, fy, MSO_SHAPE.OVAL, GREEN, size=0.022)
    point(0.54, 0.69, MSO_SHAPE.ISOSCELES_TRIANGLE, BLUE, size=0.026)
    point(0.55, 0.78, MSO_SHAPE.ISOSCELES_TRIANGLE, BLUE, size=0.026)
    connector(0.50, 0.67, 0.50, 0.81, BLUE, width=0.9)

    _rect(slide, PX(0.06), PY(0.87), WW(0.88), HH(0.10), WHITE, line=LIGHT)
    label(0.075, 0.875, 0.85, 0.085,
          "L = -log  exp(s(q,p⁺)/τ) / [ exp(s(q,p⁺)/τ) + Σₙ∈N_keep exp(s(q,n)/τ) ]",
          size=8.4, bold=True, color=NAVY, name="Cambria Math")


def s_unknown_simclr_hardneg(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, "P1 | Unknown 대조 학습 ①", "그림으로 설명: positive는 붙이고, false negative는 뺀다")

    _rect(slide, Inches(0.80), Inches(1.70), Inches(11.80), Inches(0.46),
          RGBColor(0xF4,0xF6,0xF8), line=LINE)
    _rect(slide, Inches(0.80), Inches(1.70), Inches(0.10), Inches(0.46), ACCENT)
    _text(slide, Inches(1.00), Inches(1.70), Inches(11.35), Inches(0.46),
          [[("말할 순서: 같은 wafer 두 view를 positive로 묶고, 다른 wafer 중 유용한 hard negative만 남기며, 너무 가까운 후보는 false negative로 제외합니다.",
             dict(size=13.2, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)

    # SimCLR panel
    lx, ly, lw, lh = Inches(0.80), Inches(2.38), Inches(5.70), Inches(3.82)
    _rect(slide, lx, ly, lw, lh, WHITE, line=LINE)
    _rect(slide, lx, ly, lw, Inches(0.42), NAVY)
    _text(slide, lx, ly, lw, Inches(0.42),
          [[("A. SimCLR / InfoNCE", dict(size=13.5, bold=True, color=WHITE))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _img_fit(slide, "ref_simclr_with_loss.png", lx+Inches(0.28), ly+Inches(0.58),
             lw-Inches(0.56), Inches(2.18), frame=False)
    _number_chip(slide, lx+Inches(0.25), ly+Inches(2.94), 1, "same wafer",
                 "두 augmentation이 positive pair", w=Inches(1.65), h=Inches(0.58))
    _number_chip(slide, lx+Inches(2.05), ly+Inches(2.94), 2, "encoder",
                 "wafer를 embedding q, k로 변환", w=Inches(1.70), h=Inches(0.58))
    _number_chip(slide, lx+Inches(3.90), ly+Inches(2.94), 3, "InfoNCE",
                 "positive는 가깝게, 나머지는 멀게", w=Inches(1.55), h=Inches(0.58))

    # Hard negative panel
    rx, ry, rw, rh = Inches(6.78), Inches(2.38), Inches(5.82), Inches(3.82)
    _rect(slide, rx, ry, rw, rh, WHITE, line=LINE)
    _draw_hardneg_half(slide, rx, ry, rw, rh)

    _rect(slide, Inches(0.95), Inches(6.52), Inches(11.45), Inches(0.36),
          RGBColor(0xFF,0xF7,0xE8), line=RGBColor(0xF2,0xC4,0x6D))
    _text(slide, Inches(1.05), Inches(6.52), Inches(11.25), Inches(0.36),
          [[("발표 멘트: 'false negative는 loss 분모에서 빼고, margin 밖 semi-hard negative만 남겨 같은 불량군은 모으고 다른 불량군 경계는 세웠습니다.'",
             dict(size=11.5, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def s_unknown_moco_explain(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, "P1 | Unknown 대조 학습 ②", "MoCo queue: 작은 batch에서도 negative dictionary를 크게 유지")

    _text(slide, Inches(0.85), Inches(1.72), Inches(11.55), Inches(0.42),
          [[("MoCo는 '새로운 모델'이 아니라 contrastive learning에서 negative를 많이, 안정적으로 유지하는 장치입니다.",
             dict(size=15.0, bold=True, color=NAVY))]])

    ix, iy, iw, ih = Inches(0.85), Inches(2.28), Inches(6.50), Inches(3.62)
    _rect(slide, ix, iy, iw, ih, WHITE, line=LINE)
    _img_fit(slide, "ref_moco_fig1.png", ix+Inches(0.22), iy+Inches(0.22),
             iw-Inches(0.44), ih-Inches(0.44), frame=False)

    rx = Inches(7.65)
    _rect(slide, rx, Inches(2.28), Inches(4.72), Inches(3.62), PANEL, line=LINE)
    _rect(slide, rx, Inches(2.28), Inches(0.10), Inches(3.62), ACCENT)
    _text(slide, rx+Inches(0.25), Inches(2.48), Inches(4.25), Inches(0.32),
          [[("그림 읽는 순서", dict(size=15.0, bold=True, color=NAVY))]])
    _number_chip(slide, rx+Inches(0.25), Inches(2.95), 1, "query encoder",
                 "현재 학습되는 wafer encoder", w=Inches(4.05), h=Inches(0.50))
    _number_chip(slide, rx+Inches(0.25), Inches(3.55), 2, "momentum encoder",
                 "query encoder를 천천히 따라가는 key encoder", w=Inches(4.05), h=Inches(0.56))
    _number_chip(slide, rx+Inches(0.25), Inches(4.22), 3, "queue",
                 "이전 batch wafer embedding을 negative dictionary로 저장", w=Inches(4.05), h=Inches(0.62))
    _number_chip(slide, rx+Inches(0.25), Inches(4.98), 4, "contrastive loss",
                 "q는 positive key와 가깝게, queue negatives와 멀게", w=Inches(4.05), h=Inches(0.58))

    _rect(slide, Inches(0.95), Inches(6.28), Inches(5.45), Inches(0.52),
          RGBColor(0xF4,0xF7,0xFB), line=LINE)
    _text(slide, Inches(1.05), Inches(6.28), Inches(5.25), Inches(0.52),
          [[("InfoNCE: q^T k+는 키우고, q^T k_i는 낮춤", dict(size=12.5, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(6.82), Inches(6.28), Inches(5.45), Inches(0.52),
          RGBColor(0xF4,0xF7,0xFB), line=LINE)
    _text(slide, Inches(6.92), Inches(6.28), Inches(5.25), Inches(0.52),
          [[("momentum update: θ_k ← mθ_k + (1-m)θ_q", dict(size=12.5, bold=True, color=NAVY))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def s_unknown_moco_neco(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, d.get("kicker"), d["title"])
    _text(slide, Inches(0.85), Inches(1.58), Inches(11.55), Inches(0.40),
          [[("MoCo expands the negative dictionary; NeCo regularizes patch-neighbor order to preserve local failure-pattern structure.",
             dict(size=13.8, bold=True, color=NAVY))]],
          anchor=MSO_ANCHOR.MIDDLE)

    y = Inches(2.12)
    h = Inches(4.42)
    x0 = Inches(0.78)
    gap = Inches(0.28)
    left_w = Inches(3.58)
    right_w = Inches(7.89)
    right_x = Emu(int(x0) + int(left_w) + int(gap))

    # 1:2 visual weight: MoCo is the supporting queue mechanism, NeCo is the main patch-level figure.
    _rect(slide, x0, y, left_w, h, WHITE, line=LINE)
    _rect(slide, x0, y, left_w, Inches(0.38), NAVY)
    _text(slide, x0, y, left_w, Inches(0.38),
          [[("MoCo queue", dict(size=13.2, bold=True, color=WHITE))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _img_fit(slide, "ref_moco_fig1.png", x0+Inches(0.16), y+Inches(0.56),
             left_w-Inches(0.32), Inches(2.22), frame=False)
    _rect(slide, x0+Inches(0.18), y+Inches(3.03), left_w-Inches(0.36), Inches(1.10),
          PANEL, line=LINE)
    _text(slide, x0+Inches(0.30), y+Inches(3.13), left_w-Inches(0.60), Inches(0.28),
          [[("Role in loss", dict(size=11.2, bold=True, color=NAVY))]])
    _text(slide, x0+Inches(0.30), y+Inches(3.44), left_w-Inches(0.60), Inches(0.54),
          [[("maintains queue negatives beyond the current minibatch",
             dict(size=10.6, color=INK))]],
          anchor=MSO_ANCHOR.MIDDLE)

    _rect(slide, right_x, y, right_w, h, WHITE, line=LINE)
    _rect(slide, right_x, y, right_w, Inches(0.38), NAVY)
    _text(slide, right_x, y, right_w, Inches(0.38),
          [[("NeCo patch consistency", dict(size=13.2, bold=True, color=WHITE))]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _img_fit(slide, "ref_neco_method.png", right_x+Inches(0.14), y+Inches(0.50),
             right_w-Inches(0.28), Inches(3.08), frame=False)

    chips = [
        ("ROI-aligned patches", "same original region"),
        ("neighbor ordering", "patch-level similarity"),
        ("consistency loss", "compact local structure"),
    ]
    chip_y = y + Inches(3.74)
    chip_gap = Inches(0.12)
    chip_w = Emu((int(right_w) - int(chip_gap) * 2 - int(Inches(0.36))) // 3)
    for i, (head, sub) in enumerate(chips):
        cx = Emu(int(right_x) + int(Inches(0.18)) + i*(int(chip_w)+int(chip_gap)))
        _rect(slide, cx, chip_y, chip_w, Inches(0.52), PANEL, line=LINE)
        _text(slide, cx+Inches(0.08), chip_y+Inches(0.06), chip_w-Inches(0.16), Inches(0.18),
              [[(head, dict(size=9.9, bold=True, color=NAVY))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, cx+Inches(0.08), chip_y+Inches(0.27), chip_w-Inches(0.16), Inches(0.18),
              [[(sub, dict(size=8.9, color=MUTED))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


def s_unknown_neco_explain(slide, d, idx):
    _bg(slide, WHITE)
    _title_block(slide, "P1 | Unknown 대조 학습 ③", "NeCo patch consistency: 같은 위치의 patch-neighbor order를 맞춤")

    _text(slide, Inches(0.85), Inches(1.62), Inches(11.7), Inches(0.46),
          [[("NeCo는 wafer 전체 embedding이 아니라 patch 단위 경계를 안정화합니다. 핵심은 두 positive view에서 같은 원본 좌표 영역만 ROI Align으로 맞추는 것입니다.",
             dict(size=14.4, bold=True, color=NAVY))]])

    ix, iy, iw, ih = Inches(0.80), Inches(2.18), Inches(11.75), Inches(2.92)
    _rect(slide, ix, iy, iw, ih, WHITE, line=LINE)
    _img_fit(slide, "ref_neco_method.png", ix+Inches(0.12), iy+Inches(0.16),
             iw-Inches(0.24), ih-Inches(0.32), frame=False)

    y = Inches(5.45)
    x0 = Inches(0.78)
    w = Inches(2.33)
    gap = Inches(0.13)
    steps = [
        ("two views", "같은 wafer에서 두 crop/augmentation 생성"),
        ("ROI Align", "두 view가 공유하는 원본 좌표 patch만 정렬"),
        ("reference patches", "batch patch들과 cosine distance 계산"),
        ("sorting", "nearest-neighbor 순서를 differentiable sorting"),
        ("NeCo loss", "teacher/student view의 neighbor order를 일치"),
    ]
    for i, (head, body) in enumerate(steps):
        x = Emu(int(x0) + i*(int(w)+int(gap)))
        _rect(slide, x, y, w, Inches(0.86), RGBColor(0xF4,0xF7,0xFB), line=LINE)
        _rect(slide, x, y, Inches(0.08), Inches(0.86), ACCENT)
        _text(slide, x+Inches(0.18), y+Inches(0.08), w-Inches(0.26), Inches(0.22),
              [[(f"{i+1}. {head}", dict(size=10.8, bold=True, color=NAVY))]])
        _text(slide, x+Inches(0.18), y+Inches(0.34), w-Inches(0.26), Inches(0.42),
              [[(body, dict(size=9.2, color=INK))]])
        if i < len(steps)-1:
            _rect(slide, Emu(int(x)+int(w)-int(Inches(0.03))), y+Inches(0.34),
                  Inches(0.20), Inches(0.20), ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)

    _rect(slide, Inches(1.05), Inches(6.55), Inches(11.20), Inches(0.34),
          RGBColor(0xF4,0xF6,0xF8), line=LINE)
    _text(slide, Inches(1.15), Inches(6.55), Inches(11.0), Inches(0.34),
          [[("주의: 여기서 ROI는 detection label box가 아니라, 두 augmentation이 공통으로 보는 원본 wafer 좌표 영역입니다.",
             dict(size=11.2, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx)


DISPATCH = {"title": s_title, "section": s_section, "stats": s_stats, "bullets": s_bullets,
            "two_col": s_two_col, "image_grid": s_image_grid, "table": s_table, "closing": s_closing,
            "flow": s_flow, "timeline": s_timeline, "cards": s_cards, "pipeline": s_pipeline,
            "papertext": s_papertext, "archflow": s_archflow,
            "p1_known_perf": s_p1_known_perf, "unknown_ablation": s_unknown_ablation,
            "unknown_simclr_hardneg": s_unknown_simclr_hardneg,
            "unknown_moco_neco": s_unknown_moco_neco,
            "unknown_moco_explain": s_unknown_moco_explain,
            "unknown_neco_explain": s_unknown_neco_explain,
            "p2_intro": s_p2_intro, "p2_fcmpm": s_p2_fcmpm, "p2_validation": s_p2_validation,
            "p2_selection": s_p2_selection,
            "p3_intro": s_p3_intro, "p3_generator": s_p3_generator, "p3_result": s_p3_result}


def build(spec_path, out_path):
    spec = json.load(open(spec_path, encoding="utf-8"))
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    for i, sd in enumerate(spec["slides"], 1):
        slide = prs.slides.add_slide(blank)
        DISPATCH.get(sd["type"], s_bullets)(slide, sd, i)
    prs.save(out_path)
    return len(spec["slides"])


if __name__ == "__main__":
    sp = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "spec.json")
    op = sys.argv[2] if len(sys.argv) > 2 else os.path.join(HERE, "발표_AI인증_최호길.pptx")
    n = build(sp, op)
    print(f"built {n} slides -> {op}")
